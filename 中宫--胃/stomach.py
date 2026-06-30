#!/usr/bin/env python3
"""
中宫--胃 stomach.py v5.0 — 知识消化中枢
部署位置: F:\公司知识平台\中宫--胃\
功能: 扫描原始文件 → 解析 → 清洗 → 脱敏 → 语义分块 → 分类 → 蒸馏 → 实体提取 → 推送 kb-server

整合来源:
  清洗程序 v15: 解析器链/PDF结构化/OCR降级/语义分块/脱敏/17类分类/去重/守护轮询/推送格式
  胃 v4.2:     Wiki蒸馏/实体提取/本地存储
"""

import os, sys, re, json, time, hashlib, logging
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Tuple
from logging.handlers import RotatingFileHandler

# ==================== 路径配置 ====================
BASE_DIR = Path(__file__).resolve().parent  # F:\公司知识平台\中宫--胃
RAW_DIR = Path(os.getenv("STOMACH_RAW_DIR", r"F:\公司知识平台\传入数据\原始文件"))
STORAGE_DIR = BASE_DIR / "storage"
KB_SERVER = os.getenv("KB_SERVER", "http://localhost:8080")
ADMIN_TOKEN = os.getenv("KB_ADMIN_TOKEN", "")

# 确保目录存在
for d in [STORAGE_DIR / "chunks", STORAGE_DIR / "wiki", STORAGE_DIR / "graph"]:
    d.mkdir(parents=True, exist_ok=True)

# ==================== 日志 ====================
LOG_FILE = BASE_DIR / "胃运行日志.log"
logger = logging.getLogger("stomach")
logger.setLevel(logging.INFO)
formatter = logging.Formatter("%(asctime)s [%(levelname)s] %(message)s", datefmt="%Y-%m-%d %H:%M:%S")

fh = RotatingFileHandler(LOG_FILE, maxBytes=10_000_000, backupCount=5, encoding="utf-8")
fh.setFormatter(formatter)
logger.addHandler(fh)

ch = logging.StreamHandler()
ch.setFormatter(formatter)
logger.addHandler(ch)

# ==================== 配置参数 ====================
MAX_FILE_MB = int(os.getenv("STOMACH_MAX_FILE_MB", "200"))
CHUNK_TARGET_TOKENS = int(os.getenv("STOMACH_CHUNK_TOKENS", "800"))
CHUNK_OVERLAP_PCT = float(os.getenv("STOMACH_CHUNK_OVERLAP", "0.15"))
DAEMON_INTERVAL = int(os.getenv("STOMACH_DAEMON_INTERVAL", "300"))  # 守护轮询间隔(秒)


# ╔══════════════════════════════════════════════════════════════════╗
# ║                    1. 解析器 (Parser Layer)                     ║
# ╚══════════════════════════════════════════════════════════════════╝

def parse_file(file_path: str) -> Dict:
    """解析 docx/pdf/xlsx/txt/md/json/cfg/html → 结构化文本 + 元数据 + 表格"""
    path = Path(file_path)
    ext = path.suffix.lower()
    fname = path.name

    if ext in ('.txt', '.log'):
        text = path.read_text(encoding='utf-8', errors='ignore')
        return {"text": text, "metadata": {"category": _classify_filename(fname), "type": "txt"}, "tables": [], "type": "text"}

    elif ext == '.md':
        text = path.read_text(encoding='utf-8', errors='ignore')
        return {"text": text, "metadata": {"category": _classify_filename(fname), "type": "md"}, "tables": [], "type": "markdown"}

    elif ext in ('.docx', '.doc'):
        return _parse_docx(file_path)

    elif ext == '.pdf':
        return _parse_pdf(file_path)

    elif ext in ('.xlsx', '.xls'):
        return _parse_excel(file_path)

    elif ext == '.csv':
        return _parse_csv(file_path)

    elif ext in ('.json', '.xml'):
        text = path.read_text(encoding='utf-8', errors='ignore')
        return {"text": text[:5000], "metadata": {"category": "配置数据", "type": ext}, "tables": [], "type": "data"}

    elif ext in ('.cfg', '.conf', '.ini', '.yaml', '.yml'):
        text = path.read_text(encoding='utf-8', errors='ignore')
        return {"text": text, "metadata": {"category": "网络配置", "type": "config"}, "tables": [], "type": "config"}

    elif ext in ('.html', '.htm'):
        return _parse_html(file_path)

    elif ext in ('.pptx', '.ppt'):
        return _parse_pptx(file_path)

    else:
        return {"text": "", "metadata": {"category": _classify_filename(fname), "type": ext}, "tables": [], "type": "unknown"}


def _parse_docx(file_path: str) -> Dict:
    """DOCX 解析 + pandoc 降级"""
    try:
        import docx
        doc = docx.Document(file_path)
        paras = []
        for p in doc.paragraphs:
            t = p.text.strip()
            if t:
                paras.append(t)

        # 提取表格
        tables_data = []
        for tbl in doc.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
            if rows and rows[0]:
                md = '| ' + ' | '.join(rows[0]) + ' |\n|' + '|'.join(['---'] * len(rows[0])) + '|\n'
                for row in rows[1:]:
                    md += '| ' + ' | '.join(row) + ' |\n'
                tables_data.append({"markdown": md, "rows": len(rows)})

        text = '\n\n'.join(paras)
        return {"text": text, "metadata": {"category": _classify_filename(Path(file_path).name), "type": "docx"},
                "tables": tables_data, "type": "document"}
    except Exception as e:
        logger.debug(f"python-docx failed, trying pandoc: {e}")
        try:
            import subprocess, tempfile
            with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as tmp:
                subprocess.run(['pandoc', file_path, '-t', 'plain', '-o', tmp.name], timeout=30)
                text = Path(tmp.name).read_text(encoding='utf-8', errors='ignore')
                os.unlink(tmp.name)
                return {"text": text, "metadata": {"category": _classify_filename(Path(file_path).name), "type": "docx"},
                        "tables": [], "type": "document"}
        except:
            return {"text": f"[无法解析: {file_path}]", "metadata": {"category": "未解析"}, "tables": [], "type": "error"}


def _parse_pdf(file_path: str) -> Dict:
    """PDF 解析：结构化 → 文字层 → OCR 降级 → pdftotext"""
    fname = Path(file_path).name

    # Layer 1: PyMuPDF 结构化
    try:
        import fitz
        doc = fitz.open(file_path)

        # 大 PDF 走结构化路径
        page_count = len(doc)
        if page_count > 50:
            blocks_list = []
            for page in doc:
                blocks = page.get_text("blocks")
                block_texts = [b[4].strip() for b in blocks if b[4].strip()]
                blocks_list.extend(block_texts)
            text = '\n\n'.join(blocks_list)
        else:
            text = '\n'.join(page.get_text() for page in doc)

        if text.strip() and len(text.strip()) > 50:
            return {"text": text, "metadata": {"category": _classify_filename(fname), "type": "pdf", "pages": page_count},
                    "tables": [], "type": "document"}
    except Exception as e:
        logger.debug(f"PyMuPDF failed: {e}")

    # Layer 2: pdfplumber 文字层提取
    try:
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            pages_text = []
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    pages_text.append(t)
            text = '\n\n'.join(pages_text)
            if text.strip() and len(text.strip()) > 50:
                return {"text": text, "metadata": {"category": _classify_filename(fname), "type": "pdf"},
                        "tables": [], "type": "document"}
    except Exception as e:
        logger.debug(f"pdfplumber failed: {e}")

    # Layer 3: OCR 降级（仅小文件）
    try:
        size_mb = Path(file_path).stat().st_size / (1024 * 1024)
        if size_mb < 20:
            from PIL import Image
            import fitz
            doc = fitz.open(file_path)
            if len(doc) <= 30:
                texts = []
                for page in doc:
                    pix = page.get_pixmap(dpi=200)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    # 简单的 pytesseract 调用
                    try:
                        import pytesseract
                        t = pytesseract.image_to_string(img, lang='chi_sim+eng')
                        if t.strip():
                            texts.append(t)
                    except:
                        pass
                text = '\n\n'.join(texts)
                if text.strip() and len(text.strip()) > 50:
                    return {"text": text, "metadata": {"category": _classify_filename(fname), "type": "pdf"},
                            "tables": [], "type": "document"}
    except Exception as e:
        logger.debug(f"OCR fallback failed: {e}")

    # Layer 4: pdftotext 命令行
    try:
        import subprocess, tempfile
        with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as tmp:
            subprocess.run(['pdftotext', '-layout', file_path, tmp.name], timeout=60)
            text = Path(tmp.name).read_text(encoding='utf-8', errors='ignore')
            os.unlink(tmp.name)
            if text.strip():
                return {"text": text, "metadata": {"category": _classify_filename(fname), "type": "pdf"},
                        "tables": [], "type": "document"}
    except Exception as e:
        logger.debug(f"pdftotext failed: {e}")

    return {"text": f"[无法解析: {file_path}]", "metadata": {"category": "未解析"}, "tables": [], "type": "error"}


def _parse_excel(file_path: str) -> Dict:
    """Excel 解析：xlsx/xls → pandas"""
    try:
        import pandas as pd
        dfs = pd.read_excel(file_path, sheet_name=None)
        tables = []
        all_text = []
        for sheet, df in dfs.items():
            df = df.fillna('')
            md = df.to_markdown(index=False) if hasattr(df, 'to_markdown') else df.to_string(index=False)
            tables.append({"markdown": md, "sheet": sheet, "rows": len(df)})
            all_text.append(f"### {sheet}\n{md}")
        return {"text": '\n\n'.join(all_text),
                "metadata": {"category": _classify_filename(Path(file_path).name), "type": "xlsx"},
                "tables": tables, "type": "spreadsheet"}
    except Exception as e:
        logger.warning(f"Excel parse failed: {e}")
        return {"text": f"[无法解析: {file_path}]", "metadata": {"category": "未解析"}, "tables": [], "type": "error"}


def _parse_csv(file_path: str) -> Dict:
    """CSV 解析"""
    try:
        import pandas as pd
        df = pd.read_csv(file_path, encoding='utf-8')
        md = df.to_markdown(index=False) if hasattr(df, 'to_markdown') else df.to_string(index=False)
        return {"text": md, "metadata": {"category": "表格数据", "type": "csv"},
                "tables": [{"markdown": md, "rows": len(df)}], "type": "spreadsheet"}
    except UnicodeDecodeError:
        try:
            import pandas as pd
            df = pd.read_csv(file_path, encoding='gbk')
            md = df.to_markdown(index=False) if hasattr(df, 'to_markdown') else df.to_string(index=False)
            return {"text": md, "metadata": {"category": "表格数据", "type": "csv"},
                    "tables": [{"markdown": md, "rows": len(df)}], "type": "spreadsheet"}
        except Exception as e:
            logger.warning(f"CSV parse failed: {e}")
            return {"text": f"[无法解析: {file_path}]", "metadata": {"category": "未解析"}, "tables": [], "type": "error"}


def _parse_html(file_path: str) -> Dict:
    """HTML 解析"""
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(Path(file_path).read_text(encoding='utf-8', errors='ignore'), 'html.parser')
        for tag in soup(['script', 'style', 'nav', 'footer']):
            tag.decompose()
        text = soup.get_text('\n')
        return {"text": text, "metadata": {"category": "网页文档", "type": "html"}, "tables": [], "type": "web"}
    except:
        text = Path(file_path).read_text(encoding='utf-8', errors='ignore')
        return {"text": text[:10000], "metadata": {"category": "网页文档", "type": "html"}, "tables": [], "type": "web"}


def _parse_pptx(file_path: str) -> Dict:
    """PPTX 解析"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides_text.append(f"--- 幻灯片 {i+1} ---\n" + '\n'.join(texts))
        text = '\n\n'.join(slides_text)
        return {"text": text, "metadata": {"category": "演示文稿", "type": "pptx"}, "tables": [], "type": "presentation"}
    except Exception as e:
        logger.warning(f"PPTX parse failed: {e}")
        return {"text": f"[无法解析: {file_path}]", "metadata": {"category": "未解析"}, "tables": [], "type": "error"}


# ╔══════════════════════════════════════════════════════════════════╗
# ║              2. Smart Chunk 语义分块器                           ║
# ╚══════════════════════════════════════════════════════════════════╝

def chunk_semantic(text: str, target_size: int = 800, overlap_pct: float = 0.15,
                   file_name: str = "", file_hash: str = "", category: str = "",
                   source: str = "") -> List[Dict]:
    """
    语义分块：按标题/段落边界拆分，控制每块大小
    优先级：H1/H2 标题 > 段落 > 固定长度
    """
    if not text or not text.strip():
        return []

    chunks = []
    # 中文标点
    text = text.replace('\r\n', '\n').replace('\r', '\n')

    # Step 1: 按标题分割
    title_pattern = re.compile(r'(#{1,4}\s+.+?)(?=\n#{1,4}\s+|\Z)', re.DOTALL)
    sections = title_pattern.findall(text)

    if not sections:
        # 没有标题，按空行分
        sections = [s.strip() for s in re.split(r'\n{2,}', text) if s.strip()]

    # Step 2: 合并短段落，拆分超长块
    overlap_size = max(50, int(target_size * overlap_pct))
    chunk_index = 0

    buffer = ""
    for section in sections:
        section = section.strip()
        if not section:
            continue

        if len(buffer) + len(section) < target_size * 1.5:
            buffer = (buffer + '\n\n' + section).strip()
        else:
            # 先保存 buffer
            if buffer:
                chunks.append({
                    "chunk_id": f"{file_hash}_c{chunk_index}",
                    "text": buffer,
                    "chunk_index": chunk_index,
                    "chunk_type": "text",
                    "category": category,
                    "file_name": file_name,
                    "file_hash": file_hash,
                    "source": source,
                })
                chunk_index += 1
                buffer = section
            else:
                # section 本身超长，按固定大小拆分
                sub_parts = _split_long_text(section, target_size, overlap_size)
                for sp in sub_parts:
                    chunks.append({
                        "chunk_id": f"{file_hash}_c{chunk_index}",
                        "text": sp,
                        "chunk_index": chunk_index,
                        "chunk_type": "text",
                        "category": category,
                        "file_name": file_name,
                        "file_hash": file_hash,
                        "source": source,
                    })
                    chunk_index += 1

    if buffer.strip():
        chunks.append({
            "chunk_id": f"{file_hash}_c{chunk_index}",
            "text": buffer,
            "chunk_index": chunk_index,
            "chunk_type": "text",
            "category": category,
            "file_name": file_name,
            "file_hash": file_hash,
            "source": source,
        })
        chunk_index += 1

    # 更新 total_chunks
    for c in chunks:
        c["total_chunks"] = len(chunks)

    return chunks


def _split_long_text(text: str, target_size: int, overlap: int) -> List[str]:
    """将超长文本按句子边界拆分"""
    parts = []
    current = ""
    # 按句子拆分
    sentences = re.split(r'(?<=[。！？.!?\n])', text)

    for sent in sentences:
        if len(current) + len(sent) < target_size:
            current += sent
        else:
            if current:
                parts.append(current)
            # 重叠：保留当前块的末尾
            if parts and overlap > 0:
                current = current[-overlap:] + sent if len(current) > overlap else sent
            else:
                current = sent

    if current:
        parts.append(current)

    return parts


# ╔══════════════════════════════════════════════════════════════════╗
# ║              3. 文本清洗 + 脱敏                                  ║
# ╚══════════════════════════════════════════════════════════════════╝

def clean_text(text: str) -> str:
    """清洗文本：去控制字符、统一空白、修复常见问题"""
    if not text:
        return ""
    # 去除 NULL 和控制字符
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
    # 统一换行
    text = text.replace('\r\n', '\n').replace('\r', '\n')
    # 压缩多余空行
    text = re.sub(r'\n{3,}', '\n\n', text)
    # 压缩多余空格
    text = re.sub(r' {2,}', ' ', text)
    # 去除首尾空白
    text = text.strip()
    return text


def desensitize(text: str) -> Tuple[str, List[Dict]]:
    """
    脱敏处理：遮掩身份证/手机号/银行卡/邮箱等敏感信息
    返回: (脱敏后文本, 脱敏统计)
    """
    if not text:
        return text, []

    stats = []

    # 手机号 13812345678 → 138****5678
    count_mobile = len(re.findall(r'\b1[3-9]\d{9}\b', text))
    text = re.sub(r'(1[3-9]\d)\d{4}(\d{4})', r'\1****\2', text)
    if count_mobile > 0:
        stats.append({"type": "mobile", "count": count_mobile})

    # 身份证 18位
    count_id = len(re.findall(r'\b\d{17}[\dXx]\b', text))
    text = re.sub(r'(\d{6})\d{8}(\d{4})', r'\1********\2', text)
    if count_id > 0:
        stats.append({"type": "id_card", "count": count_id})

    # 银行卡 16-19位
    count_bank = len(re.findall(r'\b\d{16,19}\b', text))
    text = re.sub(r'\b(\d{4})\d{8,11}(\d{4})\b', r'\1********\2', text)
    if count_bank > 0:
        stats.append({"type": "bank_card", "count": count_bank})

    # 邮箱
    count_email = len(re.findall(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}', text))
    text = re.sub(r'([a-zA-Z0-9._%+-]{3})[a-zA-Z0-9._%+-]*(@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,})',
                  r'\1***\2', text)
    if count_email > 0:
        stats.append({"type": "email", "count": count_email})

    return text, stats


# ╔══════════════════════════════════════════════════════════════════╗
# ║              4. 17 类智能分类引擎                                ║
# ╚══════════════════════════════════════════════════════════════════╝

def classify_document(file_path: Path, raw_text: str = "") -> str:
    """
    三要素分类：关键词(最高) > 扩展名 > 兜底
    返回 17 个标准分类之一
    """
    fname = file_path.name
    ext = file_path.suffix.lower()

    combined = (fname + " " + raw_text[:8000]).lower()

    # ======== 第一层：关键词匹配 ========
    keyword_rules = [
        (["模具", "注塑", "成型", "型腔", "分模", "滑块", "顶出", "模架", "模仁", "浇口"], "模具设计"),
        (["连接器", "connector", "端子", "针座", "接插件", "线束", "fpc", "ffc"], "连接器设计"),
        (["机械设计", "机械结构", "轴承", "齿轮", "联轴器", "减速器", "传动",
          "丝杆", "导轨", "气缸", "液压", "气动", "电磁阀"], "机械设计"),
        (["标准件", "紧固件", "gb/t", "螺钉", "螺栓", "螺母", "垫圈", "卡簧",
          "挡圈", "销", "键", "弹簧", "o型圈", "misumi"], "标准件库"),
        (["plc", "传感器", "伺服", "步进", "变频", "hmi", "继电器",
          "断路器", "接触器", "接线端子", "配电", "西门子", "三菱", "欧姆龙"], "电气自动化"),
        (["产线", "自动化线", "装配线", "流水线", "皮带线", "辊筒", "倍速链",
          "机械手", "机器人", "agv", "输送"], "自动化产线"),
        (["vlan", "路由", "交换机", "dhcp", "dns", "子网", "拓扑", "acl",
          "防火墙", "ip地址", "网段", "端口映射", "vpn", "ap", "无线", "stp",
          "802", "trunk", "nps", "域控"], "网络配置"),
        (["服务器", "运维", "部署", "linux", "ubuntu", "centos", "windows server",
          "备份", "docker", "k8s", "kubernetes", "nginx", "apache"], "服务器运维"),
        (["ecology", "泛微", "oa", "门户引擎", "流程引擎", "组织权限", "公文管理",
          "资产管理", "人事管理", "建模"], "OA系统"),
        (["规范", "标准", "gb", "jb", "qc", "技术要求", "验收", "规程",
          "astm", "asme", "iso", "din", "jis", "规格书"], "工程技术规范"),
        (["品质", "质量", "spc", "fmea", "cpk", "良率", "不良",
          "检验", "测量", "卡尺", "三次元", "三坐标", "gd&t"], "品质检测"),
        (["供应商", "采购", "供应链", "询价", "比价", "po"], "供应商管理"),
        (["制度", "规定", "办法", "流程", "通知", "公告", "考勤",
          "薪资", "人事", "入职", "离职", "年假"], "公司制度"),
        (["行政", "后勤", "车辆", "接待", "前台"], "行政人事"),
        (["财务", "成本", "报价", "合同", "预算", "费用", "发票",
          "付款", "收款", "报销", "税务", "审计", "购销"], "财务文档"),
        (["项目管理", "甘特图", "里程碑", "wbs", "项目计划", "进度",
          "交付", "验收", "bom"], "项目管理"),
        (["手册", "指南", "教程", "说明", "安装", "维护", "操作",
          "编程", "调试", "api", "sdk", "代码", "配置"], "技术文档"),
        (["会议", "周报", "月报", "总结", "报告", "申请"], "办公文档"),
    ]

    for keywords, category in keyword_rules:
        for kw in keywords:
            if kw.lower() in combined:
                return category

    # ======== 第二层：扩展名映射 ========
    ext_map = {
        ".dwg": "模具设计", ".dxf": "模具设计",
        ".stp": "机械设计", ".step": "机械设计",
        ".ipt": "机械设计", ".iam": "机械设计",
        ".sldprt": "机械设计", ".sldasm": "机械设计",
        ".prt": "机械设计", ".iges": "机械设计", ".igs": "机械设计",
        # 元数据
        ".zip": "元数据", ".rar": "元数据", ".7z": "元数据",
        ".tar": "元数据", ".gz": "元数据", ".iso": "元数据",
        ".exe": "元数据", ".dll": "元数据", ".sys": "元数据",
        ".bin": "元数据", ".msi": "元数据",
        ".mp4": "元数据", ".avi": "元数据", ".mov": "元数据",
        ".mp3": "元数据", ".wav": "元数据",
        ".psd": "元数据", ".ai": "元数据",
    }
    if ext in ext_map:
        return ext_map[ext]

    # ======== 第三层：兜底 ========
    if ext in ('.pdf', '.docx', '.doc'):
        return "技术文档"
    elif ext in ('.xlsx', '.xls', '.csv'):
        return "办公文档"
    elif ext in ('.pptx', '.ppt'):
        return "演示文稿"
    elif ext in ('.html', '.htm'):
        return "网页文档"
    elif ext in ('.txt', '.md'):
        return "通用文档"
    else:
        return "其他文档"


def _classify_filename(filename: str) -> str:
    """快速文件名分类（解析阶段用）"""
    name = filename.lower()
    for kw in ['网络', 'vlan', 'ip', '交换机', '路由', 'dhcp', 'stp', 'acl', '拓扑', '端口', 'cfg']:
        if kw in name:
            return '网络配置'
    for kw in ['服务器', '运维', '部署', '备份', 'linux']:
        if kw in name:
            return '服务器运维'
    for kw in ['oa', 'ecology', '泛微', '流程', '门户', '组织']:
        if kw in name:
            return 'OA系统'
    for kw in ['模具', '注塑', '成型', '导柱']:
        if kw in name:
            return '模具设计'
    for kw in ['制度', '规范', '流程', '操作']:
        if kw in name:
            return '公司制度'
    for kw in ['合同', '报价', '采购']:
        if kw in name:
            return '财务文档'
    return '通用文档'


# ╔══════════════════════════════════════════════════════════════════╗
# ║              5. Wiki 蒸馏器                                      ║
# ╚══════════════════════════════════════════════════════════════════╝

def distill_wiki(chunks: List[Dict], file_name: str, category: str, raw_text: str = "") -> List[Dict]:
    """从 chunks 提炼 Wiki 摘要条目"""
    wiki_entries = []

    if not chunks:
        return wiki_entries

    # 合并前 10 个 chunk 的前 300 字符作为摘要素材
    full_text = raw_text if raw_text else ' '.join(c.get('text', '')[:300] for c in chunks[:15])
    title = Path(file_name).stem.replace('(', '').replace(')', '').replace('_', ' ').strip()

    # 提取关键信息
    summary = full_text[:400].strip()
    if len(full_text) > 400:
        # 断句在标点处
        cut = max(summary.rfind('。'), summary.rfind('；'), summary.rfind('，'), summary.rfind('. '))
        if cut > 200:
            summary = summary[:cut + 1]
        else:
            summary = summary[:400] + '…'

    wiki_entries.append({
        'id': hashlib.md5(file_name.encode()).hexdigest()[:12],
        'title': title,
        'category': category,
        'summary': summary,
        'content': full_text[:5000],  # 全文前 5000 字符
        'tags': ','.join([category, Path(file_name).suffix.replace('.', '')]),
        'quality_score': _estimate_quality(full_text),
        'chunk_count': len(chunks),
        'file_name': file_name,
    })

    return wiki_entries


def _estimate_quality(text: str) -> float:
    """评估内容质量 (0-1)"""
    score = 0.3
    if len(text) > 500: score += 0.1
    if len(text) > 2000: score += 0.1
    if len(text) > 5000: score += 0.1
    if re.search(r'\d+', text): score += 0.1  # 含数字=数据型文档
    if re.search(r'[，。；：、]', text): score += 0.1  # 中文标点
    if re.search(r'(#|##|###)', text): score += 0.05  # 含标题结构
    if len(text) > 100: score += 0.05
    return min(score, 1.0)


# ╔══════════════════════════════════════════════════════════════════╗
# ║              6. 实体关系提取                                      ║
# ╚══════════════════════════════════════════════════════════════════╝

def extract_entities(chunks: List[Dict], file_name: str, category: str = "") -> Dict:
    """从 chunks 提取实体 → 知识图谱节点"""
    entities = {}
    relations = []
    full_text = ' '.join(c.get('text', '') for c in chunks)

    # IP 地址
    ip_pattern = r'\b(?:[0-9]{1,3}\.){3}[0-9]{1,3}(?:/\d{1,2})?\b'
    for ip in set(re.findall(ip_pattern, full_text)):
        entities[ip] = {'name': ip, 'type': 'ip_address', 'category': '网络'}

    # 设备型号 (LSW/AR/ES/RG-...)
    device_pattern = r'\b(?:LSW\d+|AR\d+|ES\d+[A-Z]*|RG-[A-Z0-9-]+|S\d{4}[A-Z]*|WS-C\d{4}[A-Z]*)\b'
    for dev in set(re.findall(device_pattern, full_text)):
        entities[dev] = {'name': dev, 'type': 'device', 'category': '网络设备'}

    # 材料牌号
    material_pattern = r'\b(?:PA\d{1,2}|POM|PC|ABS|PPS|PBT|PEI|PEEK|SUJ\d|S\d{3}|SKD\d{2})\b'
    for mat in set(re.findall(material_pattern, full_text)):
        entities[mat] = {'name': mat, 'type': 'material', 'category': '材料'}

    # 标准号 (GB/JB/ISO/DIN/...)
    standard_pattern = r'\b(?:GB/T?\s*\d+[\.\-]?\d*|JB/T?\s*\d+[\.\-]?\d*|ISO\s*\d+[\.\-]?\d*|DIN\s*\d+[\.\-]?\d*|JIS\s*\w+\s*\d+)\b'
    for std in set(re.findall(standard_pattern, full_text)):
        entities[std] = {'name': std, 'type': 'standard', 'category': '标准规范'}

    # 供应商/品牌
    vendor_keywords = ['MISUMI', '三菱', 'Mitsubishi', '欧姆龙', 'Omron', '西门子', 'Siemens',
                       '基恩士', 'Keyence', 'SMC', 'THK', 'NSK', 'HIWIN', '上银', '亚德客', 'AirTAC']
    for kw in vendor_keywords:
        if kw.lower() in full_text.lower():
            entities[kw] = {'name': kw, 'type': 'vendor', 'category': '供应商'}

    return {
        'nodes': entities,
        'edges': relations,
        'source': file_name,
        'category': category,
    }


# ╔══════════════════════════════════════════════════════════════════╗
# ║              7. 推送 kb-server                                   ║
# ╚
# ╔══════════════════════════════════════════════════════════════════╗
# ║              7. 推送 kb-server                                   ║
# ╚══════════════════════════════════════════════════════════════════╝

def push_to_server(chunks: List[Dict], wiki: List[Dict], entities: Dict,
                   file_name: str, file_hash: str, category: str,
                   md_path: str = "") -> Dict:
    """推送到 kb-server，3 次重试"""
    import requests
    results = {'chunks': 0, 'wiki': 0, 'graph': 0}

    headers = {
        'Content-Type': 'application/json',
        'X-Admin-Token': ADMIN_TOKEN
    }

    # 推送 chunks → /api/ingest-batch (脾)
    if chunks:
        payload = {
            "file_name": file_name,
            "file_hash": file_hash,
            "category": category,
            "chunks": chunks,
            "md_path": md_path,
        }
        for attempt in range(3):
            try:
                resp = requests.post(
                    f"{KB_SERVER}/api/ingest-batch",
                    json=payload,
                    timeout=30,
                    headers=headers
                )
                if resp.status_code == 200:
                    results['chunks'] = len(chunks)
                    logger.info(f"  Push chunks OK ({len(chunks)} chunks)")
                    break
                elif resp.status_code == 422:
                    logger.warning(f"  Push chunks 422 (format): {resp.text[:150]}")
                    break  # 格式错误不重试
                else:
                    if attempt < 2:
                        logger.warning(f"  Push chunks {resp.status_code}, retry {attempt+2}/3...")
                        time.sleep(2 ** attempt)
                    else:
                        logger.error(f"  Push chunks FAILED: {resp.status_code} {resp.text[:200]}")
            except Exception as e:
                if attempt < 2:
                    logger.warning(f"  Push chunks error: {e}, retry {attempt+2}/3...")
                    time.sleep(2 ** attempt)
                else:
                    logger.error(f"  Push chunks FINAL FAIL: {e}")

    # 推送 Wiki → /api/wiki/page (肺)
    if wiki:
        for w in wiki:
            for attempt in range(2):
                try:
                    resp = requests.post(f"{KB_SERVER}/api/wiki/page", json=w, timeout=15, headers=headers)
                    if resp.status_code == 200:
                        results['wiki'] += 1
                        break
                except:
                    if attempt == 1:
                        logger.warning(f"  Push wiki failed for {w.get('title','')}")
        if results['wiki'] > 0:
            logger.info(f"  Push wiki OK ({results['wiki']} entries)")

    # 推送实体 → /api/graph/build (肝)
    if entities.get('nodes'):
        payload = {
            "entities": entities['nodes'],
            "edges": entities.get('edges', []),
            "source": entities.get('source', ''),
            "category": entities.get('category', ''),
        }
        for attempt in range(2):
            try:
                resp = requests.post(f"{KB_SERVER}/api/graph/build", json=payload, timeout=15, headers=headers)
                if resp.status_code == 200:
                    results['graph'] = len(entities['nodes'])
                    logger.info(f"  Push graph OK ({results['graph']} nodes)")
                    break
            except:
                if attempt == 1:
                    logger.warning(f"  Push graph failed")

    return results


# ╔══════════════════════════════════════════════════════════════════╗
# ║              8. 去重引擎                                         ║
# ╚══════════════════════════════════════════════════════════════════╝

class DedupEngine:
    """MD5 去重 + 日志持久化"""

    def __init__(self, log_path: Path):
        self.log_path = log_path
        self.processed: set = set()
        self._load()

    def _load(self):
        if self.log_path.exists():
            try:
                data = json.loads(self.log_path.read_text(encoding='utf-8'))
                self.processed = set(data.get('processed', []))
                logger.info(f"[Dedup] Loaded {len(self.processed)} hashes")
            except:
                pass

    def _save(self):
        self.log_path.parent.mkdir(parents=True, exist_ok=True)
        data = {
            'processed': list(self.processed),
            'last_save': datetime.now().isoformat(),
            'count': len(self.processed),
        }
        tmp = self.log_path.with_suffix('.tmp')
        tmp.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding='utf-8')
        tmp.replace(self.log_path)

    def is_processed(self, file_hash: str) -> bool:
        return file_hash in self.processed

    def mark_processed(self, file_hash: str):
        self.processed.add(file_hash)

    def save(self):
        self._save()


# ╔══════════════════════════════════════════════════════════════════╗
# ║              9. 主消化流程                                       ║
# ╚══════════════════════════════════════════════════════════════════╝

def digest_file(file_path: str, dedup: DedupEngine = None, push: bool = True) -> Optional[Dict]:
    """
    完整消化一个文件:
      解析 → 清洗 → 脱敏 → 语义分块 → 分类 → 蒸馏 → 实体提取 → 推送
    """
    t0 = time.time()
    path = Path(file_path)
    file_name = path.name

    # 跳过临时文件
    if file_name.startswith('~$') or file_name.startswith('.'):
        return None

    logger.info(f'[Stomach] Digesting: {file_name}')

    try:
        # 0. 快速哈希 + 去重
        fsize = path.stat().st_size
        file_hash = _quick_hash(file_path, fsize)

        if dedup and dedup.is_processed(file_hash):
            logger.info(f'  Already processed, skip')
            return {'ok': True, 'file_name': file_name, 'skipped': True, 'reason': 'dedup'}

        # 大小检查
        size_mb = fsize / (1024 * 1024)
        if size_mb > MAX_FILE_MB:
            logger.warning(f'  Too large ({size_mb:.1f}MB > {MAX_FILE_MB}MB), skip')
            return {'ok': False, 'file_name': file_name, 'error': f'file too large ({size_mb:.1f}MB)'}

        # 1. 解析
        parsed = parse_file(file_path)
        raw_text = parsed.get('text', '')
        tables = parsed.get('tables', [])
        char_count = len(raw_text)
        logger.info(f'  Parsed: {char_count} chars, {len(tables)} tables')

        if not raw_text or char_count < 10:
            # 空文件或极短内容 → 元数据模式
            raw_text = f"[元数据/空文件] {file_name}\n大小: {size_mb:.1f}MB\n格式: {path.suffix}"
            logger.info(f'  Empty/minimal content, using metadata mode')

        # 2. 分类
        category = classify_document(path, raw_text)
        logger.info(f'  Classified: {category}')

        # 3. 清洗
        cleaned = clean_text(raw_text)
        logger.info(f'  Cleaned: {len(cleaned)} chars')

        # 4. 脱敏
        cleaned_safe, desens_stats = desensitize(cleaned)
        if desens_stats:
            logger.info(f'  Desensitized: {sum(s["count"] for s in desens_stats)} items')
        else:
            logger.info(f'  Desensitized: none found')

        # 5. 语义分块
        chunks = chunk_semantic(
            cleaned_safe,
            target_size=CHUNK_TARGET_TOKENS,
            overlap_pct=CHUNK_OVERLAP_PCT,
            file_name=file_name,
            file_hash=file_hash,
            category=category,
            source=file_path,
        )
        logger.info(f'  Chunked: {len(chunks)} semantic chunks')

        # 表格分块 (追加)
        tbl_idx = 0
        for tbl in tables:
            md = tbl.get('markdown', '')
            if md.strip():
                chunks.append({
                    "chunk_id": f"{file_hash}_t{tbl_idx}",
                    "text": md,
                    "chunk_index": tbl_idx,
                    "chunk_type": "table",
                    "category": category,
                    "sheet_name": tbl.get("sheet", ""),
                    "file_name": file_name,
                    "file_hash": file_hash,
                    "source": file_path,
                })
                tbl_idx += 1

        if not chunks:
            logger.warning(f'  No chunks produced, skip')
            return {'ok': False, 'file_name': file_name, 'error': 'no chunks'}

        # 6. Wiki 蒸馏
        wiki = distill_wiki(chunks, file_name, category, cleaned_safe)
        logger.info(f'  Distilled: {len(wiki)} wiki entries')

        # 7. 实体提取
        entities = extract_entities(chunks, file_name, category)
        logger.info(f'  Entities: {len(entities.get("nodes", {}))} nodes')

        # 8. 本地存储
        _save_local(chunks, wiki, entities, file_hash, file_name)

        # 9. 推送 kb-server
        push_result = {}
        if push:
            push_result = push_to_server(chunks, wiki, entities, file_name, file_hash, category)
            logger.info(f'  Push: chunks={push_result.get("chunks",0)}, wiki={push_result.get("wiki",0)}, graph={push_result.get("graph",0)}')

        # 10. 标记已处理
        if dedup:
            dedup.mark_processed(file_hash)

        elapsed = round((time.time() - t0) * 1000, 1)
        logger.info(f'  DONE {file_name} in {elapsed}ms [{category}]')

        return {
            'ok': True,
            'file_name': file_name,
            'file_hash': file_hash,
            'category': category,
            'chunks': len(chunks),
            'wiki': len(wiki),
            'entities': len(entities.get('nodes', {})),
            'size_mb': round(size_mb, 2),
            'push': push_result,
            'elapsed_ms': elapsed,
        }

    except Exception as e:
        import traceback
        logger.error(f'  FAILED {file_name}: {e}\n{traceback.format_exc()[-500:]}')
        return {'ok': False, 'file_name': file_name, 'error': str(e)}


def _quick_hash(file_path: str, fsize: int) -> str:
    """快速 SHA256（仅读前 1MB + 文件大小）"""
    with open(file_path, 'rb') as f:
        head = f.read(min(fsize, 1_000_000))
    return hashlib.sha256(head + str(fsize).encode()).hexdigest()[:16]


def _save_local(chunks: List[Dict], wiki: List[Dict], entities: Dict,
                file_hash: str, file_name: str):
    """保存消化产物到本地 storage"""
    # chunks
    cpath = STORAGE_DIR / 'chunks' / f'{file_hash}.json'
    _atomic_write_json(cpath, {'file_name': file_name, 'hash': file_hash, 'chunks': chunks})
    # wiki
    wpath = STORAGE_DIR / 'wiki' / f'{file_hash}.json'
    _atomic_write_json(wpath, {'file_name': file_name, 'hash': file_hash, 'wiki': wiki})
    # graph
    gpath = STORAGE_DIR / 'graph' / f'{file_hash}.json'
    _atomic_write_json(gpath, {'file_name': file_name, 'hash': file_hash, 'entities': entities})


def _atomic_write_json(path: Path, data):
    tmp = path.with_suffix('.tmp')
    tmp.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding='utf-8')
    tmp.replace(path)


# ╔══════════════════════════════════════════════════════════════════╗
# ║              10. 守护模式：自动消化新文件                          ║
# ╚══════════════════════════════════════════════════════════════════╝

def scan_raw_files(raw_dir: Path) -> List[Path]:
    """扫描原始文件目录，返回待处理文件列表"""
    if not raw_dir.exists():
        logger.warning(f"[Scan] 原始文件目录不存在: {raw_dir}")
        return []

    all_files = list(raw_dir.rglob('*'))
    files = [f for f in all_files if f.is_file()
             and not f.name.startswith('~$')
             and not f.name.startswith('.')
             and '.tmp' not in f.suffix
             and 'thumbs' not in str(f)
             and '.images' not in str(f)]

    # 仅处理可消化格式
    digestable_exts = {
        '.docx', '.doc', '.pdf', '.xlsx', '.xls', '.csv',
        '.txt', '.md', '.json', '.xml', '.cfg', '.conf',
        '.ini', '.yaml', '.yml', '.html', '.htm', '.pptx', '.ppt',
        '.log',
    }
    files = [f for f in files if f.suffix.lower() in digestable_exts]

    return files


def daemon_loop(raw_dir: Path = None, dedup_log: Path = None):
    """守护循环：定期扫描新文件，自动消化"""
    if raw_dir is None:
        raw_dir = RAW_DIR
    if dedup_log is None:
        dedup_log = STORAGE_DIR / 'dedup_log.json'

    logger.info("=" * 60)
    logger.info("  中宫--胃 消化中枢 v5.0 启动")
    logger.info(f"  原始文件: {raw_dir}")
    logger.info(f"  存储路径: {STORAGE_DIR}")
    logger.info(f"  推送目标: {KB_SERVER}")
    logger.info(f"  守护间隔: {DAEMON_INTERVAL}s")
    logger.info(f"  参数: MAX_FILE={MAX_FILE_MB}MB, CHUNK={CHUNK_TARGET_TOKENS}, OVERLAP={CHUNK_OVERLAP_PCT}")
    logger.info("=" * 60)

    dedup = DedupEngine(dedup_log)

    total_processed = 0
    total_skipped = 0
    total_failed = 0

    try:
        while True:
            cycle_start = time.time()
            files = scan_raw_files(raw_dir)

            new_files = 0
            skipped = 0
            failed = 0

            for fp in files:
                file_hash = None
                try:
                    fsize = fp.stat().st_size
                    file_hash = _quick_hash(str(fp), fsize)
                except:
                    pass

                if file_hash and dedup.is_processed(file_hash):
                    skipped += 1
                    continue

                result = digest_file(str(fp), dedup=dedup, push=True)
                if result:
                    if result.get('skipped'):
                        skipped += 1
                    elif result.get('ok'):
                        new_files += 1
                    else:
                        failed += 1

            # 保存去重日志（每轮）
            dedup.save()

            total_processed += new_files
            total_skipped += skipped
            total_failed += failed

            cycle_time = round(time.time() - cycle_start, 1)
            if new_files > 0 or failed > 0:
                logger.info(f"[Daemon] Cycle done in {cycle_time}s: "
                            f"new={new_files}, skip={skipped}, fail={failed} "
                            f"(total: proc={total_processed}, skip={total_skipped}, fail={total_failed})")
            else:
                logger.debug(f"[Daemon] Idle cycle ({cycle_time}s)")

            time.sleep(DAEMON_INTERVAL)

    except KeyboardInterrupt:
        logger.info("[Daemon] Received stop signal")
        dedup.save()
        logger.info(f"Final stats: processed={total_processed}, skipped={total_skipped}, failed={total_failed}")


# ╔══════════════════════════════════════════════════════════════════╗
# ║              11. 一次性批量消化                                   ║
# ╚══════════════════════════════════════════════════════════════════╝

def batch_digest(directory: str = None, push: bool = True,
                 dedup_log: Path = None, reset_dedup: bool = False):
    """一次性消化目录中所有文件"""
    if directory is None:
        directory = str(RAW_DIR)

    raw_dir = Path(directory)
    if dedup_log is None:
        dedup_log = STORAGE_DIR / 'dedup_log.json'

    if reset_dedup and dedup_log.exists():
        dedup_log.unlink()
        logger.info("[Batch] 已重置去重日志")

    dedup = DedupEngine(dedup_log)
    files = scan_raw_files(raw_dir)

    logger.info(f"[Batch] 扫描到 {len(files)} 个文件，开始消化...")
    logger.info("=" * 60)

    ok_count = 0
    skip_count = 0
    fail_count = 0
    results = []

    for i, fp in enumerate(files):
        result = digest_file(str(fp), dedup=dedup, push=push)

        if result:
            if result.get('skipped'):
                skip_count += 1
            elif result.get('ok'):
                ok_count += 1
            else:
                fail_count += 1

        results.append(result)

        # 每 50 个文件保存去重 + 打印进度
        if (i + 1) % 50 == 0:
            dedup.save()
            logger.info(f"[Batch] Progress: {i+1}/{len(files)} (ok={ok_count}, skip={skip_count}, fail={fail_count})")

    dedup.save()

    logger.info("=" * 60)
    logger.info(f"[Batch] 消化完成: ok={ok_count}, skip={skip_count}, fail={fail_count}, total={len(files)}")

    return results


# ╔══════════════════════════════════════════════════════════════════╗
# ║              12. 入口                                            ║
# ╚══════════════════════════════════════════════════════════════════╝

if __name__ == '__main__':
    import argparse

    parser = argparse.ArgumentParser(description='中宫--胃 消化中枢 v5.0')
    parser.add_argument('--daemon', '-d', action='store_true', help='启动守护模式（持续监控新文件）')
    parser.add_argument('--batch', '-b', action='store_true', help='一次性消化所有文件')
    parser.add_argument('--file', '-f', type=str, help='消化单个文件')
    parser.add_argument('--dir', type=str, help='指定原始文件目录')
    parser.add_argument('--reset-dedup', action='store_true', help='重置去重记录')
    parser.add_argument('--no-push', action='store_true', help='仅本地消化，不推送服务器')
    args = parser.parse_args()

    if args.daemon:
        # 守护模式
        raw_dir = Path(args.dir) if args.dir else RAW_DIR
        daemon_loop(raw_dir=raw_dir)

    elif args.file:
        # 单文件模式
        dedup = DedupEngine(STORAGE_DIR / 'dedup_log.json')
        result = digest_file(args.file, dedup=dedup, push=not args.no_push)
        if result:
            status = 'OK' if result.get('ok') else 'FAIL'
            if result.get('skipped'):
                status = 'SKIPPED'
            print(f"[{status}] {result.get('file_name')}: {json.dumps({k:v for k,v in result.items() if k != 'chunks'}, ensure_ascii=False)}")
        dedup.save()

    elif args.batch:
        # 批量模式
        batch_digest(directory=args.dir, push=not args.no_push,
                     reset_dedup=args.reset_dedup)

    else:
        # 默认：一次性消化
        print("中宫--胃 消化中枢 v5.0")
        print(f"默认模式：一次性消化原始文件（{RAW_DIR}）")
        print(f"守护模式：python stomach.py --daemon")
        print(f"单文件：  python stomach.py --file <文件路径>")
        print()
        batch_digest(push=not args.no_push)
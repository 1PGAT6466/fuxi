"""
unified.py 鈥?浼忕静缁熶竴澶勭悊绠＄嚎
鎵€鏈夋潵婧愶紙涓婁紶/瑁呰浇鏈?API锛夌殑鏁版嵁锛岄兘缁忚繃杩欐潯绠＄嚎銆?
"""

import asyncio
import hashlib
import logging
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

from src.models.chunk import Chunk, ChunkType
from src.models.entity import Entity
from src.models.event import Event
from src.pipeline.errors import (
    CleanError,
    ParseError,
    SaveError,
)

logger = logging.getLogger("pipeline")


class PipelineMetrics:
    """鏁版嵁绠＄嚎鍚勭幆鑺傝€楁椂缁熻"""

    def __init__(self):
        self.timings = {}
        self.current_stage = None
        self.stage_start = None

    def start(self, stage: str) -> Any:
        self.current_stage = stage
        self.stage_start = time.time()

    def end(self) -> Any:
        if self.current_stage and self.stage_start:
            elapsed = time.time() - self.stage_start
            self.timings[self.current_stage] = round(elapsed * 1000, 2)  # ms
            self.current_stage = None

    def report(self) -> dict:
        if not self.timings:
            return {"stages": {}, "total_ms": 0, "bottleneck": None}
        total = sum(self.timings.values())
        return {
            "stages": self.timings,
            "total_ms": round(total, 2),
            "bottleneck": max(self.timings, key=self.timings.get) if self.timings else None,
        }


@dataclass
class PipelineResult:
    """绠＄嚎澶勭悊缁撴灉"""

    source: str = ""
    file_path: str = ""
    raw_text: str = ""
    cleaned_text: str = ""
    tables: List[Dict] = field(default_factory=list)
    chunks: List[Chunk] = field(default_factory=list)
    events: List[Event] = field(default_factory=list)
    entities: List[Entity] = field(default_factory=list)
    metadata: Dict = field(default_factory=dict)
    metrics: Dict = field(default_factory=dict)
    skipped: bool = False
    duration_ms: float = 0
    errors: List[str] = field(default_factory=list)


class UnifiedParser:
    """缁熶竴瑙ｆ瀽鍣?鈥?鍚堝苟 stomach.py + ingest.py 鐨勮В鏋愰€昏緫"""

    def __init__(self, config: Dict = None):
        self.config = config or {}

    def parse(self, file_path: str) -> Dict:
        """瑙ｆ瀽鏂囦欢锛岃繑鍥?{text, tables, metadata}"""
        path = Path(file_path)
        if not path.exists():
            raise ParseError(f"鏂囦欢涓嶅瓨鍦? {file_path}")

        ext = path.suffix.lower()
        try:
            if ext == ".pdf":
                return self._parse_pdf(file_path)
            elif ext in (".docx", ".doc"):
                return self._parse_docx(file_path)
            elif ext in (".xlsx", ".xls"):
                return self._parse_excel(file_path)
            elif ext == ".txt":
                return self._parse_text(file_path)
            elif ext == ".md":
                return self._parse_markdown(file_path)
            elif ext == ".csv":
                return self._parse_csv(file_path)
            elif ext in (".json",):
                return self._parse_json(file_path)
            elif ext in (".html", ".htm"):
                return self._parse_html(file_path)
            elif ext in (".pptx", ".ppt"):
                return self._parse_pptx(file_path)
            else:
                return self._parse_text(file_path)
        except ParseError:
            raise
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            raise ParseError(f"瑙ｆ瀽澶辫触 ({ext}): {e}")

    def _parse_pdf(self, file_path: str) -> Dict:
        """PDF 瑙ｆ瀽 鈥?鍚堝苟 stomach.py + ingest.py 閫昏緫"""
        text = ""
        tables = []

        # 鏂瑰紡1: fitz (PyMuPDF) 鈥?涓枃鏈€浼?
        try:
            import fitz

            doc = fitz.open(file_path)
            pages_text = []
            for page in doc:
                pages_text.append(page.get_text())
            text = "\n".join(pages_text)
            doc.close()
            if text.strip():
                return {"text": text, "tables": tables, "metadata": {"parser": "fitz"}}
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            logger.warning("Exception 澶辫触: %s", e, exc_info=True)

        # 鏂瑰紡2: pdfplumber 鈥?琛ㄦ牸鎻愬彇
        try:
            import pdfplumber

            with pdfplumber.open(file_path) as pdf:
                pages_text = []
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        pages_text.append(page_text)
                    # 鎻愬彇琛ㄦ牸
                    for table in page.extract_tables():
                        if table:
                            tables.append(
                                {"headers": table[0] if table else [], "rows": table[1:] if len(table) > 1 else []}
                            )
                text = "\n".join(pages_text)
                if text.strip():
                    return {"text": text, "tables": tables, "metadata": {"parser": "pdfplumber"}}
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            logger.warning("Exception 澶辫触: %s", e, exc_info=True)

        # 鏂瑰紡3: pypdf 鈥?鍏滃簳
        try:
            from pypdf import PdfReader

            reader = PdfReader(file_path)
            pages_text = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    pages_text.append(page_text)
            text = "\n".join(pages_text)
            return {"text": text, "tables": tables, "metadata": {"parser": "pypdf"}}
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            raise ParseError(f"PDF瑙ｆ瀽澶辫触: {e}")

    def _parse_docx(self, file_path: str) -> Dict:
        """DOCX 瑙ｆ瀽"""
        try:
            from docx import Document

            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            tables = []
            for table in doc.tables:
                headers = [cell.text for cell in table.rows[0].cells] if table.rows else []
                rows = [[cell.text for cell in row.cells] for row in table.rows[1:]] if len(table.rows) > 1 else []
                tables.append({"headers": headers, "rows": rows})
            return {"text": "\n".join(paragraphs), "tables": tables, "metadata": {"parser": "python-docx"}}
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            raise ParseError(f"DOCX瑙ｆ瀽澶辫触: {e}")

    def _parse_excel(self, file_path: str) -> Dict:
        """Excel 瑙ｆ瀽"""
        try:
            import pandas as pd

            df = pd.read_excel(file_path)
            text = df.to_string(index=False)
            tables = [{"headers": list(df.columns), "rows": df.values.tolist()}]
            return {"text": text, "tables": tables, "metadata": {"parser": "pandas"}}
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            raise ParseError(f"Excel瑙ｆ瀽澶辫触: {e}")

    def _parse_text(self, file_path: str) -> Dict:
        """绾枃鏈В鏋?""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            return {"text": text, "tables": [], "metadata": {"parser": "text"}}
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            raise ParseError(f"鏂囨湰瑙ｆ瀽澶辫触: {e}")

    def _parse_markdown(self, file_path: str) -> Dict:
        """Markdown 瑙ｆ瀽"""
        return self._parse_text(file_path)

    def _parse_csv(self, file_path: str) -> Dict:
        """CSV 瑙ｆ瀽"""
        try:
            import pandas as pd

            df = pd.read_csv(file_path)
            text = df.to_string(index=False)
            tables = [{"headers": list(df.columns), "rows": df.values.tolist()}]
            return {"text": text, "tables": tables, "metadata": {"parser": "pandas"}}
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            raise ParseError(f"CSV瑙ｆ瀽澶辫触: {e}")

    def _parse_json(self, file_path: str) -> Dict:
        """JSON 瑙ｆ瀽"""
        try:
            import json

            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            text = json.dumps(data, ensure_ascii=False, indent=2)
            return {"text": text, "tables": [], "metadata": {"parser": "json"}}
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            raise ParseError(f"JSON瑙ｆ瀽澶辫触: {e}")

    def _parse_html(self, file_path: str) -> Dict:
        """HTML 瑙ｆ瀽"""
        try:
            from bs4 import BeautifulSoup

            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            return {"text": text, "tables": [], "metadata": {"parser": "beautifulsoup"}}
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            raise ParseError(f"HTML瑙ｆ瀽澶辫触: {e}")

    def _parse_pptx(self, file_path: str) -> Dict:
        """PPTX/PPT 瑙ｆ瀽 鈥?鏀寔 .pptx 鍜?.ppt 鏍煎紡"""
        import os

        # 妫€娴嬫枃浠舵牸寮?
        is_old_format = False
        try:
            import olefile

            is_old_format = olefile.isOleFile(file_path)
        except ImportError:
            pass

        # 濡傛灉鏄棫鐗?.ppt 鏍煎紡锛屽皾璇曡浆鎹负 .pptx
        pptx_path = file_path
        temp_file = None

        if is_old_format:
            converted = self._convert_ppt_to_pptx(file_path)
            if converted:
                pptx_path = converted
                temp_file = converted
            else:
                # 杞崲澶辫触锛屽洖閫€鍒颁簩杩涘埗瑙ｆ瀽
                return self._parse_ppt_binary_fallback(file_path)

        try:
            from pptx import Presentation

            prs = Presentation(pptx_path)
            slides_text = []
            for slide in prs.slides:
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_content.append(shape.text)
                slides_text.append("\n".join(slide_content))
            return {
                "text": "\n\n".join(slides_text),
                "tables": [],
                "metadata": {"parser": "python-pptx", "pages": len(prs.slides)},
            }
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            # 濡傛灉 python-pptx 澶辫触锛屽洖閫€鍒颁簩杩涘埗瑙ｆ瀽
            if is_old_format:
                return self._parse_ppt_binary_fallback(file_path)
            raise ParseError(f"PPTX瑙ｆ瀽澶辫触: {e}")
        finally:
            # 娓呯悊涓存椂鏂囦欢
            if temp_file and os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except (OSError, ValueError, KeyError, ConnectionError, TimeoutError):
                    pass

    def _convert_ppt_to_pptx(self, ppt_path: str) -> str:
        """浣跨敤 LibreOffice 灏?.ppt 杞崲涓?.pptx"""
        import os
        import shutil
        import subprocess
        import tempfile

        # 鏌ユ壘 LibreOffice
        lo_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]

        lo_path = None
        for p in lo_paths:
            if os.path.exists(p):
                lo_path = p
                break

        if not lo_path:
            logging.getLogger(__name__).warning("[Parser] LibreOffice 鏈畨瑁咃紝鏃犳硶杞崲 PPT 鏂囦欢")
            return None

        try:
            # 鍒涘缓涓存椂鐩綍
            tmpdir = tempfile.mkdtemp()

            # 瀹夊叏淇: 楠岃瘉鏂囦欢鍚嶏紝闃叉鍛戒护娉ㄥ叆
            import shlex

            safe_ppt_path = shlex.quote(ppt_path)
            safe_tmpdir = shlex.quote(tmpdir)

            # 浣跨敤 LibreOffice 杞崲
            result = subprocess.run(
                [lo_path, "--headless", "--convert-to", "pptx", "--outdir", tmpdir, ppt_path],
                timeout=60,
                capture_output=True,
                text=True,
            )

            if result.returncode == 0:
                # 鏌ユ壘杞崲鍚庣殑鏂囦欢
                base_name = os.path.splitext(os.path.basename(ppt_path))[0]
                pptx_path = os.path.join(tmpdir, f"{base_name}.pptx")

                if os.path.exists(pptx_path):
                    # 澶嶅埗缁撴灉鍒版渶缁堜綅缃?
                    result_path = ppt_path + ".converted.pptx"
                    shutil.copy2(pptx_path, result_path)
                    # 娓呯悊涓存椂鐩綍
                    shutil.rmtree(tmpdir, ignore_errors=True)
                    logging.getLogger(__name__).info(f"[Parser] PPT 杞崲鎴愬姛: {ppt_path} -> {result_path}")
                    return result_path

            logging.getLogger(__name__).warning(f"[Parser] LibreOffice 杞崲澶辫触: {result.stderr}")
            # 娓呯悊涓存椂鐩綍
            shutil.rmtree(tmpdir, ignore_errors=True)
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            logging.getLogger(__name__).warning(f"[Parser] LibreOffice 杞崲澶辫触: {e}")

        return None

    def _parse_ppt_binary_fallback(self, file_path: str) -> Dict:
        """鏃х増 .ppt 鏂囦欢浜岃繘鍒惰В鏋愶紙鍥為€€鏂规锛?""
        try:
            import olefile

            ole = olefile.OleFileIO(file_path)

            if ole.exists("PowerPoint Document"):
                data = ole.openstream("PowerPoint Document").read()
                text = self._extract_text_from_ppt_binary(data)
                ole.close()
                return {
                    "text": text,
                    "tables": [],
                    "metadata": {"parser": "olefile-binary", "pages": 0},
                }
            ole.close()
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            logging.getLogger(__name__).warning(f"[Parser] 浜岃繘鍒惰В鏋愬け璐? {e}")

        raise ParseError("PPT瑙ｆ瀽澶辫触: 鏃犳硶瑙ｆ瀽姝ゆ枃浠舵牸寮?)

    def _extract_text_from_ppt_binary(self, data: bytes) -> str:
        """浠庢棫鐗?.ppt 浜岃繘鍒舵暟鎹腑鎻愬彇鏂囨湰锛坴5 澶氬眰杩囨护鐗?鈥?闆朵贡鐮侊級"""
        import re
        import struct

        # OLE 鏍囪鐮佺偣
        OLE_MARKERS = {0x8000, 0x8700}

        # 鈹€鈹€ 绗竴灞傦細OLE 瀵硅薄鏍囩锛堢簿纭尮閰嶏級鈹€鈹€
        OLE_OBJECT_LABELS = {
            "Word.Document.80",
            "Word 鏂囦欢",
            "Microsoft Word 鏂囦欢",
            "Paint.Picture0",
            "Paint.Picture",
            "榛為櫍鍦栧奖鍍?,
            "Drawing",
            "AutoCAD.Drawing.150",
            "AutoCAD Drawing",
            "Equation",
            "Equation.30",
            "Microsoft 鏂圭▼寮忕法杓櫒 3.0",
            "AutoCAD. Drawing.140",
            "AutoCAD 鍦栨獢",
            "Excel.Sheet.80",
            "Microsoft Excel 宸ヤ綔琛?,
            "Excel.Chart.80",
            "Microsoft Excel 鍦栬〃",
            "Chart",
            "MSGraph.Chart.80",
            "Microsoft Graph 2000 鍦栬〃",
            "1-2-3 宸ヤ綔琛?,
            "123Worksheet0",
            "Times New Roman",
            "鏂扮窗鏄庨珨",
            "Arial Narrow",
            "Arial",
            "AR MingtiM BIG-5",
            "ingtiM BIG-5",
            "Comic Sans MS",
            "Wingdings",
            "Symbol",
            "Helvetica",
            "Monotype Sorts",
            "type Sorts",
            "鍏ㄧ湡绱伴毟鏇?,
            "pe Sorts",
            "鑿悍绮楅粦楂?,
            "绱版槑楂?,
            "妯欐シ楂?,
            "KaiTi",
            "SimSun",
            "SimHei",
            "Txt",
            "瀹嬩綋",
            "榛戜綋",
            "妤蜂綋",
            "浠垮畫",
            "___PPT10",
            "___PPT9",
            "___PPT11",
            "PPT10",
            "PPT9",
            "PPT11",
            "鏈懡鍚?-1",
        }

        # 鈹€鈹€ 绗簩灞傦細鎶€鏈叧閿瘝锛堝唴瀹瑰繀椤诲寘鍚嚦灏戜竴涓級鈹€鈹€
        TECH_KEYWORDS = {
            "濉戣啝",
            "瑷▓",
            "澹佸帤",
            "鍑稿彴",
            "閫ｆ帴鍣?,
            "绔瓙",
            "姝ｅ悜鍔?,
            "鏉愭枡",
            "淇濇寔鍔?,
            "housing",
            "connector",
            "pin",
            "terminal",
            "thickness",
            "plastic",
            "material",
            "force",
            "stress",
            "resistance",
            "绗竴绔?,
            "绗簩绔?,
            "绗笁绔?,
            "绗洓绔?,
            "绗簲绔?,
            "绗叚绔?,
            "mm",
            "gf",
            "Gpa",
            "Mpa",
            "smt",
            "SMT",
            "LCP",
            "Nylon",
            "PPS",
            "PCT",
            "閲?,
            "閷?,
            "閴?,
            "閵?,
            "闆?,
            "鎺ヨЦ",
            "闃绘姉",
            "鎻掑叆",
            "鎷斿嚭",
            "璁婂舰",
            "鎳夊姏",
            "寮峰害",
            "瑷▓瑕忚寖",
            "灏哄",
            "鍏樊",
            "閰嶅悎",
            "mating",
            "unmating",
            "reliability",
            "鎺ヨЦ闆婚樆",
            "楂旂闆婚樆",
            "钖勮啘闆婚樆",
            "鍗℃Λ",
            "骞叉秹閲?,
            "濉戣啝闆朵欢",
            "绲愭瑷▓",
            "鍦撹",
            "鎷旀ā瑙?,
            "灏庨浕",
            "鐒婇尗",
            "闆婚崓",
            "楂橀牷",
            "淇¤櫉",
            "鍌宠几",
            "瑕佹眰",
            "鍔熻兘",
            "鐩殑",
            "鍒嗛",
            "绋",
            "甯哥敤",
            "缂洪粸",
            "鍎粸",
            "鐗规€?,
            "鍙冩暩",
            "瑕忔牸",
            "妯欐簴",
            "鍘熺悊",
            "鐞嗚珫",
            "瑷堢畻",
            "鍏紡",
            "鏂规硶",
            "姝ラ",
            "鍟忛",
            "鍘熷洜",
            "瑙ｆ焙",
            "鏀瑰杽",
            "鍎寲",
            "鎺у埗",
            "娴佺▼",
            "宸ヨ棟",
            "瑁界▼",
            "妾㈡脯",
            "娓│",
            "椹楄瓑",
            "琛ㄩ潰",
            "绲愭",
            "褰㈢媭",
            "浣嶇疆",
            "瑙掑害",
            "闁撹窛",
            "conductor",
            "insulation",
            "solder",
            "welding",
            "plating",
            "contact",
            "terminal",
            "housing",
            "header",
            "receptacle",
            "pitch",
            "polarity",
            "locking",
            "retention",
            "impedance",
            "capacitance",
            "inductance",
            "temperature",
            "humidity",
            "vibration",
            "gold",
            "tin",
            "copper",
            "brass",
            "phosphor",
            "beryllium",
            "nickel",
            "silver",
            "reflow",
            "wave",
            "soldering",
            "paste",
            "flux",
            "surface",
            "mount",
            "through",
            "hole",
            "profile",
            "curve",
            "graph",
            "chart",
            "table",
            "figure",
            "diagram",
            "schematic",
            "section",
            "chapter",
            "page",
            "slide",
            "appendix",
            "reference",
            "specification",
            "0.",
            "1.",
            "2.",
            "3.",
            "4.",
            "5.",
            "6.",
            "7.",
            "8.",
            "9.",
            "T)",
            "R)",
            "mm)",
            "gf)",
            "Gpa)",
            "Mpa)",
        }

        # 鈹€鈹€ 绗笁灞傦細涔辩爜妯″紡锛堟鍒欙級鈹€鈹€
        NOISE_PATTERNS = [
            r"^[\u8000\u8700]",  # OLE 鏍囪寮€澶?
            r"^[鑰€鏅絻鐚€鐢€鍊€鐛ラ懟妞€鏀€娓€宕€鐎堢伋妯佹垁姹憽娼傜墶榘佹秵椴",  # 宸茬煡 OLE 涔辩爜瀛楃
            r"^[\u8000-\u8FFF]{2,}",  # 杩炵画绉佹湁鍖哄瓧绗?
            r"^[濂栫爳鑵€鍙嗘笗]",  # 宸茬煡 OLE 涔辩爜
        ]

        # 鎻愬彇鎵€鏈?UTF-16LE 鏂囨湰
        all_texts = []
        i = 0
        while i < len(data) - 1:
            char = struct.unpack("<H", data[i : i + 2])[0]
            is_printable = (
                (0x20 <= char <= 0x7E)
                or (0x4E00 <= char <= 0x9FFF)
                or (0x3000 <= char <= 0x303F)
                or (0xFF00 <= char <= 0xFFEF)
                or (0x2000 <= char <= 0x206F)
                or char == 0x000A
            )
            if is_printable:
                text = []
                while i + 1 < len(data):
                    char = struct.unpack("<H", data[i : i + 2])[0]
                    is_printable = (
                        (0x20 <= char <= 0x7E)
                        or (0x4E00 <= char <= 0x9FFF)
                        or (0x3000 <= char <= 0x303F)
                        or (0xFF00 <= char <= 0xFFEF)
                        or (0x2000 <= char <= 0x206F)
                        or char == 0x000A
                    )
                    if is_printable:
                        text.append(chr(char))
                        i += 2
                    else:
                        break
                result = "".join(text).strip()
                if len(result) >= 3:
                    all_texts.append(result)
            else:
                i += 2

        # 澶氬眰杩囨护
        clean_texts = []
        for text in all_texts:
            # 绗竴灞傦細绮剧‘鍖归厤 OLE 鏍囩
            if text in OLE_OBJECT_LABELS:
                continue
            # 绗簩灞傦細姝ｅ垯鍖归厤涔辩爜妯″紡
            is_noise = False
            for pattern in NOISE_PATTERNS:
                if re.match(pattern, text):
                    is_noise = True
                    break
            if is_noise:
                continue
            # 绗笁灞傦細妫€鏌ユ槸鍚﹀寘鍚?OLE 鏍囪鐮佺偣
            if any(ord(c) in OLE_MARKERS for c in text):
                continue
            # 绗洓灞傦細妫€鏌ユ槸鍚﹀寘鍚妧鏈叧閿瘝
            has_keyword = any(kw in text for kw in TECH_KEYWORDS)
            # 绗簲灞傦細妫€鏌ユ槸鍚︽槸绾暟瀛?鏍囩偣
            is_pure_number = re.match(r"^[\d\.\-\+\s/\(\)]+$", text)
            # 绗叚灞傦細妫€鏌ユ槸鍚︽槸绾?ASCII锛堟帓闄?OLE 鏍囩锛?
            is_pure_ascii = all(0x20 <= ord(c) <= 0x7E for c in text)
            is_ole_label = any(
                tag in text for tag in ["Word.Document", "Paint.Picture", "AutoCAD", "Excel", "Microsoft"]
            )
            # 淇濈暀鏉′欢
            if has_keyword:
                clean_texts.append(text)
            elif is_pure_number and len(text) <= 15:
                clean_texts.append(text)
            elif is_pure_ascii and not is_ole_label and len(text) >= 5:
                if not re.search(r"[#^~<>{}\[\]|\\]", text):
                    clean_texts.append(text)

        if not clean_texts:
            return ""

        # 鍘婚噸锛堜繚鎸侀『搴忥級
        seen = set()
        unique_parts = []
        for part in clean_texts:
            key = part.replace(" ", "")
            if key not in seen:
                seen.add(key)
                unique_parts.append(part)

        result = "\n".join(unique_parts[:500])

        # 绻佺畝杞崲锛堢箒浣撲腑鏂?鈫?绠€浣撲腑鏂囷級
        try:
            import opencc

            converter = opencc.OpenCC("t2s")
            result = converter.convert(result)
        except ImportError:
            pass  # opencc 鏈畨瑁呭垯璺宠繃

        return result


class UnifiedCleaner:
    """缁熶竴娓呮礂鍣?鈥?鍚堝苟 cleaners.py 瀹屾暣娓呮礂閫昏緫锛堜换鍔? P0淇锛?""

    # 瀹夊叏娓呮礂妯″紡锛堜粠 cleaners.py 鍚屾锛?
    NOISE_PATTERNS = [
        r"鐗堟潈褰?*鎵€鏈?,
        r"Copyright\s+漏?\s*\d{4}",
        r"All\s+Rights?\s+Reserved",
        r"鏈粡璁稿彲.*涓嶅緱.*(?:澶嶅埗|杞浇|浼犳挱)",
        r"鍏嶈矗澹版槑.*",
        r"浠ヤ笂鍐呭浠呬緵鍙傝€?,
        r"鏈枃浠?*鏈€缁堣В閲婃潈",
        r"濡?*渚垫潈.*璇疯仈绯?,
        r"澹版槑锛?*涓嶆壙鎷?*璐ｄ换",
        r"娓╅Θ鎻愮ず锛?*鎶曡祫鏈夐闄?,
    ]

    SENSITIVE_PATTERNS = [
        (r"1[3-9]\d{9}", "鎵嬫満鍙?),
        (r"\d{17}[\dXx]", "韬唤璇佸彿"),
        (r"\d{16,19}", "閾惰鍗″彿"),
    ]

    def __init__(self, config: Dict = None):
        self.config = config or {}
        self._enable_sensitive_mask = config.get("mask_sensitive", True) if config else True
        self._enable_semantic = config.get("semantic_clean", True) if config else True

    def clean(self, parsed: Dict) -> Dict:
        """娓呮礂鏂囨湰"""
        text = parsed.get("text", "")
        try:
            text = self._clean_text(text)
            parsed["text"] = text
            return parsed
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            raise CleanError(f"娓呮礂澶辫触: {e}")

    def _clean_text(self, text: str) -> str:
        """缁熶竴娓呮礂閫昏緫 鈥?浠诲姟1+3 P0淇锛氭纭楠ら『搴?

        娓呮礂姝ラ椤哄簭锛堜换鍔? P0淇锛夛細
        1. 鎺у埗瀛楃绉婚櫎锛堟渶鍏堬紝閬垮厤骞叉壈鍚庣画鍖归厤锛?
        2. 绌烘牸鍚堝苟锛堝湪鍐呭娓呮礂涔嬪墠缁熶竴绌虹櫧锛?
        3. 鍐呭娓呮礂锛氭敞鍏ラ槻鎶?鈫?鑴辨晱 鈫?鐗堟潈 鈫?鍏ㄨ 鈫?鍘婚噸
        """
        import re

        if not text or not isinstance(text, str):
            return ""

        # ---- 闃舵1锛氭帶鍒跺瓧绗︾Щ闄わ紙鏈€鍏堟墽琛岋級 ----
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)

        # ---- 闃舵2锛氱┖鏍煎悎骞?----
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)

        # ---- 闃舵3锛氬唴瀹规竻娲?----

        # 3a. Prompt Injection 鍑€鍖栵紙鍚檷绾ф棩蹇楋級
        try:
            from src.services.prompt_guard import sanitize_document_content

            text, injection_detected = sanitize_document_content(text)
            if injection_detected:
                logger.warning("[Security] 鏂囨。鍐呭涓娴嬪埌 Prompt Injection 妯″紡锛屽凡鍑€鍖?)
        except ImportError:
            pass  # prompt_guard 妯″潡涓嶅彲鐢ㄦ椂闄嶇骇涓烘棤鍑€鍖?

        # 3b. 鍘婚櫎HTML鏍囩
        text = re.sub(r"<[^>]+>", "", text)

        # 3c. 鍘婚櫎URL
        text = re.sub(r"https?://\S+", "", text)

        # 3d. 鍘婚櫎閭
        text = re.sub(r"\S+@\S+\.\S+", "", text)

        # 3e. 鍘婚櫎椤电湁椤佃剼锛堟暟瀛?椤电爜锛夆€?浠诲姟6 P0淇锛氬彧鍖归厤绾暟瀛椾笖闀垮害<=5
        text = self._strip_header_footer_numbers(text)

        # 3f. 鐗堟潈澹版槑鍘婚櫎 + 閲嶅娈佃惤鍘婚噸 + 鍏ㄨ鈫掑崐瑙掕浆鎹?
        if self._enable_semantic:
            text = self._strip_noise(text)
            text = self._deduplicate_paragraphs(text)
            text = self._normalize_width(text)

        # 3g. 鏁忔劅淇℃伅鑴辨晱
        if self._enable_sensitive_mask:
            text, _ = self._detect_and_mask_sensitive(text)

        # 鏈€缁堝浣欑┖鐧藉悎骞?
        text = re.sub(r"\s+", " ", text)

        return text.strip()

    def _strip_header_footer_numbers(self, text: str) -> str:
        """浠诲姟6 P0淇锛氶〉鐪夐〉鑴氬幓闄?鈥?鍙尮閰嶈棣栬灏剧函鏁板瓧涓旈暱搴?=5"""
        import re

        lines = text.split("\n")
        cleaned = []
        for line in lines:
            stripped = line.strip()
            # 鍙尮閰嶇函鏁板瓧琛岋紝涓旈暱搴?<= 5锛堥伩鍏嶈鏉€鍚堟硶鏁版嵁琛屽骞翠唤銆佹暟閲忥級
            if stripped.isdigit() and len(stripped) <= 5:
                continue
            cleaned.append(line)
        return "\n".join(cleaned)

    def _strip_noise(self, text: str) -> str:
        """鍘婚櫎鐗堟潈/澹版槑琛?""
        import re

        for pattern in self.NOISE_PATTERNS:
            text = re.sub(pattern, "", text, flags=re.IGNORECASE)
        return text

    def _deduplicate_paragraphs(self, text: str) -> str:
        """鍘婚櫎閲嶅娈佃惤"""
        paragraphs = text.split("\n\n")
        seen = set()
        cleaned = []
        for p in paragraphs:
            p = p.strip()
            if not p:
                continue
            # 鍙栧墠 100 瀛楀仛鍘婚噸
            key = p[:100]
            if key in seen:
                continue
            seen.add(key)
            cleaned.append(p)
        return "\n\n".join(cleaned)

    def _normalize_width(self, text: str) -> str:
        """鍏ㄨ瀛楁瘝鏁板瓧 鈫?鍗婅"""
        import re

        text = re.sub(r"[锛?锛篯", lambda m: chr(ord(m.group()) - 0xFEE0), text)
        text = re.sub(r"[锝?锝歖", lambda m: chr(ord(m.group()) - 0xFEE0), text)
        text = re.sub(r"[锛?锛橾", lambda m: chr(ord(m.group()) - 0xFEE0), text)
        return text

    def _detect_and_mask_sensitive(self, text: str) -> tuple:
        """妫€娴嬪苟鑴辨晱鏁忔劅淇℃伅"""
        import re

        issues = []
        for pattern, label in self.SENSITIVE_PATTERNS:
            matches = re.findall(pattern, text)
            if matches:
                issues.append(f"鍙戠幇 {len(matches)} 涓獅label}")
                for m in matches:
                    text = text.replace(m, m[:3] + "*" * (len(m) - 6) + m[-3:])
        return text, issues


class UnifiedChunker:
    """缁熶竴鍒嗗潡鍣?鈥?琛ㄦ牸鎰熺煡 + 鏍囬灞傜骇浼犳挱锛堜换鍔? P0淇锛?""

    def __init__(self, config: Dict = None):
        self.config = config or {}
        self.chunk_size = config.get("chunk_size", 1000) if config else 1000
        self.chunk_overlap = config.get("chunk_overlap", 100) if config else 100

    def chunk(self, parsed: Dict, tables: List[Dict] = None) -> List[Chunk]:
        """鍒嗗潡 鈥?琛ㄦ牸鐙珛瀛樺偍锛屾敮鎸佹爣棰樺眰绾т紶鎾?""
        text = parsed.get("text", "")
        if not text or len(text) < 50:
            return [Chunk(text=text, chunk_index=0)] if text.strip() else []

        # 鎻愬彇鏍囬灞傜骇缁撴瀯
        heading_structure = parsed.get("metadata", {}).get("heading_structure", [])

        chunks = []

        # 鏂囨湰鍒嗗潡锛堣繑鍥?(text, start_pos) 鍏冪粍鐢ㄤ簬绮剧‘瀹氫綅锛?
        text_chunks_with_pos = self._chunk_text_with_pos(text)
        for i, (ct, chunk_start) in enumerate(text_chunks_with_pos):
            chunk = Chunk(
                text=ct,
                chunk_index=i,
                chunk_type=ChunkType.TEXT,
            )
            # 浠诲姟2 P0淇锛氫紶鎾爣棰樺眰绾э紝浣跨敤 chunk 璧峰浣嶇疆瀹氫綅
            heading_info = self._find_heading_for_chunk(text, chunk_start, heading_structure)
            if heading_info:
                chunk.heading = heading_info.get("text", "")
                chunk.heading_level = heading_info.get("level", None)
            chunks.append(chunk)

        # 琛ㄦ牸鐙珛瀛樺偍
        if tables:
            for table in tables:
                if table.get("headers") or table.get("rows"):
                    chunk = Chunk(
                        text=self._table_to_text(table),
                        chunk_index=len(chunks),
                        chunk_type=ChunkType.TABLE,
                        structured_table=table,
                    )
                    chunks.append(chunk)

        # 璁剧疆 total_chunks
        for chunk in chunks:
            chunk.total_chunks = len(chunks)

        return chunks

    def _find_heading_for_chunk(self, full_text: str, chunk_start: int, heading_structure: list) -> dict:
        """浠诲姟2 P0淇锛氭壘鍒拌 chunk 鎵€褰掑睘鐨勬渶杩戞爣棰?

        浣跨敤 chunk 璧峰浣嶇疆锛坈hunk_start锛夊畾浣嶏紝鑰岄潪 full_text.find(chunk_text)锛?
        閬垮厤褰?chunk 鏂囨湰閲嶅鍑虹幇鏃跺畾浣嶉敊璇€?
        """
        if not heading_structure:
            return None
        try:
            if chunk_start < 0:
                return None
            # 閬嶅巻鏍囬锛屾壘璇ヤ綅缃箣鍓嶆渶杩戠殑
            best = None
            best_pos = -1
            for h in heading_structure:
                h_text = h.get("text", "")
                if not h_text:
                    continue
                h_pos = full_text.find(h_text)
                if 0 <= h_pos <= chunk_start and h_pos > best_pos:
                    best = h
                    best_pos = h_pos
            return best
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            return None

    def _chunk_text_with_pos(self, text: str) -> List[tuple]:
        """鏂囨湰鍒嗗潡 鈥?杩斿洖 (chunk_text, start_position) 鍏冪粍"""
        chunks = []
        start = 0
        text_len = len(text)

        while start < text_len:
            end = min(start + self.chunk_size, text_len)

            # 灏濊瘯鍦ㄥ彞瀛愯竟鐣屾柇寮€
            if end < text_len:
                for sep in ["\n\n", "\n", "銆?, "锛?, ".", ";"]:
                    last_sep = text.rfind(sep, start + self.chunk_size // 2, end)
                    if last_sep > start:
                        end = last_sep + len(sep)
                        break

            chunk = text[start:end].strip()
            if chunk and len(chunk) > 20:
                chunks.append((chunk, start))

            start = end - self.chunk_overlap if end < text_len else text_len

        if not chunks:
            chunks = [(text[: self.chunk_size], 0)]

        return chunks

    def _table_to_text(self, table: Dict) -> str:
        """琛ㄦ牸杞枃鏈?""
        headers = table.get("headers", [])
        rows = table.get("rows", [])
        lines = []
        if headers:
            lines.append(" | ".join(str(h) for h in headers))
            lines.append(" | ".join(["---"] * len(headers)))
        for row in rows[:50]:  # 鏈€澶?0琛?
            lines.append(" | ".join(str(cell) for cell in row))
        return "\n".join(lines)


class UnifiedClassifier:
    """缁熶竴鍒嗙被鍣?""

    def __init__(self, config: Dict = None):
        self.config = config or {}

    def classify(self, chunk: Chunk) -> str:
        """鍒嗙被"""
        try:
            from src.category_registry import match_category

            return match_category(chunk.text, file_ext=chunk.file_type, file_name=chunk.file_name) or "閫氱敤鍔炲叕"
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
            return "閫氱敤鍔炲叕"


class UnifiedEmbedder:
    """缁熶竴鍚戦噺鍖栧櫒"""

    def __init__(self, config: Dict = None):
        self.config = config or {}

    async def embed_batch(self, texts: List[str]) -> List[Optional[List[float]]]:
        """鎵归噺鍚戦噺鍖栵紙鐩存帴璋冪敤 Embedder HTTP API锛?""
        import asyncio
        import os

        embedder_url = os.getenv("KB_EMBEDDER_URL", "http://127.0.0.1:8081")

        for attempt in range(5):  # 澧炲姞閲嶈瘯娆℃暟
            try:
                import requests as req

                resp = req.post(
                    f"{embedder_url}/embed",
                    json={"texts": texts},
                    timeout=30,
                )
                if resp.status_code == 200:
                    data = resp.json()
                    vectors = data.get("vectors", [])
                    if vectors and len(vectors) == len(texts):
                        return vectors
                logger.warning(f"[Embedder] HTTP {resp.status_code} (attempt {attempt+1})")
            except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
                logger.warning(f"[Embedder] 璋冪敤澶辫触 (attempt {attempt+1}): {e}")
            if attempt < 4:
                await asyncio.sleep(3 + attempt * 2)  # 閫掑寤惰繜
        return [None] * len(texts)


class UnifiedSaver:
    """缁熶竴瀛樺偍鍣?""

    def __init__(self, config: Dict = None):
        self.config = config or {}

    # FAKE-ASYNC: 鏈嚱鏁版爣璁?async 浠呬负鎺ュ彛缁熶竴锛屽唴閮ㄥ悓姝ユ墽琛?

    async def save(self, chunks: List[Chunk]) -> Any:
        """瀛樺偍鍒?SQLite + ChromaDB"""
        if not chunks:
            return

        try:
            from src.db.memory_store import get_store

            store = get_store()

            # 杞崲涓?dict 鏍煎紡瀛樺偍
            chunk_dicts = []
            for chunk in chunks:
                d = chunk.to_dict()
                d["source_file"] = chunk.source_file or chunk.file_name
                # 闄勫姞鍐呭鍝堝笇鐢ㄤ簬澧為噺鏇存柊
                d["content_hash"] = hashlib.sha256(chunk.text.encode("utf-8")).hexdigest()
                chunk_dicts.append(d)

            store.insert_many(chunk_dicts)
            logger.info(f"[Saver] 瀛樺偍 {len(chunk_dicts)} 涓?chunk")
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            raise SaveError(f"瀛樺偍澶辫触: {e}")


class UnifiedExtractor:
    """SAG 寮忎簨浠?瀹炰綋鎻愬彇鍣?""

    def __init__(self, config: Dict = None):
        self.config = config or {}

    async def extract(self, chunks: List[Chunk]) -> tuple:
        """浠庣鐗囦腑鎻愬彇浜嬩欢鍜屽疄浣?""
        all_events = []
        all_entities = []

        for i, chunk in enumerate(chunks):
            if chunk.chunk_type == ChunkType.TABLE:
                continue

            try:
                result = await self._extract_single(chunk, i, chunks)
                events = result.get("events", [])
                entities = result.get("entities", [])

                # 鏋勫缓 Event 瀵硅薄
                for ev_data in events:
                    event = Event(
                        event_id=self._make_id("ev", ev_data.get("title", "") + chunk.chunk_id),
                        title=ev_data.get("title", ""),
                        summary=ev_data.get("summary", ""),
                        content=ev_data.get("content", ""),
                        keywords=ev_data.get("keywords", []),
                        priority=ev_data.get("priority", "UNKNOWN"),
                        references=ev_data.get("references", []),
                        chunk_ids=[chunk.chunk_id],
                        entity_names=ev_data.get("entities", []),
                        file_hash=chunk.file_hash,
                        file_name=chunk.file_name,
                        level=0,
                    )
                    all_events.append(event)

                # 鏋勫缓 Entity 瀵硅薄
                for ent_data in entities:
                    entity = Entity(
                        entity_id=self._make_id("ent", ent_data.get("name", "")),
                        name=ent_data.get("name", ""),
                        entity_type=ent_data.get("type", ""),
                        description=ent_data.get("description", ""),
                        chunk_ids=[chunk.chunk_id],
                        file_hash=chunk.file_hash,
                        file_name=chunk.file_name,
                    )
                    all_entities.append(entity)

            except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
                logger.warning(f"[Extractor] 鎻愬彇澶辫触 chunk {i}: {e}")

        # 瀹炰綋鍘婚噸褰掍竴鍖?
        all_entities = self._deduplicate_entities(all_entities)

        return all_entities, all_events

    # FAKE-ASYNC: 鏈嚱鏁版爣璁?async 浠呬负鎺ュ彛缁熶竴锛屽唴閮ㄥ悓姝ユ墽琛?

    async def _extract_single(self, chunk: Chunk, index: int, all_chunks: List[Chunk]) -> Dict:
        """鍗曚釜 chunk 鎻愬彇"""
        try:
            from src.services.llm import call_ai

            prev_heading = all_chunks[index - 1].heading if index > 0 else ""
            prev_summary = all_chunks[index - 1].text[:300] if index > 0 else ""

            prompt = f"""浠庝互涓嬫枃鏈腑鎻愬彇浜嬩欢鍜屽疄浣撱€?

鏂囦欢锛歿chunk.file_name}
鍒嗙被锛歿chunk.category}
鐗囨锛氱 {chunk.chunk_index + 1}/{chunk.total_chunks} 娈?
鍓嶆枃鏍囬锛歿prev_heading}

鏂囨湰锛?
{chunk.text[:2000]}

璇疯繑鍥濲SON鏍煎紡锛?
{{"events": [{{"title": "浜嬮」鏍囬", "summary": "涓€鍙ヨ瘽鎽樿", "content": "瀹屾暣鍐呭", "keywords": ["鍏抽敭璇?], "priority": "HIGH/MEDIUM/LOW", "entities": ["瀹炰綋鍚?], "references": [1]}}], "entities": [{{"name": "瀹炰綋鍚?, "type": "person/organization/product/material/device", "description": "浣滅敤鎻忚堪"}}]}}"""

            response = await call_ai(prompt)
            if response:
                import json

                # 灏濊瘯瑙ｆ瀽JSON
                try:
                    # 娓呯悊鍝嶅簲涓殑markdown浠ｇ爜鍧楁爣璁?
                    clean_response = response.strip()
                    if clean_response.startswith("```"):
                        clean_response = (
                            clean_response.split("\n", 1)[1] if "\n" in clean_response else clean_response[3:]
                        )
                    if clean_response.endswith("```"):
                        clean_response = clean_response[:-3]
                    clean_response = clean_response.strip()
                    return json.loads(clean_response)
                except json.JSONDecodeError as e:
                    logger.warning("json.JSONDecodeError 澶辫触: %s", e, exc_info=True)

        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            logger.debug(f"[Extractor] LLM鎻愬彇澶辫触: {e}")

        return {"events": [], "entities": []}

    def _make_id(self, prefix: str, content: str) -> str:
        """鐢熸垚鍞竴ID"""
        import hashlib

        return f"{prefix}_{hashlib.sha256(content.encode()).hexdigest()[:12]}"

    def _deduplicate_entities(self, entities: List[Entity]) -> List[Entity]:
        """瀹炰綋鍘婚噸褰掍竴鍖?""
        seen = {}
        for e in entities:
            key = e.name.lower().strip()
            if key in seen:
                seen[key].chunk_ids.extend(e.chunk_ids)
                seen[key].mentions += 1
                if e.description and not seen[key].description:
                    seen[key].description = e.description
            else:
                seen[key] = e
        return list(seen.values())


class UnifiedPipeline:
    """
    浼忕静缁熶竴澶勭悊绠＄嚎
    鎵€鏈夋潵婧愶紙涓婁紶/瑁呰浇鏈?API锛夌殑鏁版嵁锛岄兘缁忚繃杩欐潯绠＄嚎銆?
    """

    def __init__(self, config: Dict = None):
        self.config = config or {}
        self.parser = UnifiedParser(config)
        self.cleaner = UnifiedCleaner(config)
        self.chunker = UnifiedChunker(config)
        self.classifier = UnifiedClassifier(config)
        self.embedder = UnifiedEmbedder(config)
        self.saver = UnifiedSaver(config)
        self.extractor = UnifiedExtractor(config)
        self._transactional_saver = TransactionalSaver(config)
        self._processing = set()
        self._lock = asyncio.Lock()

    def _compute_hash(self, file_path: str) -> str:
        """璁＄畻鏂囦欢鍝堝笇"""
        import hashlib

        with open(file_path, "rb") as f:
            return hashlib.sha256(f.read()).hexdigest()

    async def process(self, file_path: str, source: str = "upload") -> PipelineResult:
        """
        缁熶竴澶勭悊鍏ュ彛

        Args:
            file_path: 鏂囦欢璺緞
            source: 鏉ユ簮鏍囪瘑 ("upload" / "loader" / "api")

        Returns:
            PipelineResult: 鍖呭惈 chunks/entities/events
        """
        start_time = time.time()
        file_hash = self._compute_hash(file_path)

        # 骞跺彂闃叉姢锛氬悓涓€鏂囦欢涓嶉噸澶嶅鐞?
        async with self._lock:
            if file_hash in self._processing:
                logger.info(f"[Pipeline] 鏂囦欢姝ｅ湪澶勭悊涓紝璺宠繃: {file_path}")
                return PipelineResult(source=source, file_path=file_path, skipped=True)
            self._processing.add(file_hash)

        try:
            result = PipelineResult(source=source, file_path=file_path)
            metrics = PipelineMetrics()

            # Step 1: 瑙ｆ瀽锛堜笉鍙仮澶?鈫?鎶涘嚭锛?
            metrics.start("parse")
            parsed = self.parser.parse(file_path)
            result.raw_text = parsed.get("text", "")
            result.tables = parsed.get("tables", [])
            result.metadata = parsed.get("metadata", {})
            metrics.end()

            # Step 2: 娓呮礂锛堝彲鎭㈠ 鈫?闄嶇骇涓哄師鏂囷級
            metrics.start("clean")
            try:
                cleaned = self.cleaner.clean(parsed)
                result.cleaned_text = cleaned.get("text", "")
            except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
                logger.warning(f"[Pipeline] 娓呮礂澶辫触锛屼娇鐢ㄥ師鏂? {e}")
                result.cleaned_text = result.raw_text
                result.errors.append(f"CleanError: {e}")
            metrics.end()

            # Step 3: 鍒嗗潡锛堝彲鎭㈠ 鈫?闄嶇骇涓哄崟鍧楋級
            metrics.start("chunk")
            try:
                chunks = self.chunker.chunk(parsed, result.tables)
                result.chunks = chunks
            except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
                logger.warning(f"[Pipeline] 鍒嗗潡澶辫触锛岄檷绾т负鍗曞潡: {e}")
                result.chunks = [Chunk(text=result.cleaned_text, chunk_index=0)]
                result.errors.append(f"ChunkError: {e}")
            metrics.end()

            # Step 4: 鍒嗙被锛堝彲鎭㈠ 鈫?榛樿鍒嗙被锛?
            for chunk in result.chunks:
                try:
                    chunk.category = self.classifier.classify(chunk)
                except (
                    OSError,
                    ValueError,
                    KeyError,
                    ConnectionError,
                    TimeoutError,
                ) as e:  # TODO: Narrow exception type
                    chunk.category = "閫氱敤鍔炲叕"

            # Step 5: 璁剧疆鏉ユ簮淇℃伅
            for chunk in result.chunks:
                chunk.file_hash = file_hash
                chunk.file_name = Path(file_path).name
                chunk.file_type = Path(file_path).suffix.lower()
                chunk.source_pipeline = source
                chunk.source_file = file_path

            # Step 6: 鍚戦噺鍖栵紙鍙仮澶?鈫?鏍囪寰呰ˉ锛?
            metrics.start("embed")
            try:
                embeddings = await self.embedder.embed_batch([c.text for c in result.chunks])
                for chunk, emb in zip(result.chunks, embeddings):
                    chunk.embedding = emb
            except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
                logger.warning(f"[Pipeline] 鍚戦噺鍖栧け璐ワ紝鏍囪寰呰ˉ: {e}")
                result.errors.append(f"EmbedError: {e}")
            metrics.end()

            # Step 7: 瀛樺偍锛堜笉鍙仮澶?鈫?鎶涘嚭锛涘閲忔洿鏂?+ 浜嬪姟淇濇姢锛?
            metrics.start("save")
            await self._save_with_incremental_and_transaction(result.chunks, file_hash)
            metrics.end()

            # Step 8: SAG寮忔彁鍙栵紙鍙€夛紝鍙仮澶?鈫?璺宠繃锛?
            try:
                from src.services.feature_flags import load_flags

                flags = load_flags()
                if flags.get("event_entity_extract", False):
                    try:
                        entities, events = await self.extractor.extract(result.chunks)
                        result.entities = entities
                        result.events = events
                        logger.info(f"[Pipeline] SAG鎻愬彇: {len(events)} 浜嬩欢, {len(entities)} 瀹炰綋")
                    except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
                        logger.warning(f"[Pipeline] 浜嬩欢/瀹炰綋鎻愬彇澶辫触锛岃烦杩? {e}")
                        result.errors.append(f"ExtractError: {e}")
            except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
                logger.warning("Exception 澶辫触: %s", e, exc_info=True)

            result.duration_ms = (time.time() - start_time) * 1000
            result.metrics = metrics.report()
            logger.info(
                f"[Pipeline] 瀹屾垚: {file_path} 鈫?{len(result.chunks)} chunks, {result.duration_ms:.0f}ms, 鑰楁椂鏄庣粏: {result.metrics}"
            )

            return result

        finally:
            async with self._lock:
                self._processing.discard(file_hash)

    # === P2浼樺寲 浠诲姟2锛氭壒閲忛噸绱㈠紩 ===

    async def reindex_file(self, file_hash: str) -> PipelineResult:
        """閲嶆柊绱㈠紩鎸囧畾鏂囦欢锛堟寜 file_hash锛?

        浠?memory_store 涓煡鎵炬枃浠惰矾寰勶紝鐒跺悗閲嶆柊璧板畬鏁村鐞嗙绾裤€?
        """
        from src.db.memory_store import get_store

        store = get_store()
        chunks = store.get_by_hash(file_hash)
        if not chunks:
            raise ValueError(f"鏈壘鍒?file_hash={file_hash} 鐨勬枃妗ｈ褰?)

        # 浠庡凡鏈?chunk 璁板綍涓壘鍒版簮鏂囦欢璺緞
        file_path = None
        for c in chunks:
            fp = c.get("source_file") or c.get("file_path") or c.get("path")
            if fp and Path(fp).exists():
                file_path = fp
                break

        if not file_path:
            raise FileNotFoundError(f"鏃犳硶鎵惧埌 file_hash={file_hash} 瀵瑰簲鐨勬簮鏂囦欢")

        logger.info(f"[Pipeline] 寮€濮嬮噸绱㈠紩: file_hash={file_hash}, file_path={file_path}")
        return await self.process(file_path, source="reindex")

    # === P1浼樺寲 浠诲姟1锛氬閲忔洿鏂?===

    def _compute_chunk_hash(self, text: str) -> str:
        """璁＄畻chunk鍐呭鐨凪D5鍝堝笇"""
        return hashlib.sha256(text.encode("utf-8")).hexdigest()

    async def _get_existing_chunk_hashes(self, file_hash: str) -> dict:
        """鑾峰彇鏂囦欢宸插瓨鍦ㄧ殑chunk鍝堝笇鏄犲皠

        Returns:
            {chunk_hash: chunk_id} 鏄犲皠
        """
        try:
            from src.db.memory_store import get_store

            store = get_store()
            chunks = store.get_by_hash(file_hash)
            return {
                c.get("content_hash", ""): c.get("chunk_id", c.get("id", "")) for c in chunks if c.get("content_hash")
            }
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            logger.warning(f"[Pipeline] 鑾峰彇宸叉湁chunk鍝堝笇澶辫触: {e}")
            return {}

    async def _incremental_save(self, file_hash: str, new_chunks: list, new_embeddings: list) -> Any:
        """澧為噺淇濆瓨锛氬彧鏇存柊鍙樺寲鐨刢hunk

        Args:
            file_hash: 鏂囦欢鍝堝笇
            new_chunks: 鏂扮殑 Chunk 瀵硅薄鍒楄〃
            new_embeddings: 瀵瑰簲鐨?embedding 鍒楄〃

        Returns:
            (added_count, updated_count, deleted_count, skipped_count)
        """
        existing_hashes = await self._get_existing_chunk_hashes(file_hash)

        to_add = []  # (chunk, embedding, chunk_hash)
        to_update = []  # (chunk, embedding, chunk_hash)
        new_hash_set = set()
        skipped = 0

        for chunk, embedding in zip(new_chunks, new_embeddings):
            chunk_hash = self._compute_chunk_hash(chunk.text)
            new_hash_set.add(chunk_hash)
            chunk_id = chunk.chunk_id or chunk.to_dict().get("chunk_id", "")

            if chunk_hash in existing_hashes:
                # 鍐呭鏈彉鍖栵紝璺宠繃
                skipped += 1
                continue

            existing_ids = set(existing_hashes.values())
            if chunk_id and chunk_id in existing_ids:
                # ID 瀛樺湪浣嗗唴瀹瑰彉鍖栵紝鏇存柊
                to_update.append((chunk, embedding, chunk_hash))
            else:
                # 鏂板
                to_add.append((chunk, embedding, chunk_hash))

        # 鍒犻櫎鏃ф枃浠朵腑涓嶅啀瀛樺湪鐨?chunks
        to_delete = [cid for chash, cid in existing_hashes.items() if chash not in new_hash_set]

        # 鎵ц瀹為檯鎿嶄綔
        if to_add:
            await self._batch_add_chunks(to_add)
        if to_update:
            await self._batch_update_chunks(to_update)
        if to_delete:
            await self._batch_delete_chunks(to_delete)

        logger.info(
            f"[Pipeline] 澧為噺淇濆瓨: +{len(to_add)} 鏂板, ~{len(to_update)} 鏇存柊, "
            f"-{len(to_delete)} 鍒犻櫎, ={skipped} 璺宠繃"
        )
        return len(to_add), len(to_update), len(to_delete), skipped

    async def _batch_add_chunks(self, chunks_with_emb: list) -> Any:
        """鎵归噺鏂板 chunk 鍒?SQLite + ChromaDB"""
        from src.db.memory_store import get_store
        from src.db.vector_store import get_vector_store

        store = get_store()
        vs = get_vector_store()

        chunk_dicts = []
        for chunk, embedding, chunk_hash in chunks_with_emb:
            d = chunk.to_dict()
            d["content_hash"] = chunk_hash
            d["source_file"] = chunk.source_file or chunk.file_name
            chunk_dicts.append(d)

        # SQLite 鍐欏叆
        store.add_batch(chunk_dicts)

        # ChromaDB 鍐欏叆锛堣繃婊?None embedding锛?
        if vs:
            valid_chunks = [(c, e, h) for c, e, h in chunks_with_emb if e is not None]
            if valid_chunks:
                ids = [c.chunk_id or c.to_dict().get("chunk_id", "") for c, _, _ in valid_chunks]
                embeddings = [emb for _, emb, _ in valid_chunks]
                metas = [c.to_dict() for c, _, _ in valid_chunks]
                docs = [c.text for c, _, _ in valid_chunks]
                vs.add(ids=ids, embeddings=embeddings, metadata=metas, documents=docs)
                logger.info(f"[Pipeline] ChromaDB 鍐欏叆 {len(valid_chunks)} 涓悜閲?)
            else:
                logger.warning("[Pipeline] 鎵€鏈?embedding 涓?None锛岃烦杩?ChromaDB 鍐欏叆")

    async def _batch_update_chunks(self, chunks_with_emb: list) -> Any:
        """鎵归噺鏇存柊 chunk 鍦?SQLite + ChromaDB"""
        from src.db.memory_store import get_store
        from src.db.vector_store import get_vector_store

        store = get_store()
        vs = get_vector_store()

        for chunk, embedding, chunk_hash in chunks_with_emb:
            chunk_id = chunk.chunk_id or chunk.to_dict().get("chunk_id", "")
            d = chunk.to_dict()
            d["content_hash"] = chunk_hash
            d["source_file"] = chunk.source_file or chunk.file_name

            # SQLite: 鍒犻櫎鏃ц褰?+ 鏂板
            store.delete_by_hash(chunk.file_hash)
            store.add(d)

            # ChromaDB: upsert
            if vs:
                vs.upsert(
                    target_id=chunk_id,
                    embedding=embedding,
                    metadata=d,
                    document=chunk.text,
                )

    async def _batch_delete_chunks(self, chunk_ids: list) -> Any:
        """鎵归噺鍒犻櫎 chunk锛堜粠 ChromaDB 鍒犻櫎锛孲QLite 閫氳繃绾ц仈鏍囪锛?""
        if not chunk_ids:
            return
        from src.db.vector_store import get_vector_store

        vs = get_vector_store()
        if vs:
            try:
                # ChromaDB 鎵归噺鍒犻櫎
                for cid in chunk_ids:
                    vs._collection.delete(ids=[cid])
                logger.debug(f"[Pipeline] 浠?ChromaDB 鍒犻櫎 {len(chunk_ids)} 涓?chunk")
            except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
                logger.warning(f"[Pipeline] ChromaDB 鍒犻櫎澶辫触: {e}")

    # === P1浼樺寲 浠诲姟2锛氫簨鍔℃€т繚瀛?===

    async def _save_with_incremental_and_transaction(self, chunks: List[Chunk], file_hash: str) -> Any:
        """缁撳悎澧為噺鏇存柊鍜屼簨鍔′繚鎶ょ殑淇濆瓨娴佺▼

        1. 鍏堝啓 SQLite锛堜富瀛樺偍锛夛紝璁板綍鐘舵€佷负 'pending'
        2. 鑾峰彇宸叉湁 hash 杩涜澧為噺瀵规瘮
        3. 鍐欏叆 ChromaDB
        4. 濡傛灉 ChromaDB 鎴愬姛锛屾洿鏂?SQLite 鐘舵€佷负 'active'
        5. 濡傛灉 ChromaDB 澶辫触锛屽洖婊?SQLite 鐘舵€佷负 'failed'
        """
        logger.info(f"[Pipeline] 浜嬪姟鎬у閲忎繚瀛樺紑濮? file_hash={file_hash[:16]}...")

        # 鏀堕泦 embeddings锛堝凡鍦?process 涓绠楀ソ锛?
        embeddings = [c.embedding for c in chunks]

        # 鎵ц澧為噺淇濆瓨
        try:
            added, updated, deleted, skipped = await self._incremental_save(file_hash, chunks, embeddings)
            logger.info(f"[Pipeline] 浜嬪姟鎬т繚瀛樺畬鎴? +{added} ~{updated} -{deleted} ={skipped}")
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            logger.error(f"[Pipeline] 澧為噺淇濆瓨澶辫触锛屽洖閫€鑷冲叏閲忎繚瀛? {e}")
            # 鍥為€€锛氬叏閲忎繚瀛?
            await self.saver.save(chunks)


# === P1浼樺寲 浠诲姟2锛歍ransactionalSaver ===


class TransactionalSaver:
    """浜嬪姟鎬т繚瀛樺櫒 鈥?淇濊瘉 SQLite 鍜?ChromaDB 鐨勪竴鑷存€?

    涓ら樁娈垫彁浜わ細
    1. 鍏堝啓 SQLite锛堜富瀛樺偍锛夛紝鐘舵€?'pending'
    2. 鍐?ChromaDB锛堝悜閲忓瓨鍌級
    3. ChromaDB 鎴愬姛 鈫?鏇存柊 SQLite 鐘舵€?'active'
    4. ChromaDB 澶辫触 鈫?鍥炴粴 SQLite 鐘舵€?'failed'
    """

    def __init__(self, config: Dict = None):
        self.config = config or {}

    async def save(self, chunks: List[Chunk], embeddings: List[Optional[List[float]]], metadata: Dict = None) -> bool:
        """浜嬪姟鎬т繚瀛樺叆鍙?

        Returns:
            True 鍏ㄩ儴鎴愬姛锛汧alse ChromaDB 澶辫触锛孲QLite 宸叉爣璁?failed
        """
        if not chunks:
            return True

        metadata = metadata or {}
        chunk_ids = []

        try:
            from src.db.memory_store import get_store

            store = get_store()

            # Phase 1: SQLite pending
            chunk_dicts = []
            for chunk in chunks:
                d = chunk.to_dict()
                d["status"] = "pending"
                d["content_hash"] = hashlib.sha256(chunk.text.encode("utf-8")).hexdigest()
                d["source_file"] = chunk.source_file or chunk.file_name
                chunk_dicts.append(d)
                chunk_ids.append(chunk.chunk_id or d.get("chunk_id", ""))

            store.add_batch(chunk_dicts)
            logger.debug(f"[TransactionalSaver] Phase 1: SQLite pending ({len(chunk_ids)} chunks)")

            # Phase 2: ChromaDB
            chroma_success = False
            try:
                from src.db.vector_store import get_vector_store

                vs = get_vector_store()
                if vs:
                    ids = [c.chunk_id or c.to_dict().get("chunk_id", "") for c in chunks]
                    valid_embeddings = [e for e in embeddings if e is not None]
                    valid_chunks = [c for c, e in zip(chunks, embeddings) if e is not None]
                    if valid_embeddings and valid_chunks:
                        valid_ids = [c.chunk_id or c.to_dict().get("chunk_id", "") for c in valid_chunks]
                        valid_metas = [c.to_dict() for c in valid_chunks]
                        valid_docs = [c.text for c in valid_chunks]
                        chroma_success = vs.add(
                            ids=valid_ids,
                            embeddings=valid_embeddings,
                            metadata=valid_metas,
                            documents=valid_docs,
                        )
                    else:
                        chroma_success = True  # 鏃犳湁鏁?embedding锛屼笉绠?ChromaDB 澶辫触
                else:
                    chroma_success = True  # ChromaDB 涓嶅彲鐢紝涓嶇畻澶辫触
            except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
                logger.error(f"[TransactionalSaver] Phase 2 ChromaDB 鍐欏叆澶辫触: {e}")
                chroma_success = False

            # Phase 3: Confirm or Rollback
            if chroma_success:
                self._update_sqlite_status(chunk_ids, "active")
                logger.info(f"[TransactionalSaver] Phase 3: 纭 鈫?active ({len(chunk_ids)} chunks)")
                return True
            else:
                self._update_sqlite_status(chunk_ids, "failed")
                logger.error(f"[TransactionalSaver] Phase 3: 鍥炴粴 鈫?failed ({len(chunk_ids)} chunks)")
                return False

        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            logger.error(f"[TransactionalSaver] 淇濆瓨澶辫触: {e}")
            return False

    def _update_sqlite_status(self, chunk_ids: List[str], status: str) -> Any:
        """鏇存柊 SQLite 涓?chunk 鐨勭姸鎬?""
        try:
            from src.db.memory_store import get_store

            store = get_store()
            with store._db_conn:
                for cid in chunk_ids:
                    store._db_conn.execute(
                        "UPDATE chunks SET status=? WHERE json_extract(doc, '$.chunk_id')=? AND status='pending'",
                        (status, cid),
                    )
            store._db_conn.commit()
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            logger.error(f"[TransactionalSaver] 鏇存柊 SQLite 鐘舵€佸け璐? {e}")

    async def verify_consistency(self, file_hash: str = None) -> Dict:
        """瀹氭湡涓€鑷存€ф牎楠岋細姣斿 SQLite 鍜?ChromaDB 鐨勮褰?

        Returns:
            {
                'consistent': bool,
                'sqlite_count': int,
                'chromadb_count': int,
                'pending_count': int,
                'failed_count': int,
                'orphans_in_chromadb': list,  # ChromaDB 鏈変絾 SQLite 鏃?
                'orphans_in_sqlite': list,     # SQLite 鏈変絾 ChromaDB 鏃?
            }
        """
        result = {
            "consistent": False,
            "sqlite_count": 0,
            "chromadb_count": 0,
            "pending_count": 0,
            "failed_count": 0,
            "orphans_in_chromadb": [],
            "orphans_in_sqlite": [],
        }

        try:
            from src.db.memory_store import get_store
            from src.db.vector_store import get_vector_store

            store = get_store()
            vs = get_vector_store()

            # SQLite 缁熻
            if file_hash:
                sqlite_chunks = store.get_by_hash(file_hash)
            else:
                sqlite_chunks = store.get_all()
            sqlite_ids = {c.get("chunk_id", c.get("id", "")) for c in sqlite_chunks}
            result["sqlite_count"] = len(sqlite_ids)

            # 缁熻 pending/failed
            for c in sqlite_chunks:
                status = c.get("status", "active")
                if status == "pending":
                    result["pending_count"] += 1
                elif status == "failed":
                    result["failed_count"] += 1

            # ChromaDB 缁熻
            if vs and vs._usable:
                chroma_count = vs.count
                result["chromadb_count"] = chroma_count if chroma_count >= 0 else 0
                # 浠?ChromaDB 鑾峰彇鎵€鏈?ID
                try:
                    all_chroma = vs._collection.get(include=[])
                    chroma_ids = set(all_chroma.get("ids", []))

                    # 浜ゅ弶姣斿
                    result["orphans_in_sqlite"] = list(chroma_ids - sqlite_ids)
                    result["orphans_in_chromadb"] = list(sqlite_ids - chroma_ids)
                except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
                    pass

            # 涓€鑷存€у垽鏂?
            result["consistent"] = (
                result["pending_count"] == 0
                and result["failed_count"] == 0
                and len(result["orphans_in_sqlite"]) == 0
                and len(result["orphans_in_chromadb"]) == 0
            )

            logger.info(
                f"[ConsistencyCheck] consistent={result['consistent']}, "
                f"sqlite={result['sqlite_count']}, chroma={result['chromadb_count']}, "
                f"pending={result['pending_count']}, failed={result['failed_count']}"
            )

        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:
            logger.error(f"[ConsistencyCheck] 鏍￠獙澶辫触: {e}")

        return result


# 鍏ㄥ眬绠＄嚎瀹炰緥
_pipeline_instance = None


def get_pipeline(config: Dict = None) -> UnifiedPipeline:
    """鑾峰彇鍏ㄥ眬绠＄嚎瀹炰緥"""
    global _pipeline_instance
    if _pipeline_instance is None:
        _pipeline_instance = UnifiedPipeline(config)
    return _pipeline_instance

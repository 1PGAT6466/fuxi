"""
parsers.py — 统一解析器
处理各种文件格式的解析逻辑。
"""
from pathlib import Path
from typing import Dict

from src.pipeline.errors import ParseError


class UnifiedParser:
    """统一解析器 — 合并 stomach.py + ingest.py 的解析逻辑"""

    # 支持的文件格式
    SUPPORTED_TYPES = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".doc": "docx",
        ".xlsx": "xlsx",
        ".xls": "xlsx",
        ".txt": "text",
        ".md": "markdown",
        ".html": "html",
        ".htm": "html",
        ".csv": "csv",
        ".json": "json",
        ".pptx": "pptx",
        ".ppt": "pptx",  # 添加 .ppt 支持
    }

    def __init__(self, config: Dict = None):
        self.config = config or {}

    def parse(self, file_path: str) -> Dict:
        """解析文件，返回 {text, tables, metadata}"""
        path = Path(file_path)
        if not path.exists():
            raise ParseError(f"文件不存在: {file_path}")

        ext = path.suffix.lower()
        try:
            if ext == ".pdf":
                return self._parse_pdf(file_path)
            elif ext in (".docx", ".doc"):
                return self._parse_docx(file_path)
            elif ext in (".xlsx", ".xls"):
                return self._parse_excel(file_path)
            elif ext in (".pptx", ".ppt"):  # 修改以支持 .ppt
                return self._parse_pptx(file_path)
            elif ext == ".txt":
                return self._parse_text(file_path)
            elif ext == ".md":
                return self._parse_markdown(file_path)
            elif ext == ".csv":
                return self._parse_csv(file_path)
            elif ext in (".json",):
                return self._parse_json(file_path)
            elif ext in (".html", ".htm"):
                return self._parse_html(file_path)
            else:
                return self._parse_text(file_path)
        except ParseError:
            raise
        except Exception as e:  # TODO: Narrow exception type
            raise ParseError(f"解析失败 ({ext}): {e}")

    def _parse_pdf(self, file_path: str) -> Dict:
        """PDF 解析 — 合并 stomach.py + ingest.py 逻辑"""
        text = ""
        tables = []

        # 方式1: fitz (PyMuPDF) — 中文最优
        try:
            import fitz
            doc = fitz.open(file_path)
            pages_text = []
            for page in doc:
                pages_text.append(page.get_text())
            text = "\n".join(pages_text)
            doc.close()
            if text.strip():
                return {"text": text, "tables": tables, "metadata": {"parser": "fitz"}}
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
            pass  # 静默：Exception 失败不影响主流程

        # 方式2: pdfplumber — 表格提取
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                pages_text = []
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        pages_text.append(page_text)
                    # 提取表格
                    for table in page.extract_tables():
                        if table:
                            tables.append({"headers": table[0] if table else [], "rows": table[1:] if len(table) > 1 else []})
                text = "\n".join(pages_text)
                if text.strip():
                    return {"text": text, "tables": tables, "metadata": {"parser": "pdfplumber"}}
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
            pass  # 静默：Exception 失败不影响主流程

        # 方式3: pypdf — 兜底
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            pages_text = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    pages_text.append(page_text)
            text = "\n".join(pages_text)
            return {"text": text, "tables": tables, "metadata": {"parser": "pypdf"}}
        except Exception as e:  # TODO: Narrow exception type
            raise ParseError(f"PDF解析失败: {e}")

    def _parse_docx(self, file_path: str) -> Dict:
        """DOCX 解析"""
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            tables = []
            for table in doc.tables:
                headers = [cell.text for cell in table.rows[0].cells] if table.rows else []
                rows = [[cell.text for cell in row.cells] for row in table.rows[1:]] if len(table.rows) > 1 else []
                tables.append({"headers": headers, "rows": rows})
            return {"text": "\n".join(paragraphs), "tables": tables, "metadata": {"parser": "python-docx"}}
        except Exception as e:  # TODO: Narrow exception type
            raise ParseError(f"DOCX解析失败: {e}")

    def _parse_excel(self, file_path: str) -> Dict:
        """Excel 解析"""
        try:
            import pandas as pd
            df = pd.read_excel(file_path)
            text = df.to_string(index=False)
            tables = [{"headers": list(df.columns), "rows": df.values.tolist()}]
            return {"text": text, "tables": tables, "metadata": {"parser": "pandas"}}
        except Exception as e:  # TODO: Narrow exception type
            raise ParseError(f"Excel解析失败: {e}")

    def _parse_text(self, file_path: str) -> Dict:
        """纯文本解析"""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            return {"text": text, "tables": [], "metadata": {"parser": "text"}}
        except Exception as e:  # TODO: Narrow exception type
            raise ParseError(f"文本解析失败: {e}")

    def _parse_markdown(self, file_path: str) -> Dict:
        """Markdown 解析 — 保留标题层级结构"""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                raw = f.read()
        except Exception as e:  # TODO: Narrow exception type
            raise ParseError(f"Markdown读取失败: {e}")

        # 尝试使用 mistune 解析结构
        structure = []
        try:
            import mistune
            md_parser = mistune.create_markdown(renderer=mistune.AstRenderer())
            ast = md_parser(raw)
            structure = self._extract_md_structure(ast)
        except ImportError:
            pass  # 静默：ImportError 失败不影响主流程

        # 兜底：正则提取标题层级
        if not structure:
            structure = self._regex_md_structure(raw)

        return {
            "text": raw,
            "tables": [],
            "metadata": {
                "parser": "markdown",
                "heading_structure": structure,
            },
        }

    def _extract_md_structure(self, ast) -> list:
        """从 mistune AST 提取标题结构"""
        import re
        headings = []
        for node in ast:
            if node.get("type") == "heading":
                headings.append({
                    "level": node.get("attrs", {}).get("level", 1),
                    "text": self._strip_markdown_text(node),
                })
        return headings

    def _strip_markdown_text(self, node) -> str:
        """从 AST 节点提取纯文本"""
        if not node:
            return ""
        if isinstance(node, str):
            return node
        if isinstance(node, dict):
            if node.get("type") == "text":
                return node.get("raw", "")
            children = node.get("children", [])
            return "".join(self._strip_markdown_text(c) for c in children)
        if isinstance(node, list):
            return "".join(self._strip_markdown_text(c) for c in node)
        return ""

    def _regex_md_structure(self, raw: str) -> list:
        """正则兜底：从 Markdown 文本提取标题"""
        import re
        headings = []
        for line in raw.split("\n"):
            m = re.match(r"^(#{1,6})\s+(.+)", line.strip())
            if m:
                headings.append({
                    "level": len(m.group(1)),
                    "text": m.group(2).strip(),
                })
        return headings

    def _parse_pptx(self, file_path: str) -> Dict:
        """PPTX/PPT 解析 — 支持 .pptx 和 .ppt 格式"""
        import os

        # 检测文件格式
        is_old_format = False
        try:
            import olefile
            is_old_format = olefile.isOleFile(file_path)
        except ImportError:
            pass
        
        # 如果是旧版 .ppt 格式，尝试转换为 .pptx
        pptx_path = file_path
        temp_file = None
        
        if is_old_format:
            converted = self._convert_ppt_to_pptx(file_path)
            if converted:
                pptx_path = converted
                temp_file = converted
            else:
                # 转换失败，回退到二进制解析
                return self._parse_ppt_binary_fallback(file_path)
        
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            slides_text = []
            for slide in prs.slides:
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_content.append(shape.text)
                slides_text.append("\n".join(slide_content))
            return {
                "text": "\n\n".join(slides_text),
                "tables": [],
                "metadata": {"parser": "python-pptx", "pages": len(prs.slides)},
            }
        except Exception as e:
            # 如果 python-pptx 失败，回退到二进制解析
            if is_old_format:
                return self._parse_ppt_binary_fallback(file_path)
            raise ParseError(f"PPTX解析失败: {e}")
        finally:
            # 清理临时文件
            if temp_file and os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except Exception:
                    pass
    
    def _convert_ppt_to_pptx(self, ppt_path: str) -> str:
        """使用 PowerPoint 将 .ppt 转换为 .pptx"""
        import os
        import tempfile
        
        try:
            import win32com.client

            # 创建临时目录
            tmpdir = tempfile.mkdtemp()
            pptx_path = os.path.join(tmpdir, os.path.splitext(os.path.basename(ppt_path))[0] + ".pptx")
            
            # 使用 PowerPoint 转换
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = False
            
            try:
                presentation = powerpoint.Presentations.Open(ppt_path)
                presentation.SaveAs(pptx_path, 24)  # 24 = ppSaveAsOpenXMLPresentation
                presentation.Close()
            finally:
                powerpoint.Quit()
            
            if os.path.exists(pptx_path):
                logger.info(f"[Parser] PPT 转换成功: {ppt_path} -> {pptx_path}")
                return pptx_path
        except ImportError:
            logger.warning("[Parser] pywin32 未安装，无法转换 PPT 文件")
        except Exception as e:
            logger.warning(f"[Parser] PowerPoint 转换失败: {e}")
        
        return None
    
    def _parse_ppt_binary_fallback(self, file_path: str) -> Dict:
        """旧版 .ppt 文件二进制解析（回退方案）"""
        try:
            import olefile
            ole = olefile.OleFileIO(file_path)
            
            if ole.exists('PowerPoint Document'):
                data = ole.openstream('PowerPoint Document').read()
                text = self._extract_text_from_ppt_binary(data)
                ole.close()
                return {
                    "text": text,
                    "tables": [],
                    "metadata": {"parser": "olefile-binary", "pages": 0},
                }
            ole.close()
        except Exception as e:
            logger.warning(f"[Parser] 二进制解析失败: {e}")
        
        raise ParseError("PPT解析失败: 无法解析此文件格式")

    def _extract_text_from_ppt_binary(self, data: bytes) -> str:
        """从旧版 .ppt 二进制数据中提取文本"""
        import re

        # 尝试多种编码
        encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
        text = ""
        
        for encoding in encodings:
            try:
                text = data.decode(encoding, errors='ignore')
                break
            except Exception:
                continue
        
        # 移除控制字符
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
        
        # 提取中文句子
        chinese_sentences = []
        for match in re.finditer(r'[\u4e00-\u9fff]+[^\u4e00-\u9fff]*[\u4e00-\u9fff]+', text):
            sentence = match.group().strip()
            if len(sentence) > 5:  # 只保留长度大于 5 的句子
                chinese_sentences.append(sentence)
        
        # 提取英文句子
        english_sentences = []
        for match in re.finditer(r'[a-zA-Z]+[^a-zA-Z]*[a-zA-Z]+', text):
            sentence = match.group().strip()
            if len(sentence) > 10:  # 只保留长度大于 10 的句子
                english_sentences.append(sentence)
        
        # 合并句子
        all_sentences = chinese_sentences + english_sentences
        
        # 如果没有找到句子，返回空字符串
        if not all_sentences:
            return ""
        
        # 返回前 100 个句子
        return '\n'.join(all_sentences[:100])

    def _parse_csv(self, file_path: str) -> Dict:
        """CSV 解析"""
        try:
            import pandas as pd
            df = pd.read_csv(file_path)
            text = df.to_string(index=False)
            tables = [{"headers": list(df.columns), "rows": df.values.tolist()}]
            return {"text": text, "tables": tables, "metadata": {"parser": "pandas"}}
        except Exception as e:  # TODO: Narrow exception type
            raise ParseError(f"CSV解析失败: {e}")

    def _parse_json(self, file_path: str) -> Dict:
        """JSON 解析"""
        try:
            import json
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            text = json.dumps(data, ensure_ascii=False, indent=2)
            return {"text": text, "tables": [], "metadata": {"parser": "json"}}
        except Exception as e:  # TODO: Narrow exception type
            raise ParseError(f"JSON解析失败: {e}")

    def _parse_html(self, file_path: str) -> Dict:
        """HTML 解析"""
        try:
            from bs4 import BeautifulSoup
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            return {"text": text, "tables": [], "metadata": {"parser": "beautifulsoup"}}
        except Exception as e:  # TODO: Narrow exception type
            raise ParseError(f"HTML解析失败: {e}")

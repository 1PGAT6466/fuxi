"""
ingest.py - 文本提取、分类、分块、清洗模块
从 server.py 分离，供主路由导入使用
"""
import re, html, os
from pathlib import Path
import logging; logger = logging.getLogger(__name__)

try:
    import jieba
    jieba.setLogLevel(20)
except ImportError:
    jieba = None


def _sanitize_filename(fn: str) -> str:
    # 去除非法字符 + 截断到 180 字符（留 buffer 给 file_hash 前缀）
    safe = re.sub(r'[<>:"/\\|?*]', '_', fn).replace("..", "_").strip()
    # 分离名称和扩展名，分别截断
    name, dot, ext = safe.rpartition('.')
    if len(safe) > 180:
        # 保留前 120 字符 + 扩展名后 40 字符
        name = name[:120]
        ext = ext[-40:] if len(ext) > 40 else ext
        safe = f"{name}.{ext}"
    return safe


def _classify_text(text: str, file_type: str = "") -> str:
    """统一分类：基于文件名+扩展名+内容的加权打分（C1-C4 内联版本）"""
    t = text.lower() if text else ""
    ft = (file_type or "").lower()
    # 快速规则：扩展名直接映射
    EXT_MAP = {
        ".cfg": "IT网络", ".conf": "IT网络", ".ini": "IT网络",
        ".stp": "模具设计", ".step": "模具设计", ".dwg": "模具设计", ".dxf": "模具设计",
        ".awl": "电气自动化", ".scl": "电气自动化",
        ".xlsx": "供应商管理", ".xls": "供应商管理",
    }
    if ft in EXT_MAP:
        return EXT_MAP[ft]
    # 内容关键词匹配
    if any(k in t for k in ["vlan","子网","网关","dhcp","路由","交换机","拓扑","interface","trunk"]):
        return "IT网络"
    if any(k in t for k in ["模具","导柱","导套","顶针","滑块","浇口","分型面","模架","型腔"]):
        return "模具设计"
    if any(k in t for k in ["plc","传感器","伺服","变频","电磁阀","气缸","profinet","npn","pnp"]):
        return "电气自动化"
    if any(k in t for k in ["三坐标","蔡司","contura","calypso","位置度","平面度","gr&r","cpk"]):
        return "品质测量"
    if any(k in t for k in ["采购","报价","合同","供应商","交货","付款"]):
        return "供应商管理"
    if any(k in t for k in ["注塑","成型","模温","料筒","干燥","收缩率","lcp","pa66","pbt","pom"]):
        return "工程技术规范"
    if any(k in t for k in ["考勤","请假","薪资","人事","报销","制度"]):
        return "行政人事"
    if any(k in t for k in ["财务","预算","税务","审计","发票"]):
        return "财务文档"
    return "通用办公"


def _audit_text(text: str, sensitive_patterns: list = None) -> tuple:
    if sensitive_patterns is None:
        sensitive_patterns = []
    safe = text
    flags = []
    for pat in sensitive_patterns:
        if pat.search(safe):
            flags.append("检测到敏感信息")
            safe = pat.sub("[已脱敏]", safe)
    return safe, flags


def _clean_text(raw: str) -> str:
    """
    数据清洗（v13.0 增强）：
    - 去除 HTML 实体编码
    - 去除 URL 链接
    - 去除不可见控制字符
    - 去除常见页眉页脚（页码、路径、邮件签名）
    - 去除连续空行
    """
    raw = html.unescape(raw)
    raw = re.sub(r'https?://\S+', '', raw)
    raw = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', raw)
    
    # 页眉页脚清理
    raw = re.sub(r'^\s*第\s*\d+\s*页\s*(共\s*\d+\s*页)?\s*$', '', raw, flags=re.MULTILINE)
    raw = re.sub(r'^\s*Page\s*\d+\s*(of\s*\d+)?\s*$', '', raw, flags=re.MULTILINE | re.IGNORECASE)
    raw = re.sub(r'^\s*\d{4}-\d{2}-\d{2}\s*\d{1,2}:\d{2}(:\d{2})?\s*$', '', raw, flags=re.MULTILINE)  # 独立时间戳行
    raw = re.sub(r'^\s*[A-Z]:\\[^\n]{10,80}\s*$', '', raw, flags=re.MULTILINE)  # 文件路径行
    # 邮件签名
    raw = re.sub(r'--+\.?\s*$', '', raw, flags=re.MULTILINE)  # -- 分隔线
    raw = re.sub(r'(Best Regards|Sincerely|此致|敬礼|顺祝商祺)[\s\S]{0,200}$', '', raw, flags=re.IGNORECASE)
    # 乱码字符（连续 3+ 个非中英数字符）
    raw = re.sub(r'[^\w\s\u4e00-\u9fff\u3000-\u303f，。！？；：“”‘’（）【】《》、…—·\-\+\.\/]{3,}', ' ', raw)
    
    raw = re.sub(r'\n{4,}', '\n\n\n', raw)
    return raw.strip()


def _generate_summary(text: str, max_len: int = 200) -> str:
    first_para = ""
    for para in text.split("\n\n"):
        para = para.strip()
        if len(para) > 30 and not para.startswith("#"):
            first_para = para[:max_len]
            break
    if not first_para:
        first_para = text[:max_len].strip()
    keywords = []
    try:
        if jieba:
            import jieba.analyse
            keywords = jieba.analyse.extract_tags(text, topK=5)
    except Exception:
        logger.warning(f"[ingest] suppressed exception", exc_info=True)
        pass
    kw_str = "、".join(keywords[:5]) if keywords else ""
    summary = f"[文档摘要] {first_para}"
    if kw_str:
        summary += f"\n[关键词] {kw_str}"
    return summary


# v15.0: 使用 services/chunking.py 的 Markdown-AST 智能分块
def _smart_chunk(text: str, size: int = 1200, overlap: int = 100) -> list:
    """智能分块"""
    if not text or len(text) < 50:
        return [text] if text and len(text.strip()) > 10 else []
    
    chunks = []
    start = 0
    text_len = len(text)
    
    while start < text_len:
        end = min(start + size, text_len)
        
        # 尝试在句号、换行处断开
        if end < text_len:
            for sep in ['\n\n', '\n', '。', '；', '.', ';']:
                last_sep = text.rfind(sep, start + size // 2, end)
                if last_sep > start:
                    end = last_sep + len(sep)
                    break
        
        chunk = text[start:end].strip()
        if chunk and len(chunk) > 20:
            # v4.0 表格结构化提取
            structured = None
            try:
                from src.services.table_parser import extract_tables_from_markdown
                if "|" in chunk:
                    tables = extract_tables_from_markdown(chunk)
                    if tables:
                        structured = tables[0]
            except Exception:
                logger.debug("[suppressed] structured = tables[0]")
                pass
            if structured:
                chunks.append({"text": chunk, "structured_table": structured})
            else:
                chunks.append(chunk)
        
        start = end - overlap if end < text_len else text_len
    
    return chunks

def _extract_pdf_dual(file_path: str) -> str:
    """
    v10.0: PDF 双轨解析
    - pdfplumber 主力（保留表格/双栏/页码信息）
    - PyPDF2 回退（兼容旧 PDF）
    """
    lines = []
    
    # 路0: fitz (PyMuPDF) - 中文最优，仅对 < 50MB 文件使用
    file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
    if file_size_mb < 50:
        try:
            import fitz
            doc = fitz.open(file_path)
            total = doc.page_count
            for i in range(total):
                try:
                    page = doc[i]
                    text = page.get_text()
                    if text and text.strip():
                        lines.append("[Page %d/%d]\n%s" % (i+1, total, text.strip()))
                except Exception:
                    logger.warning(f"[ingest] suppressed exception", exc_info=True)
                    pass
            doc.close()
            if lines:
                return "\n".join(lines)
        except Exception:
            logger.warning(f"[ingest] suppressed exception", exc_info=True)
            pass
    
    # 路1: pdfplumber
    try:
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            total = len(pdf.pages)
            for i, page in enumerate(pdf.pages):
                try:
                    # 提取文本（保留布局）
                    text = page.extract_text() or ""
                    # 提取表格（结构化标注）
                    from src.services.multimodal import enhance_table_extraction
                    combined = enhance_table_extraction(text, tables)
                    if combined.strip():
                        lines.append(f"[Page {i+1}/{total}]\n{combined.strip()}")
                except Exception:
                    # 单页失败，回退
                    try:
                        from PyPDF2 import PdfReader
                        reader = PdfReader(file_path)
                        if i < len(reader.pages):
                            txt = reader.pages[i].extract_text() or ""
                            if txt.strip():
                                lines.append(f"[Page {i+1}/{total}]\n{txt}")
                    except Exception:
                        lines.append(f"[Page {i+1}/{total}] skipped")
        if lines:
            return "\n".join(lines)
    except Exception:
        logger.warning(f"[ingest] suppressed exception", exc_info=True)
        pass
    
    # 路2: PyPDF2 回退
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        total = len(reader.pages)
        for i, page in enumerate(reader.pages):
            try:
                txt = page.extract_text() or ""
                if txt.strip():
                    lines.append(f"[Page {i+1}/{total}]\n{txt}")
            except Exception:
                lines.append(f"[Page {i+1}/{total}] skipped")
        return "\n".join(lines)
    except Exception as e:
        return f"[PDF 解析失败: {e}]"
    
    return ""


def _compute_file_hash(file_path: str) -> str:
    """
    v10.0: 计算文件 SHA256（MD5 去重用）
    """
    import hashlib
    sha = hashlib.sha256()
    with open(file_path, "rb") as f:
        while chunk := f.read(8192):
            sha.update(chunk)
    return sha.hexdigest()


def _extract_text(file_path: str, ext: str) -> str:
    path = Path(file_path)
    try:
        if ext in [".txt",".md",".csv",".cfg",".log",".ini",".conf",".json",".xml",".html",".htm",
                  ".py",".js",".ts",".java",".c",".cpp",".h",".sh",".bat",".ps1",".yaml",".yml"]:
            return path.read_text(encoding="utf-8", errors="ignore")
        elif ext == ".docx":
            from docx import Document
            return "\n".join(p.text for p in Document(str(path)).paragraphs if p.text.strip())
        elif ext == ".doc":
            # 老式 .doc 二进制格式 → antiword 转换
            import subprocess
            try:
                result = subprocess.run(
                    ["antiword", "-m", "UTF-8.txt", str(path)],
                    capture_output=True, text=True, timeout=30
                )
                if result.returncode == 0 and result.stdout.strip():
                    return result.stdout.strip()
            except Exception:
                logger.warning(f"[ingest] suppressed exception", exc_info=True)
                pass
            # 回退: .doc 也可能是 docx 格式，尝试用 python-docx 打开
            try:
                from docx import Document
                return "\n".join(p.text for p in Document(str(path)).paragraphs if p.text.strip())
            except Exception:
                return ""
        elif ext == ".wps":
            # WPS 格式尝试用 python-docx（新版 WPS 可能是 docx 格式）
            try:
                from docx import Document
                return "\n".join(p.text for p in Document(str(path)).paragraphs if p.text.strip())
            except Exception:
                return ""
        elif ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            lines = []
            for sn in wb.sheetnames:
                ws = wb[sn]
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue
                lines.append(f"[Sheet: {sn}] 行数={len(rows)}")
                headers = [str(c) if c is not None else "" for c in rows[0]]
                lines.append(f"[列名] {' | '.join(headers)}")
                for row in rows[1:]:
                    vals = []
                    for j, c in enumerate(row):
                        if c is None:
                            continue
                        col_name = headers[j] if j < len(headers) else f"Col{j}"
                        vals.append(f"{col_name}={c}")
                    if vals:
                        lines.append(" | ".join(vals))
            wb.close()
            return "\n".join(lines)
        elif ext == ".xls":
            import xlrd
            wb = xlrd.open_workbook(str(path))
            lines = []
            for sn in wb.sheet_names():
                ws = wb.sheet_by_name(sn)
                if ws.nrows == 0:
                    continue
                lines.append(f"[Sheet: {sn}] 行数={ws.nrows}")
                headers = [str(ws.cell_value(0, j)) for j in range(ws.ncols)]
                lines.append(f"[列名] {' | '.join(headers)}")
                for r in range(1, ws.nrows):
                    vals = []
                    for j in range(ws.ncols):
                        v = ws.cell_value(r, j)
                        if v == "" or v is None:
                            continue
                        col_name = headers[j] if j < len(headers) else f"Col{j}"
                        vals.append(f"{col_name}={v}")
                    if vals:
                        lines.append(" | ".join(vals))
            return "\n".join(lines)
        elif ext == ".pdf":
            # v10.0: 双轨 PDF 解析 — pdfplumber 主力 + PyPDF2 回退
            return _extract_pdf_dual(str(path))
        elif ext in [".pptx",".ppt"]:
            # python-pptx for .pptx, olefile fallback for .ppt
            try:
                from pptx import Presentation
                slides = []
                for i, sl in enumerate(Presentation(str(path)).slides):
                    t = [f"[Slide {i+1}]"]
                    for sh in sl.shapes:
                        if sh.has_text_frame:
                            t.append(sh.text_frame.text)
                    if sl.has_notes_slide and sl.notes_slide.notes_text_frame:
                        notes = sl.notes_slide.notes_text_frame.text.strip()
                        if notes:
                            t.append("[Notes] " + notes)
                    slides.append("\n".join(t))
                result = "\n\n".join(slides)
                if len(result.strip()) > 50:
                    return result
            except Exception:
                logger.warning(f"[ingest] suppressed exception", exc_info=True)
                pass
            # olefile fallback for legacy .ppt
            try:
                import olefile, re
                ole = olefile.OleFileIO(str(path))
                texts = []
                for stream in ole.listdir():
                    try:
                        raw = ole.openstream(stream).read()
                        try:
                            t = raw.decode("utf-16-le", errors="ignore")
                        except Exception:
                            continue
                        t = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", t)
                        t = t.strip()
                        if len(t) > 20:
                            texts.append(t)
                    except Exception:
                        logger.warning(f"[ingest] suppressed exception", exc_info=True)
                        pass
                ole.close()
                result = "\n\n".join(texts)
                if len(result.strip()) > 100:
                    return result
            except Exception:
                logger.warning(f"[ingest] suppressed exception", exc_info=True)
                pass
            return ""
        elif ext == ".zip":
            import zipfile, tempfile, shutil
            lines = []
            try:
                with zipfile.ZipFile(str(path), 'r') as zf:
                    for name in zf.namelist():
                        if name.endswith('/'):
                            continue
                        inner_ext = os.path.splitext(name)[1].lower()
                        # 只处理可读文件格式
                        supported = ['.txt','.md','.csv','.log','.json','.xml','.html','.htm','.cfg','.ini','.conf','.docx','.doc','.xlsx','.xls','.pdf','.pptx','.ppt']
                        if inner_ext not in supported:
                            continue
                        try:
                            with zf.open(name) as inner_f:
                                content = inner_f.read()
                            # 写入临时文件用 _extract_text 递归解析
                            tmp_dir = tempfile.mkdtemp()
                            tmp_path = os.path.join(tmp_dir, os.path.basename(name))
                            with open(tmp_path, 'wb') as tmp_f:
                                tmp_f.write(content)
                            inner_text = _extract_text(tmp_path, inner_ext)
                            shutil.rmtree(tmp_dir, ignore_errors=True)
                            if inner_text.strip():
                                lines.append(f"[ZIP内: {name}]\n{inner_text[:2000]}")
                        except Exception:
                            lines.append(f"[ZIP内: {name}] 提取失败")
                return "\n\n---\n\n".join(lines[:20])  # 最多 20 个文件
            except Exception as e:
                return f"[ZIP提取失败: {e}]"
        elif ext in [".deb",".dwg",".dxf",".stp",".step",".igs",".iges",
                    ".jpg",".jpeg",".png",".gif",".bmp",".svg",
                    ".exe",".msi",".apk",".dmg",".pkg",".rpm",".7z",".rar",".tar.gz",".gz",
                    ".bin",".iso",".img",".dll",".so",".o",".a",".lib",
                    ".mp3",".mp4",".avi",".mkv",".wav",".flac",".mov",".wmv"]:
            # 二进制/图纸/图片/安装包格式 — 仅记录文件名，无文本
            return f"[文件: {path.name}] (二进制格式，无文本提取)"
        else:
            # 未知格式：尝试作为 UTF-8 文本读取
            try:
                return path.read_text(encoding="utf-8", errors="replace")
            except Exception:
                return f"[文件: {path.name}] (未知格式)"
    except Exception as e:
        return f"[提取失败: {e}]"
    # === 通用降级解析器：尝试 raw text + 二进制扫描 ===
    try:
        with open(str(path), "rb") as f:
            raw = f.read()
    except Exception:
        return ""
    
    # 1. 尝试 UTF-8 纯文本
    try:
        text = raw.decode("utf-8")
        if len(text.strip()) > 50:
            return text
    except Exception:
        logger.warning(f"[ingest] suppressed exception", exc_info=True)
        pass
    
    # 2. 尝试 UTF-16 LE
    try:
        text = raw.decode("utf-16-le")
        import re
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", text)
        if len(text.strip()) > 100:
            return text
    except Exception:
        logger.warning(f"[ingest] suppressed exception", exc_info=True)
        pass
    
    # 3. 尝试 latin-1 + 提取可读行（二进制格式最后手段）
    try:
        text = raw.decode("latin-1")
        import re
        lines = text.split("\n")
        readable = []
        for line in lines:
            alpha_ratio = sum(1 for c in line if c.isalpha() or c.isspace()) / max(len(line), 1)
            if alpha_ratio > 0.6 and len(line.strip()) > 10:
                readable.append(line.strip())
        if len(readable) > 3:
            return "\n".join(readable)
    except Exception:
        logger.warning(f"[ingest] suppressed exception", exc_info=True)
        pass
    
    # === 智能自愈层：自动搜索并安装缺失的解析器 ===
    ext_str = str(path).lower()
    if "." in ext_str:
        ext = ext_str.rsplit(".", 1)[-1]
        import subprocess as _sp
        try:
            _dn = _sp.DEVNULL if hasattr(_sp, 'DEVNULL') else open('/dev/null', 'w')
        except Exception:
            _dn = None
        known_map = {
            "msg": ["extract-msg", "msg-parser"],
            "eml": ["mail-parser"],
            "wps": ["python-docx"],
            "rtf": ["striprtf", "pyth"],
            "ods": ["odfpy", "ezodf"],
            "odt": ["odfpy", "python-docx"],
            "epub": ["ebooklib", "epub2txt"],
            "djvu": ["djvulibre-python"],
            "ps": ["ghostscript"],
            "ai": ["pdf2image"],
            "cdr": ["pdf2image"],
            "one": ["python-docx"],
            "vsd": ["python-docx"],
            "vsdx": ["python-docx"],
            "mpp": ["python-docx"],
            "pub": ["python-docx"],
        }
        candidates = known_map.get(ext, [])
        if candidates:
            import importlib as _il
            installed = None
            for pkg in candidates:
                try:
                    _il.import_module(pkg.replace("-", "_"))
                    installed = pkg
                    break
                except ImportError:
                    continue
            if not installed:
                logger.warning(f"[ingest] 缺少解析依赖 {candidates[0]} (ext={ext})，请手动安装: pip install {candidates[0]}")
            if installed:
                return _extract_text(path, ext)
    return ""


# === merged from ingestion.py ===
"""
伏羲 Fuxi · 统一入库引擎
========================
将解析器输出的标准化文档 → 按类型分块 → 向量化 → 写入 ChromaDB

流程:
  parse_result → type_router → chunker → embedder → ChromaDB
             ↘ table_extractor → kb_tables 独立索引
"""
import os
import re
import json
import uuid
import hashlib
import logging
import asyncio
from typing import List, Dict, Optional
from pathlib import Path

logger = logging.getLogger(__name__)

# ============================================================
# 语义分块
# ============================================================



def smart_chunk_semantic(text: str, chunk_size: int = 1500, overlap: int = 200) -> List[str]:
    """语义分块：按段落边界切割，保持语义完整性"""
    if not text:
        return []
    
    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
    chunks = []
    current = ""
    
    for p in paragraphs:
        # 如果是标题行（包含 # 或全部大写等），新起一个 chunk
        is_heading = p.startswith('#') or p.startswith('【') or p.startswith('[')
        
        if is_heading and current and len(current) > 100:
            chunks.append(current.strip())
            current = p
        elif len(current) + len(p) > chunk_size:
            chunks.append(current.strip())
            # 重叠：保留上一段落的最后内容
            if overlap and current:
                overlap_text = current[-overlap:] if len(current) > overlap else current[-50:]
                current = overlap_text + '\n\n' + p
            else:
                current = p
        else:
            if current:
                current += '\n\n' + p
            else:
                current = p
    
    if current.strip():
        chunks.append(current.strip())
    
    return chunks


def chunk_table(table_data: Dict) -> List[str]:
    """表格不分块，整表保留为一个 chunk"""
    if not table_data:
        return []
    
    chunks = []
    if table_data.get('markdown'):
        chunks.append(table_data['markdown'])
    else:
        # 生成 Markdown 表示
        header = table_data.get('header', [])
        rows = table_data.get('all_rows', table_data.get('sample_rows', []))
        if header:
            lines = ['| ' + ' | '.join(str(h) for h in header) + ' |']
            lines.append('|' + '|'.join(['---' for _ in header]) + '|')
            for row in rows:
                lines.append('| ' + ' | '.join(str(c) for c in row) + ' |')
            chunks.append('\n'.join(lines))
    
    return chunks


def chunk_image(file_path: str, ocr_text: str = "") -> List[str]:
    """图片：多模态转录 + OCR 文字
    优先使用 Vision 模型转录图片内容，fallback 到 OCR，最终 fallback 到文件名
    """
    # 如果已有 OCR 文本且足够长，直接用
    if ocr_text and len(ocr_text) > 50:
        return [ocr_text]
    
    # 尝试多模态 Vision 模型转录
    try:
        from src.services.multimodal import transcribe_image
        transcription = transcribe_image(file_path)
        if transcription and len(transcription) > 20:
            logger.info(f"[multimodal] 图片转录成功: {Path(file_path).name} ({len(transcription)}字)")
            return [f"[图片内容] {transcription}"]
    except Exception as e:
        logger.debug(f"[multimodal] 图片转录跳过: {e}")
    
    # Fallback: OCR 文本
    if ocr_text and len(ocr_text) > 10:
        return [ocr_text]
    
    # 最终 Fallback: 文件名占位
    return [f"[图片: {Path(file_path).name}]"]


# ============================================================
# 入库引擎
# ============================================================

async def ingest_document(
    parse_result: Dict,
    file_name: str = "",
    category: str = "",
    embed_fn=None,
    vector_store=None,
    table_store=None,
    memory_store=None
) -> Dict:
    """
    统一入库一个文档
    
    参数:
      parse_result: 解析器输出 {type, text, metadata, tables, images}
      file_name: 原始文件名
      category: 文档分类
      embed_fn: 向量化函数
      vector_store: ChromaDB kb_chunks 集合
      table_store: ChromaDB kb_tables 集合
      memory_store: BM25 全文索引
    
    返回: {chunks_added, tables_indexed, errors}
    """
    doc_type = parse_result.get('type', 'unknown')
    text = parse_result.get('text', '')
    metadata = parse_result.get('metadata', {})
    tables = parse_result.get('tables', [])
    images = parse_result.get('images', [])
    
    file_hash = hashlib.md5((file_name + str(metadata)).encode()).hexdigest()[:16]
    
    # 分类：如果没传入 category 或为"通用办公"，用 match_category 重新分类
    if not category or category == "通用办公":
        try:
            from src.category_registry import match_category as _match_cat
            ext = os.path.splitext(file_name)[1].lower() if file_name else ""
            _cat = _match_cat(text[:5000], file_ext=ext, file_name=file_name)
            if _cat:
                category = _cat
        except Exception:
            logger.debug("[suppressed] category = _cat")
            pass
    
    # 防御性校验：category 不能是 Python repr() 格式
    if category and ("[{" in category or "': " in category):
        import re
        m = re.search(r"'category':\s*'([^']+)'", category)
        category = m.group(1) if m else "通用办公"
    
    chunks = []
    result = {"chunks_added": 0, "tables_indexed": 0, "errors": [], "file_hash": file_hash}
    
    # --- 处理文本 ---
    if text and len(text) > 20:
        text_chunks = smart_chunk_semantic(text)
        for i, tc in enumerate(text_chunks):
            chunks.append({
                "file_hash": file_hash,
                "file_name": file_name,
                "category": category,
                "chunk_index": i,
                "text": tc,
                "result_type": "text",
                "doc_type": doc_type,
                "_source": f"parser:{doc_type}",
            })
    
    # --- 处理表格 ---
    for t in tables:
        table_chunks = chunk_table(t)
        for tc in table_chunks:
            chunks.append({
                "file_hash": file_hash,
                "file_name": file_name,
                "category": category,
                "chunk_index": len(chunks),
                "text": tc,
                "result_type": "table",
                "doc_type": doc_type,
                "_source": "parser:table",
                "sheet_name": t.get('sheet', ''),
                "table_rows": t.get('rows', 0),
                "table_cols": t.get('cols', 0),
            })
    
    # --- 处理图片 ---
    for img_path in images:
        img_chunks = chunk_image(img_path, text)
        for ic in img_chunks:
            chunks.append({
                "file_hash": file_hash,
                "file_name": file_name,
                "category": category,
                "chunk_index": len(chunks),
                "text": ic,
                "result_type": "image",
                "doc_type": doc_type,
                "_source": "parser:image",
            })
    
    if not chunks:
        logger.info(f"[Ingest] No content extracted from {file_name}")
        return result
    
    # --- 写入向量库 ---
    try:
        if embed_fn and vector_store:
            texts = [c['text'][:1000] for c in chunks]
            embeddings = await embed_fn(texts)
            
            if embeddings:
                ids = [f"{file_hash}_{c['chunk_index']}" for c in chunks]
                metadatas = [{k: str(v)[:512] for k, v in c.items() if k != 'text'} for c in chunks]
                documents = [c['text'] for c in chunks]
                
                vector_store.add(ids=ids, embeddings=embeddings, documents=documents, metadata=metadatas)
                result["chunks_added"] = len(chunks)
                logger.info(f"[Ingest] Added {len(chunks)} chunks to vector store")
    except Exception as e:
        result["errors"].append(f"vector_store: {str(e)}")
        logger.error(f"[Ingest] vector_store error: {e}")
    
    # --- 写入 BM25 索引 ---
    try:
        if memory_store:
            for c in chunks:
                memory_store.add_document(c)
    except Exception as e:
        result["errors"].append(f"memory_store: {str(e)}")
    
    # --- 写入表格独立索引 ---
    try:
        if table_store and tables:
            from src.services.table_view import index_tables_from_chunks
            table_result = await index_tables_from_chunks(chunks, clear_first=False)
            result["tables_indexed"] = table_result.get("tables_indexed", 0)
            
            # 1.5.3b: 从正文中移除表格原文（已被独立索引）
            if table_result.get("tables_indexed", 0) > 0:
                import re
                table_pattern = re.compile(r'\|[^\n]+\|\n\|[-\s|:]+\|\n(?:\|[^\n]+\|\n)+', re.MULTILINE)
                for c in chunks:
                    text = c.get("text", "")
                    if text and "|" in text and "---" in text:
                        cleaned = table_pattern.sub("", text).strip()
                        if len(cleaned) > 50:  # 保留有意义的剩余内容
                            c["text"] = cleaned
    except Exception as e:
        result["errors"].append(f"table_store: {str(e)}")
    
    return result


async def ingest_directory(
    dir_path: str,
    category: str = "",
    embed_fn=None,
    vector_store=None,
    table_store=None,
    memory_store=None
) -> Dict:
    """批量入库整个目录"""
    from src.services.parsers import parse_file, identify_file
    
    total = {"chunks_added": 0, "tables_indexed": 0, "files_processed": 0, "errors": []}
    
    for root, dirs, files in os.walk(dir_path):
        for fname in files:
            fpath = os.path.join(root, fname)
            ftype = identify_file(fpath)
            
            if ftype == 'unknown':
                continue
            
            try:
                parsed = await parse_file(fpath)
                rel_name = os.path.relpath(fpath, dir_path)
                result = await ingest_document(
                    parsed, rel_name, category,
                    embed_fn, vector_store, table_store, memory_store
                )
                total["chunks_added"] += result.get("chunks_added", 0)
                total["tables_indexed"] += result.get("tables_indexed", 0)
                total["files_processed"] += 1
                if result.get("errors"):
                    total["errors"].extend(result["errors"])
            except Exception as e:
                total["errors"].append(f"{fpath}: {str(e)}")
    
    return total


def minhash_dedup(texts: list, threshold: float = 0.85) -> list:
    """MinHash 近似去重：返回去重后的索引列表
    
    原理：用 MinHash 签名估计 Jaccard 相似度，跳过相似度 > 阈值的文档
    """
    import hashlib
    
    def _shingles(text, k=5):
        """生成 k-gram 集合"""
        text = text.lower().replace("\n", " ").replace("  ", " ")
        return set(text[i:i+k] for i in range(max(0, len(text) - k + 1)))
    
    def _minhash(shingles, num_hashes=128):
        """计算 MinHash 签名"""
        sig = []
        for i in range(num_hashes):
            min_val = float('inf')
            for s in shingles:
                h = int(hashlib.md5(f"{i}_{s}".encode()).hexdigest(), 16)
                min_val = min(min_val, h)
            sig.append(min_val)
        return sig
    
    def _jaccard_est(sig1, sig2):
        """估计 Jaccard 相似度"""
        matches = sum(1 for a, b in zip(sig1, sig2) if a == b)
        return matches / len(sig1)
    
    n = len(texts)
    if n <= 1:
        return list(range(n))
    
    # 计算所有签名
    sigs = [_minhash(_shingles(t)) for t in texts]
    
    # 贪心去重
    keep = [True] * n
    for i in range(n):
        if not keep[i]:
            continue
        for j in range(i + 1, n):
            if not keep[j]:
                continue
            if _jaccard_est(sigs[i], sigs[j]) > threshold:
                keep[j] = False
    
    return [i for i in range(n) if keep[i]]

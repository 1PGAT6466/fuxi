"""
parser.py — 少阳·文档解析器
PDF/DOCX/XLSX/TXT/HTML/CSV/PPTX
"""
import logging
from typing import Dict, Optional
from pathlib import Path

logger = logging.getLogger("shaoyang.parser")


class DocumentParser:
    """文档解析器"""

    SUPPORTED_TYPES = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".doc": "docx",
        ".xlsx": "xlsx",
        ".xls": "xlsx",
        ".txt": "text",
        ".md": "text",
        ".html": "html",
        ".htm": "html",
        ".csv": "csv",
        ".pptx": "pptx",
    }

    def parse(self, file_path: str) -> Dict:
        """解析文档"""
        path = Path(file_path)
        suffix = path.suffix.lower()
        file_type = self.SUPPORTED_TYPES.get(suffix)

        if not file_type:
            return {"error": f"不支持的文件类型: {suffix}", "text": ""}

        try:
            if file_type == "text":
                return self._parse_text(path)
            elif file_type == "html":
                return self._parse_html(path)
            elif file_type == "csv":
                return self._parse_csv(path)
            elif file_type == "pdf":
                return self._parse_pdf(path)
            elif file_type == "docx":
                return self._parse_docx(path)
            elif file_type == "xlsx":
                return self._parse_xlsx(path)
            elif file_type == "pptx":
                return self._parse_pptx(path)
        except Exception as e:
            logger.error(f"[Parser] 解析失败: {e}")
            return {"error": str(e), "text": ""}

        return {"error": "未知错误", "text": ""}

    def _parse_text(self, path: Path) -> Dict:
        text = path.read_text(encoding="utf-8", errors="ignore")
        return {"text": text, "file_type": "text", "pages": 1}

    def _parse_html(self, path: Path) -> Dict:
        import re
        text = path.read_text(encoding="utf-8", errors="ignore")
        text = re.sub(r"<[^>]+>", "", text)
        text = re.sub(r"\s+", " ", text).strip()
        return {"text": text, "file_type": "html", "pages": 1}

    def _parse_csv(self, path: Path) -> Dict:
        import csv
        rows = []
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append("\t".join(row))
        return {"text": "\n".join(rows), "file_type": "csv", "pages": 1}

    def _parse_pdf(self, path: Path) -> Dict:
        try:
            import fitz
            doc = fitz.open(str(path))
            text = ""
            for page in doc:
                text += page.get_text()
            return {"text": text, "file_type": "pdf", "pages": len(doc)}
        except ImportError:
            return {"error": "需要安装PyMuPDF", "text": ""}

    def _parse_docx(self, path: Path) -> Dict:
        try:
            import docx
            doc = docx.Document(str(path))
            text = "\n".join([p.text for p in doc.paragraphs])
            return {"text": text, "file_type": "docx", "pages": 1}
        except ImportError:
            return {"error": "需要安装python-docx", "text": ""}

    def _parse_xlsx(self, path: Path) -> Dict:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True)
            text = ""
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    text += "\t".join([str(c) if c else "" for c in row]) + "\n"
            return {"text": text, "file_type": "xlsx", "pages": 1}
        except ImportError:
            return {"error": "需要安装openpyxl", "text": ""}

    def _parse_pptx(self, path: Path) -> Dict:
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return {"text": text, "file_type": "pptx", "pages": len(prs.slides)}
        except ImportError:
            return {"error": "需要安装python-pptx", "text": ""}

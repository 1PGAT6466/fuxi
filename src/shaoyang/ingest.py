"""
ingest.py - 鏂囨湰鎻愬彇銆佸垎绫汇 佸垎鍧椼 佹竻娲楁 鍧?
浠?server.py 鍒嗙 锛屼緵涓昏矾鐢卞 鍏 娇鐢?
"""

import html
import os
import re
from pathlib import Path

import logging

logger = logging.getLogger(__name__)

try:
    import jieba

    jieba.setLogLevel(20)
except ImportError:
    jieba = None

def _clean_text(raw: str) -> str:
    """
    鏁版嵁娓呮礂锛坴13.0 澧炲己锛夛細
    - 鍘婚櫎 HTML 瀹炰綋缂栫爜
    - 鍘婚櫎 URL 閾炬帴
    - 鍘婚櫎涓嶅彲瑙佹帶鍒跺瓧绗?
    - 鍘婚櫎甯歌 椤电湁椤佃剼锛堥〉鐮併 佽矾寰勩 侀偖浠剁 鍚嶏級
    - 鍘婚櫎杩炵画绌鸿 
    """
    raw = html.unescape(raw)
    raw = re.sub(r"https?://\S+", "", raw)
    raw = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", raw)

    # 椤电湁椤佃剼娓呯悊
    raw = re.sub(r"^\s*绗琝s*\d+\s*椤礬s*(鍏盶s*\d+\s*椤??\s*$", "", raw, flags=re.MULTILINE)
    raw = re.sub(r"^\s*Page\s*\d+\s*(of\s*\d+)?\s*$", "", raw, flags=re.MULTILINE | re.IGNORECASE)
    raw = re.sub(r"^\s*\d{4}-\d{2}-\d{2}\s*\d{1,2}:\d{2}(:\d{2})?\s*$", "", raw, flags=re.MULTILINE)  # 鐙 珛鏃堕棿鎴宠 
    raw = re.sub(r"^\s*[A-Z]:\\[^\n]{10,80}\s*$", "", raw, flags=re.MULTILINE)  # 鏂囦欢璺 緞琛?
    # 閭 欢绛惧悕
    raw = re.sub(r"--+\.?\s*$", "", raw, flags=re.MULTILINE)  # -- 鍒嗛殧绾?
    raw = re.sub(r"(Best Regards|Sincerely|姝 嚧|鏁  |椤虹 鍟嗙 )[\s\S]{0,200}$", "", raw, flags=re.IGNORECASE)
    # 涔辩爜瀛楃 锛堣繛缁?3+ 涓 潪涓 嫳鏁板瓧绗 級
    raw = re.sub(r"[^\w\s\u4e00-\u9fff\u3000-\u303f锛屻 傦紒锛燂紱锛氣 溾 濃 樷 欙紙锛夈 愩 戙 娿 嬨 佲   斅穃-\+\.\/]{3,}", " ", raw)

    raw = re.sub(r"\n{4,}", "\n\n\n", raw)
    return raw.strip()


def _generate_summary(text: str, max_len: int = 200) -> str:
    first_para = ""
    for para in text.split("\n\n"):
        para = para.strip()
        if len(para) > 30 and not para.startswith("#"):
            first_para = para[:max_len]
            break
    if not first_para:
        first_para = text[:max_len].strip()
    keywords = []
    try:
        if jieba:
            import jieba.analyse

            keywords = jieba.analyse.extract_tags(text, topK=5)
    except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
        logger.warning(f"[ingest] suppressed exception", exc_info=True)
        pass
    kw_str = "銆?".join(keywords[:5]) if keywords else ""
    summary = f"[鏂囨。鎽樿 ] {first_para}"
    if kw_str:
        summary += f"\n[鍏抽敭璇峕 {kw_str}"
    return summary


# v15.0: 浣跨敤 services/chunking.py 鐨?Markdown-AST 鏅鸿兘鍒嗗潡
def _smart_chunk(text: str, size: int = 1200, overlap: int = 100) -> list:
    """鏅鸿兘鍒嗗潡"""
    if not text or len(text) < 50:
        return [text] if text and len(text.strip()) > 10 else []

    chunks = []
    start = 0
    text_len = len(text)

    while start < text_len:
        end = min(start + size, text_len)

        # 灏濊瘯鍦 彞鍙枫 佹崲琛屽 鏂 紑
        if end < text_len:
            for sep in ["\n\n", "\n", "。", "，", ".", ";"]:
                last_sep = text.rfind(sep, start + size // 2, end)
                if last_sep > start:
                    end = last_sep + len(sep)
                    break

        chunk = text[start:end].strip()
        if chunk and len(chunk) > 20:
            # v4.0 琛 牸缁撴瀯鍖栨彁鍙?
            structured = None
            try:
                from src.infra import extract_tables_from_markdown

                if "|" in chunk:
                    tables = extract_tables_from_markdown(chunk)
                    if tables:
                        structured = tables[0]
            except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
                logger.debug("[suppressed] structured = tables[0]")
                pass
            if structured:
                chunks.append({"text": chunk, "structured_table": structured})
            else:
                chunks.append(chunk)

        start = end - overlap if end < text_len else text_len

    return chunks


def _extract_pdf_dual(file_path: str) -> str:
    """
    v10.0: PDF 鍙岃建瑙ｆ瀽
    - pdfplumber 涓诲姏锛堜繚鐣欒鏍?鍙屾爮/椤电爜淇 伅锛?
    - PyPDF2 鍥為  锛堝吋瀹规棫 PDF锛?
    """
    lines = []

    # 璺?: fitz (PyMuPDF) - 涓 枃鏈 浼橈紝浠呭  < 50MB 鏂囦欢浣跨敤
    file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
    if file_size_mb < 50:
        try:
            import fitz

            doc = fitz.open(file_path)
            total = doc.page_count
            for i in range(total):
                try:
                    page = doc[i]
                    text = page.get_text()
                    if text and text.strip():
                        lines.append("[Page %d/%d]\n%s" % (i + 1, total, text.strip()))
                except (
                    OSError,
                    ValueError,
                    KeyError,
                    ConnectionError,
                    TimeoutError,
                ) as e:  # TODO: Narrow exception type
                    logger.warning(f"[ingest] suppressed exception", exc_info=True)
                    pass
            doc.close()
            if lines:
                return "\n".join(lines)
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
            logger.warning(f"[ingest] suppressed exception", exc_info=True)
            pass

    # 璺?: pdfplumber
    try:
        import pdfplumber

        with pdfplumber.open(file_path) as pdf:
            total = len(pdf.pages)
            for i, page in enumerate(pdf.pages):
                try:
                    # 鎻愬彇鏂囨湰锛堜繚鐣欏竷灞 锛?
                    text = page.extract_text() or ""
                    # 鎻愬彇琛 牸锛堢粨鏋勫寲鏍囨敞锛?
                    from src.infra import enhance_table_extraction

                    combined = enhance_table_extraction(text, tables)
                    if combined.strip():
                        lines.append(f"[Page {i+1}/{total}]\n{combined.strip()}")
                except (
                    OSError,
                    ValueError,
                    KeyError,
                    ConnectionError,
                    TimeoutError,
                ) as e:  # TODO: Narrow exception type
                    # 鍗曢〉澶辫触锛屽洖閫 
                    try:
                        from pypdf import PdfReader

                        reader = PdfReader(file_path)
                        if i < len(reader.pages):
                            txt = reader.pages[i].extract_text() or ""
                            if txt.strip():
                                lines.append(f"[Page {i+1}/{total}]\n{txt}")
                    except (
                        OSError,
                        ValueError,
                        KeyError,
                        ConnectionError,
                        TimeoutError,
                    ) as e:  # TODO: Narrow exception type
                        lines.append(f"[Page {i+1}/{total}] skipped")
        if lines:
            return "\n".join(lines)
    except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
        logger.warning(f"[ingest] suppressed exception", exc_info=True)
        pass

    # 璺?: pypdf 鍥為  
    try:
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        total = len(reader.pages)
        for i, page in enumerate(reader.pages):
            try:
                txt = page.extract_text() or ""
                if txt.strip():
                    lines.append(f"[Page {i+1}/{total}]\n{txt}")
            except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
                lines.append(f"[Page {i+1}/{total}] skipped")
        return "\n".join(lines)
    except Exception as e:  # TODO: Narrow exception type
        return f"[PDF 瑙ｆ瀽澶辫触: {e}]"

    return ""


def _compute_file_hash(file_path: str) -> str:
    """
    v10.0: 璁畻鏂囦欢 SHA256锛圡D5 鍘婚噸鐢 級
    """
    import hashlib

    sha = hashlib.sha256()
    with open(file_path, "rb") as f:
        while chunk := f.read(8192):
            sha.update(chunk)
    return sha.hexdigest()


# ============================================================
# v16.0: _extract_text 绛栫暐妯″紡閲嶆瀯 鈥?姣忎釜鏂囦欢绫诲瀷鐙 珛澶勭悊鍑芥暟
# ============================================================

# 绾 枃鏈 墿灞曞悕闆嗗悎锛圲TF-8 鐩存帴璇诲彇锛?
_PLAINTEXT_EXTS = frozenset(
    [
        ".txt",
        ".md",
        ".csv",
        ".cfg",
        ".log",
        ".ini",
        ".conf",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".py",
        ".js",
        ".ts",
        ".java",
        ".c",
        ".cpp",
        ".h",
        ".sh",
        ".bat",
        ".ps1",
        ".yaml",
        ".yml",
    ]
)

# 浜岃繘鍒?涓嶅彲璇绘牸寮忛泦鍚堬紙浠呰 褰曟枃浠跺悕锛?
_BINARY_EXTS = frozenset(
    [
        ".deb",
        ".dwg",
        ".dxf",
        ".stp",
        ".step",
        ".igs",
        ".iges",
        ".jpg",
        ".jpeg",
        ".png",
        ".gif",
        ".bmp",
        ".svg",
        ".exe",
        ".msi",
        ".apk",
        ".dmg",
        ".pkg",
        ".rpm",
        ".7z",
        ".rar",
        ".tar.gz",
        ".gz",
        ".bin",
        ".iso",
        ".img",
        ".dll",
        ".so",
        ".o",
        ".a",
        ".lib",
        ".mp3",
        ".mp4",
        ".avi",
        ".mkv",
        ".wav",
        ".flac",
        ".mov",
        ".wmv",
    ]
)

# ZIP 鍐呮敮鎸侀 掑綊瑙ｆ瀽鐨勬枃浠舵墿灞曞悕
_ZIP_SUPPORTED_EXTS = frozenset(
    [
        ".txt",
        ".md",
        ".csv",
        ".log",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".cfg",
        ".ini",
        ".conf",
        ".docx",
        ".doc",
        ".xlsx",
        ".xls",
        ".pdf",
        ".pptx",
        ".ppt",
    ]
)

# 鑷 剤灞傦細鎵 睍鍚?鈫?pip 瀹夎 鍊欓 夛紙鎸変紭鍏堢骇鎺掑垪锛?
_SELF_HEAL_MAP = {
    "msg": ["extract-msg", "msg-parser"],
    "eml": ["mail-parser"],
    "wps": ["python-docx"],
    "rtf": ["striprtf", "pyth"],
    "ods": ["odfpy", "ezodf"],
    "odt": ["odfpy", "python-docx"],
    "epub": ["ebooklib", "epub2txt"],
    "djvu": ["djvulibre-python"],
    "ps": ["ghostscript"],
    "ai": ["pdf2image"],
    "cdr": ["pdf2image"],
    "one": ["python-docx"],
    "vsd": ["python-docx"],
    "vsdx": ["python-docx"],
    "mpp": ["python-docx"],
    "pub": ["python-docx"],
}


def _extract_plaintext(path: Path) -> str:
    """鎻愬彇绾 枃鏈 枃浠跺唴瀹癸紙UTF-8 缂栫爜锛?"""""
    return path.read_text(encoding="utf-8", errors="ignore")


def _extract_docx(path: Path) -> str:
    """鎻愬彇 .docx 鏂囦欢鍐呭"""
    from docx import Document

    return "\n".join(p.text for p in Document(str(path)).paragraphs if p.text.strip())


def _extract_doc(path: Path) -> str:
    """鎻愬彇 .doc 鏂囦欢鍐呭锛坅ntiword 鈫?python-docx 鍥為锛?"""""
    import subprocess

    # 璺?: antiword 鍛戒护琛屽伐鍏?
    try:
        result = subprocess.run(["antiword", "-m", "UTF-8.txt", str(path)], capture_output=True, text=True, timeout=30)
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
    except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
        logger.warning(f"[ingest] suppressed exception", exc_info=True)
    # 璺?: 鍥為  鍒?python-docx
    try:
        return _extract_docx(path)
    except Exception as e:  # TODO: Narrow exception type
        logger.warning("Exception 澶辫触: %s", e, exc_info=True)
        return ""


def _extract_wps(path: Path) -> str:
    """鎻愬彇 .wps 鏂囦欢鍐呭 锛堟柊鐗?WPS 鍙 兘鏄?docx 鏍煎紡锛?"""""
    try:
        return _extract_docx(path)
    except Exception as e:  # TODO: Narrow exception type
        logger.warning("Exception 澶辫触: %s", e, exc_info=True)
        return ""


def _extract_xlsx(path: Path) -> str:
    """鎻愬彇 .xlsx 鏂囦欢鍐呭涓虹粨鏋勫寲鏂囨湰"""
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    lines = []
    for sn in wb.sheetnames:
        ws = wb[sn]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        lines.append(f"[Sheet: {sn}] 琛屾暟={len(rows)}")
        headers = [str(c) if c is not None else "" for c in rows[0]]
        lines.append(f"[鍒楀悕] {' | '.join(headers)}")
        for row in rows[1:]:
            vals = []
            for j, c in enumerate(row):
                if c is None:
                    continue
                col_name = headers[j] if j < len(headers) else f"Col{j}"
                vals.append(f"{col_name}={c}")
            if vals:
                lines.append(" | ".join(vals))
    wb.close()
    return "\n".join(lines)


def _extract_xls(path: Path) -> str:
    """鎻愬彇 .xls 鏂囦欢鍐呭涓虹粨鏋勫寲鏂囨湰"""
    import xlrd

    wb = xlrd.open_workbook(str(path))
    lines = []
    for sn in wb.sheet_names():
        ws = wb.sheet_by_name(sn)
        if ws.nrows == 0:
            continue
        lines.append(f"[Sheet: {sn}] 琛屾暟={ws.nrows}")
        headers = [str(ws.cell_value(0, j)) for j in range(ws.ncols)]
        lines.append(f"[鍒楀悕] {' | '.join(headers)}")
        for r in range(1, ws.nrows):
            vals = []
            for j in range(ws.ncols):
                v = ws.cell_value(r, j)
                if v == "" or v is None:
                    continue
                col_name = headers[j] if j < len(headers) else f"Col{j}"
                vals.append(f"{col_name}={v}")
            if vals:
                lines.append(" | ".join(vals))
    return "\n".join(lines)


def _extract_pptx_modern(path: Path) -> str:
    """鎻愬彇 .pptx 鏂囦欢鍐呭锛坧ython-pptx锛?"""""
    from pptx import Presentation

    slides = []
    for i, sl in enumerate(Presentation(str(path)).slides):
        t = [f"[Slide {i+1}]"]
        for sh in sl.shapes:
            if sh.has_text_frame:
                t.append(sh.text_frame.text)
        if sl.has_notes_slide and sl.notes_slide.notes_text_frame:
            notes = sl.notes_slide.notes_text_frame.text.strip()
            if notes:
                t.append("[Notes] " + notes)
        slides.append("\n".join(t))
    result = "\n\n".join(slides)
    if len(result.strip()) > 50:
        return result
    return ""


def _extract_ppt_legacy(path: Path) -> str:
    """鎻愬彇鏃 増 .ppt 鏂囦欢鍐呭 锛坥lefile锛?"""""
    import olefile

    ole = olefile.OleFileIO(str(path))
    texts = []
    for stream in ole.listdir():
        try:
            raw = ole.openstream(stream).read()
            try:
                t = raw.decode("utf-16-le", errors="ignore")
            except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
                continue
            t = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", t)
            t = t.strip()
            if len(t) > 20:
                texts.append(t)
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
            pass
    ole.close()
    result = "\n\n".join(texts)
    if len(result.strip()) > 100:
        return result
    return ""


def _extract_pptx(path: Path) -> str:
    """鎻愬彇 .pptx/.ppt 鏂囦欢鍐呭锛坧ptx 浼樺厛锛宭egacy ole 鍥為锛?"""""
    try:
        result = _extract_pptx_modern(path)
        if result:
            return result
    except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
        logger.warning(f"[ingest] suppressed exception", exc_info=True)
    try:
        return _extract_ppt_legacy(path)
    except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
        logger.warning(f"[ingest] suppressed exception", exc_info=True)
    return ""


def _extract_zip(path: Path) -> str:
    """鎻愬彇 .zip 鍘嬬缉鍖呭唴鍙  鏂囦欢鐨勫唴瀹癸紙閫掑綊瑙ｆ瀽锛屾渶澶?20 涓 唴閮 枃浠讹級"""
    import shutil
    import tempfile
    import zipfile

    lines = []
    try:
        with zipfile.ZipFile(str(path), "r") as zf:
            for name in zf.namelist():
                if name.endswith("/"):
                    continue
                inner_ext = os.path.splitext(name)[1].lower()
                if inner_ext not in _ZIP_SUPPORTED_EXTS:
                    continue
                try:
                    with zf.open(name) as inner_f:
                        content = inner_f.read()
                    tmp_dir = tempfile.mkdtemp()
                    tmp_path = os.path.join(tmp_dir, os.path.basename(name))
                    with open(tmp_path, "wb") as tmp_f:
                        tmp_f.write(content)
                    inner_text = _extract_text(tmp_path, inner_ext)
                    shutil.rmtree(tmp_dir, ignore_errors=True)
                    if inner_text.strip():
                        lines.append(f"[ZIP鍐? {name}]\n{inner_text[:2000]}")
                except (
                    OSError,
                    ValueError,
                    KeyError,
                    ConnectionError,
                    TimeoutError,
                ) as e:  # TODO: Narrow exception type
                    lines.append(f"[ZIP鍐? {name}] 鎻愬彇澶辫触")
        return "\n\n---\n\n".join(lines[:20])
    except Exception as e:  # TODO: Narrow exception type
        return f"[ZIP鎻愬彇澶辫触: {e}]"


def _extract_binary(path: Path) -> str:
    """杩斿洖浜岃繘鍒?涓嶅彲璇绘牸寮忕殑鍗犱綅淇 伅"""
    return f"[鏂囦欢: {path.name}] (浜岃繘鍒舵牸寮忥紝鏃犳枃鏈 彁鍙?"


def _extract_generic_fallback(path: Path) -> str:
    """閫氱敤闄嶇骇瑙ｆ瀽鍣 細灏濊瘯 raw bytes 鈫?UTF-8 鈫?UTF-16 鈫?latin-1"""
    try:
        with open(str(path), "rb") as f:
            raw = f.read()
    except Exception as e:  # TODO: Narrow exception type
        logger.warning("Exception 澶辫触: %s", e, exc_info=True)
        return ""

    # 灏濊瘯 UTF-8
    try:
        text = raw.decode("utf-8")
        if len(text.strip()) > 50:
            return text
    except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
        pass

    # 灏濊瘯 UTF-16 LE
    try:
        text = raw.decode("utf-16-le")
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", text)
        if len(text.strip()) > 100:
            return text
    except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
        pass

    # 灏濊瘯 latin-1 + 鍙  琛屾彁鍙?
    try:
        text = raw.decode("latin-1")
        lines_list = text.split("\n")
        readable = []
        for line in lines_list:
            alpha_ratio = sum(1 for c in line if c.isalpha() or c.isspace()) / max(len(line), 1)
            if alpha_ratio > 0.6 and len(line.strip()) > 10:
                readable.append(line.strip())
        if len(readable) > 3:
            return "\n".join(readable)
    except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
        pass

    return ""


def _extract_self_heal(path: Path, ext: str) -> str:
    """鏅鸿兘鑷 剤灞傦細妫 娴嬬己澶辩殑瑙ｆ瀽渚濊禆锛屾彁绀虹敤鎴峰畨瑁呭苟閲嶈瘯"""
    candidates = _SELF_HEAL_MAP.get(ext, [])
    if not candidates:
        return ""
    import importlib as _il

    installed = None
    for pkg in candidates:
        try:
            _il.import_module(pkg.replace("-", "_"))
            installed = pkg
            break
        except ImportError:
            continue
    if not installed:
        logger.warning(f"[ingest] 缂哄皯瑙ｆ瀽渚濊禆 {candidates[0]} (ext={ext})锛岃 鎵嬪姩瀹夎 : pip install {candidates[0]}")
        return ""
    # 渚濊禆宸插畨瑁咃紝閫掑綊璋冪敤 _extract_text 閲嶆柊瑙ｆ瀽
    return _extract_text(str(path), ext)


# 鎵 睍鍚?鈫?鎻愬彇鍑芥暟鐨勮皟搴 瓧鍏革紙绛栫暐妯″紡锛?
_EXTRACTOR_DISPATCH = {
    "docx": _extract_docx,
    "doc": _extract_doc,
    "wps": _extract_wps,
    "xlsx": _extract_xlsx,
    "xls": _extract_xls,
    "pdf": lambda p: _extract_pdf_dual(str(p)),
    "pptx": _extract_pptx,
    "ppt": _extract_pptx,
    "zip": _extract_zip,
}


def _extract_text(file_path: str, ext: str) -> str:
    """
    缁熶竴鏂囨鏂囨湰鎻愬彇鍏 彛锛坴16.0 绛栫暐妯紡閲嶆瀯锛夈 ?

    鏍规嵁鏂囦欢鎵 睍鍚嶅垎娲惧埌瀵瑰簲鐨勬彁鍙栧嚱鏁帮細
    - 绾 枃鏈 墿灞曞悕 鈫?UTF-8 鐩存帴璇诲彇
    - 鍔炲叕鏂囨锛坉ocx/doc/wps/xlsx/xls/pdf/pptx/ppt/zip锛夆啋 涓撶敤鎻愬彇鍣?
    - 浜岃繘鍒舵牸寮忥紙鍥剧墖/闊抽 /鍙 墽琛屾枃浠剁瓑锛夆啋 鍗犱綅淇 伅
    - 鍏朵粬 鈫?閫氱敤闄嶇骇瑙ｆ瀽鍣?鈫?鑷 剤灞?

    Args:
        file_path: 鏂囦欢璺 緞瀛楃 涓?
        ext: 鏂囦欢鎵 睍鍚嶏紙鍚 偣鍙凤級锛屽  ".pdf"銆?.docx""

    Returns:
        鎻愬彇鐨勬枃鏈 唴瀹瑰瓧绗 覆銆傚 璐 椂杩斿洖閿欒 鎻忚堪瀛楃 涓层 ?
    """
    path = Path(file_path)

    # Phase 1: 宸茬煡鏍煎紡 鈥?鐩存帴璋冨害
    try:
        # 绾 枃鏈?
        if ext in _PLAINTEXT_EXTS:
            return _extract_plaintext(path)

        # 涓撶敤鎻愬彇鍣 紙绛栫暐瀛楀吀鏌 壘锛?
        ext_key = ext.lstrip(".").lower()
        extractor = _EXTRACTOR_DISPATCH.get(ext_key)
        if extractor is not None:
            return extractor(path)

        # 浜岃繘鍒?涓嶅彲璇绘牸寮?
        if ext in _BINARY_EXTS:
            return _extract_binary(path)

        # 鏈 煡鏍煎紡锛氬皾璇?UTF-8 鏂囨湰璇诲彇
        try:
            return path.read_text(encoding="utf-8", errors="replace")
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
            return f"[鏂囦欢: {path.name}] (鏈 煡鏍煎紡)"

    except Exception as e:  # TODO: Narrow exception type
        return f"[鎻愬彇澶辫触: {e}]"

    # Phase 2: 閫氱敤闄嶇骇锛堜粎褰?Phase 1 鏈?return 鏃跺埌杈撅級
    result = _extract_generic_fallback(path)
    if result:
        return result

    # Phase 3: 鏅鸿兘鑷 剤
    ext_key = ext.lstrip(".").lower() if "." in ext else ext
    healed = _extract_self_heal(path, ext_key)
    if healed:
        return healed

    return ""


# === merged from ingestion.py ===
"""
浼忕静 Fuxi 路 缁熶竴鍏 簱寮曟搸
========================
灏嗚 鏋愬櫒杈撳嚭鐨勬爣鍑嗗寲鏂囨 鈫?鎸夌被鍨嬪垎鍧?鈫?鍚戦噺鍖?鈫?鍐欏叆 ChromaDB

娴佺 :
  parse_result 鈫?type_router 鈫?chunker 鈫?embedder 鈫?ChromaDB
             鈫?table_extractor 鈫?kb_tables 鐙 珛绱 紩
"""
import hashlib
import logging
import os
import re
from pathlib import Path
from typing import Dict, List

logger = logging.getLogger(__name__)

# ============================================================
# 璇 箟鍒嗗潡
# ============================================================


def smart_chunk_semantic(text: str, chunk_size: int = 1500, overlap: int = 200) -> List[str]:
    """璇 箟鍒嗗潡锛氭寜娈佃惤杈圭晫鍒囧壊锛屼繚鎸佽 涔夊畬鏁存 ?"""""
    if not text:
        return []

    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks = []
    current = ""

    for p in paragraphs:
        # 濡傛灉鏄 爣棰樿 锛堝寘鍚?# 鎴栧叏閮  鍐欑瓑锛夛紝鏂拌捣涓 涓?chunk
        is_heading = p.startswith("#") or p.startswith("(") or p.startswith("[")

        if is_heading and current and len(current) > 100:
            chunks.append(current.strip())
            current = p
        elif len(current) + len(p) > chunk_size:
            chunks.append(current.strip())
            # 閲嶅彔锛氫繚鐣欎笂涓 娈佃惤鐨勬渶鍚庡唴瀹?
            if overlap and current:
                overlap_text = current[-overlap:] if len(current) > overlap else current[-50:]
                current = overlap_text + "\n\n" + p
            else:
                current = p
        else:
            if current:
                current += "\n\n" + p
            else:
                current = p

    if current.strip():
        chunks.append(current.strip())

    return chunks


def chunk_table(table_data: Dict) -> List[str]:
    """琛 牸涓嶅垎鍧楋紝鏁磋淇濈暀涓轰竴涓?chunk"""
    if not table_data:
        return []

    chunks = []
    if table_data.get("markdown"):
        chunks.append(table_data["markdown"])
    else:
        # 鐢熸垚 Markdown 琛  
        header = table_data.get("header", [])
        rows = table_data.get("all_rows", table_data.get("sample_rows", []))
        if header:
            lines = ["| " + " | ".join(str(h) for h in header) + " |"]
            lines.append("|" + "|".join(["---" for _ in header]) + "|")
            for row in rows:
                lines.append("| " + " | ".join(str(c) for c in row) + " |")
            chunks.append("\n".join(lines))

    return chunks


def chunk_image(file_path: str, ocr_text: str = "") -> List[str]:
    """Image chunking with multimodal + OCR"""
    # if ocr_text exists and is long enough, use directly
    if ocr_text and len(ocr_text) > 50:
        return [ocr_text]

    # try multimodal Vision model
    try:
        from src.infra import transcribe_image

        transcription = transcribe_image(file_path)
        if transcription and len(transcription) > 20:
            logger.info(f"[multimodal] 鍥剧墖杞 綍鎴愬姛: {Path(file_path).name} ({len(transcription)}瀛?")
            return [f"[鍥剧墖鍐呭 ] {transcription}"]
    except Exception as e:  # TODO: Narrow exception type
        logger.debug(f"[multimodal] 鍥剧墖杞 綍璺宠繃: {e}")

    # Fallback: OCR 鏂囨湰
    if ocr_text and len(ocr_text) > 10:
        return [ocr_text]

    # 鏈 缁?Fallback: 鏂囦欢鍚嶅崰浣?
    return [f"[鍥剧墖: {Path(file_path).name}]"]


# ============================================================
# 鍏 簱寮曟搸
# ============================================================
# FAKE-ASYNC: 鏈 嚱鏁版爣璁?async 浠呬负鎺 彛缁熶竴锛屽唴閮 悓姝 墽琛?

# ============================================================
# v3.0: ingest_document 杈呭姪鍑芥暟 鈥?鍑嗗  / 瀛樺偍 / 绱 紩涓夐樁娈?
# ============================================================


def _resolve_category(file_name: str, text: str, category: str) -> str:
    """_resolve_category"""

    if not category or category == "閫氱敤鍔炲叕":
        try:
            from src.category_registry import match_category as _match_cat

            ext = os.path.splitext(file_name)[1].lower() if file_name else ""
            _cat = _match_cat(text[:5000], file_ext=ext, file_name=file_name)
            if _cat:
                category = _cat
        except (OSError, ValueError, KeyError, ConnectionError, TimeoutError) as e:  # TODO: Narrow exception type
            logger.debug("[suppressed] category = _cat")
    # 闃插尽鎬 牎楠岋細category 涓嶈兘鏄?Python repr() 鏍煎紡
    if category and ("[{" in category or "': " in category):
        m = re.search(r"'category':\s*'([^']+)'", category)
        category = m.group(1) if m else "閫氱敤鍔炲叕"
    return category


def _prepare_chunks(
    text: str, tables: list, images: list, file_hash: str, file_name: str, category: str, doc_type: str
) -> list:
    """Phase 1: 灏嗚 鏋愮粨鏋滃垏鍒嗕负缁熶竴鐨?chunk 鍒楄锛堟枃鏈?琛 牸/鍥剧墖锛?"""""
    chunks = []

    # 鏂囨湰鍒嗗潡
    if text and len(text) > 20:
        text_chunks = smart_chunk_semantic(text)
        for i, tc in enumerate(text_chunks):
            chunks.append(
                {
                    "file_hash": file_hash,
                    "file_name": file_name,
                    "category": category,
                    "chunk_index": i,
                    "text": tc,
                    "result_type": "text",
                    "doc_type": doc_type,
                    "_source": f"parser:{doc_type}",
                }
            )

    # 琛 牸鍒嗗潡
    for t in tables:
        table_chunks = chunk_table(t)
        for tc in table_chunks:
            chunks.append(
                {
                    "file_hash": file_hash,
                    "file_name": file_name,
                    "category": category,
                    "chunk_index": len(chunks),
                    "text": tc,
                    "result_type": "table",
                    "doc_type": doc_type,
                    "_source": "parser:table",
                    "sheet_name": t.get("sheet", ""),
                    "table_rows": t.get("rows", 0),
                    "table_cols": t.get("cols", 0),
                }
            )

    # 鍥剧墖鍒嗗潡
    for img_path in images:
        img_chunks = chunk_image(img_path, text)
        for ic in img_chunks:
            chunks.append(
                {
                    "file_hash": file_hash,
                    "file_name": file_name,
                    "category": category,
                    "chunk_index": len(chunks),
                    "text": ic,
                    "result_type": "image",
                    "doc_type": doc_type,
                    "_source": "parser:image",
                }
            )

    return chunks


async def _store_to_vector(chunks: list, file_hash: str, embed_fn, vector_store) -> int:
    """Phase 2: 鍐欏叆鍚戦噺搴?(ChromaDB)锛岃繑鍥炴垚鍔熷啓鍏 殑 chunk 鏁?"""""
    if embed_fn and vector_store:
        texts = [c["text"][:1000] for c in chunks]
        embeddings = await embed_fn(texts)
        if embeddings:
            ids = [f"{file_hash}_{c['chunk_index']}" for c in chunks]
            metadatas = [{k: str(v)[:512] for k, v in c.items() if k != "text"} for c in chunks]
            documents = [c["text"] for c in chunks]
            vector_store.add(ids=ids, embeddings=embeddings, documents=documents, metadata=metadatas)
            logger.info(f"[Ingest] Added {len(chunks)} chunks to vector store")
            return len(chunks)
    return 0


def _store_to_memory(chunks: list, memory_store) -> None:
    """Phase 2: 鍐欏叆 BM25 鍏 枃绱 紩"""
    if memory_store:
        for c in chunks:
            memory_store.add_document(c)


async def _index_tables(chunks: list, tables: list, table_store) -> int:
    """Phase 2: 鍐欏叆琛 牸鐙 珛绱 紩锛屽苟娓呯悊姝ｆ枃涓 殑琛 牸鍘熸枃锛岃繑鍥炵储寮曞叆鐨勮鏍兼暟"""
    if not (table_store and tables):
        return 0
    from src.infra import index_tables_from_chunks

    table_result = await index_tables_from_chunks(chunks, clear_first=False)
    indexed = table_result.get("tables_indexed", 0)

    # 1.5.3b: 浠庢 鏂囦腑绉婚櫎琛 牸鍘熸枃锛堝凡琚 嫭绔嬬储寮曪級
    if indexed > 0:
        table_pattern = re.compile(r"\|[^\n]+\|\n\|[\-\s|:]+\|\n(?:\|[^\n]+\|\n)+", re.MULTILINE)
        for c in chunks:
            text_content = c.get("text", "")
            if text_content and "|" in text_content and "---" in text_content:
                cleaned = table_pattern.sub("", text_content).strip()
                if len(cleaned) > 50:
                    c["text"] = cleaned
    return indexed


async def ingest_document(
    parse_result: Dict,
    file_name: str = "",
    category: str = "",
    embed_fn=None,
    vector_store=None,
    table_store=None,
    memory_store=None,
) -> Dict:
    """
    缁熶竴鍏 簱涓 涓 枃妗?鈥?v3.0 涓夐樁娈甸噸鏋勩 ?

    闃舵 :
      1. 鍑嗗  (prepare): 鍒嗙被瑙ｆ瀽 鈫?鏂囨湰/琛 牸/鍥剧墖缁熶竴鍒嗗潡
      2. 瀛樺偍 (store): 鍐欏叆鍚戦噺搴?+ BM25 绱 紩 + 琛 牸鐙 珛绱 紩
      3. 姹囨 ?(result): 杩斿洖缁熻 淇 伅

    鍙傛暟:
      parse_result: 瑙ｆ瀽鍣 緭鍑?{type, text, metadata, tables, images}
      file_name: 鍘熷 鏂囦欢鍚?
      category: 鏂囨鍒嗙被
      embed_fn: 鍚戦噺鍖栧嚱鏁?
      vector_store: ChromaDB kb_chunks 闆嗗悎
      table_store: ChromaDB kb_tables 闆嗗悎
      memory_store: BM25 鍏 枃绱 紩

    杩斿洖: {chunks_added, tables_indexed, errors, file_hash}
    """
    doc_type = parse_result.get("type", "unknown")
    text = parse_result.get("text", "")
    metadata = parse_result.get("metadata", {})
    tables = parse_result.get("tables", [])
    images = parse_result.get("images", [])

    file_hash = hashlib.sha256((file_name + str(metadata)).encode()).hexdigest()[:16]
    category = _resolve_category(file_name, text, category)

    result = {"chunks_added": 0, "tables_indexed": 0, "errors": [], "file_hash": file_hash}

    # Phase 1: 鍑嗗  chunk 鍒楄〃
    chunks = _prepare_chunks(text, tables, images, file_hash, file_name, category, doc_type)
    if not chunks:
        logger.info(f"[Ingest] No content extracted from {file_name}")
        return result

    # Phase 2: 鍐欏叆鍚戦噺搴?
    try:
        result["chunks_added"] = await _store_to_vector(chunks, file_hash, embed_fn, vector_store)
    except Exception as e:  # TODO: Narrow exception type
        result["errors"].append(f"vector_store: {str(e)}")

    # Phase 2: 鍐欏叆 BM25 绱 紩
    try:
        _store_to_memory(chunks, memory_store)
    except Exception as e:  # TODO: Narrow exception type
        result["errors"].append(f"memory_store: {str(e)}")

    # Phase 2: 鍐欏叆琛 牸鐙 珛绱 紩
    try:
        result["tables_indexed"] = await _index_tables(chunks, tables, table_store)
    except Exception as e:  # TODO: Narrow exception type
        result["errors"].append(f"table_store: {str(e)}")

    return result


async def ingest_directory(
    dir_path: str, category: str = "", embed_fn=None, vector_store=None, table_store=None, memory_store=None
) -> Dict:
    """鎵归噺鍏 簱鏁翠釜鐩 綍"""
    from src.infra import identify_file, parse_file

    total = {"chunks_added": 0, "tables_indexed": 0, "files_processed": 0, "errors": []}

    for root, dirs, files in os.walk(dir_path):
        for fname in files:
            fpath = os.path.join(root, fname)
            ftype = identify_file(fpath)

            if ftype == "unknown":
                continue

            try:
                parsed = await parse_file(fpath)
                rel_name = os.path.relpath(fpath, dir_path)
                result = await ingest_document(
                    parsed, rel_name, category, embed_fn, vector_store, table_store, memory_store
                )
                total["chunks_added"] += result.get("chunks_added", 0)
                total["tables_indexed"] += result.get("tables_indexed", 0)
                total["files_processed"] += 1
                if result.get("errors"):
                    total["errors"].extend(result["errors"])
            except Exception as e:  # TODO: Narrow exception type
                total["errors"].append(f"{fpath}: {str(e)}")

    return total


def minhash_dedup(texts: list, threshold: float = 0.85) -> list:
    """MinHash approximate dedup.

    Estimates Jaccard similarity via MinHash signatures,
    skipping documents above similarity threshold.
    """
    import hashlib

    def _shingles(text, k=5):
        """鐢熸垚 k-gram 闆嗗悎"""
        text = text.lower().replace("\n", " ").replace("  ", " ")
        return set(text[i : i + k] for i in range(max(0, len(text) - k + 1)))

    def _minhash(shingles, num_hashes=128):
        """璁畻 MinHash 绛惧悕"""
        sig = []
        for i in range(num_hashes):
            min_val = float("inf")
            for s in shingles:
                h = int(hashlib.sha256(f"{i}_{s}".encode()).hexdigest(), 16)
                min_val = min(min_val, h)
            sig.append(min_val)
        return sig

    def _jaccard_est(sig1, sig2):
        """Estimate Jaccard similarity"""
        matches = sum(1 for a, b in zip(sig1, sig2) if a == b)
        return matches / len(sig1)

    n = len(texts)
    if n <= 1:
        return list(range(n))

    # 璁＄畻鎵 鏈夌 鍚?
    sigs = [_minhash(_shingles(t)) for t in texts]

    # 璐 績鍘婚噸
    keep = [True] * n
    for i in range(n):
        if not keep[i]:
            continue
        for j in range(i + 1, n):
            if not keep[j]:
                continue
            if _jaccard_est(sigs[i], sigs[j]) > threshold:
                keep[j] = False

    return [i for i in range(n) if keep[i]]

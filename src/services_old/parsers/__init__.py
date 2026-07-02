"""
伏羲 Fuxi · 多类型解析器
========================
按照文件类型分路处理的标准化解析框架

文字类: TXT/DOCX/MD/PPTX → 纯文本 + 结构元数据
PDF类: 原生PDF → pdfplumber文本 + 表格; 扫描PDF → OCR
表格类: CSV/XLSX → Pandas 结构化 → JSON/Markdown
图片类: PNG/JPG → OCR 文字提取

统一输出: {type, text, metadata, tables, images}
"""
import os
import re
import io
import csv
import json
import logging
from typing import List, Dict, Optional, Tuple
from pathlib import Path

logger = logging.getLogger(__name__)

# ============================================================
# 文件类型识别
# ============================================================

EXT_MAP = {
    # 文本类
    '.txt': 'text', '.md': 'text', '.markdown': 'text',
    '.docx': 'text', '.doc': 'text',
    '.pptx': 'text', '.ppt': 'text',
    # PDF 类
    '.pdf': 'pdf',
    # 表格类
    '.csv': 'table', '.xlsx': 'table', '.xls': 'table',
    # 图片类
    '.png': 'image', '.jpg': 'image', '.jpeg': 'image',
    '.gif': 'image', '.bmp': 'image', '.webp': 'image',
    # 压缩包
    '.zip': 'archive', '.rar': 'archive', '.7z': 'archive',
}

def identify_file(path: str) -> str:
    """根据后缀识别文件类型"""
    ext = Path(path).suffix.lower()
    return EXT_MAP.get(ext, 'unknown')


# ============================================================
# 文本解析器
# ============================================================

async def parse_text(file_path: str) -> Dict:
    """解析 TXT/MD/DOCX/PPTX → 纯文本 + 结构元数据"""
    ext = Path(file_path).suffix.lower()
    text = ""
    metadata = {"source": file_path, "type": "text", "format": ext}
    
    try:
        if ext in ('.txt', '.md', '.markdown'):
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                text = f.read()
            metadata["paragraphs"] = len([p for p in text.split('\n\n') if p.strip()])
            
        elif ext in ('.docx', '.doc'):
            try:
                import docx
                doc = docx.Document(file_path)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                text = '\n\n'.join(paragraphs)
                metadata["paragraphs"] = len(paragraphs)
                
                # 提取标题层级
                headings = []
                for p in doc.paragraphs:
                    if p.style.name.startswith('Heading'):
                        headings.append({"level": p.style.name, "text": p.text})
                metadata["headings"] = headings
            except ImportError:
                logger.warning("python-docx not installed, trying textract...")
                text = _fallback_text_extract(file_path)
                
        elif ext in ('.pptx', '.ppt'):
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                slides = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    if slide_text:
                        slides.append(f"[Slide {i+1}]\n" + '\n'.join(slide_text))
                text = '\n\n---\n\n'.join(slides)
                metadata["slides"] = len(prs.slides)
            except ImportError:
                text = _fallback_text_extract(file_path)
    except Exception as e:
        logger.error(f"Text parse failed for {file_path}: {e}")
        text = _fallback_text_extract(file_path)
    
    return {
        "type": "text",
        "text": text.strip(),
        "metadata": metadata,
        "tables": [],
        "images": []
    }


def _fallback_text_extract(file_path: str) -> str:
    """最终兜底：尝试用 textract 或直接读二进制"""
    try:
        import textract
        return textract.process(file_path).decode('utf-8', errors='replace')
    except ImportError:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()
        except Exception:
            return ""


# ============================================================
# PDF 解析器
# ============================================================

async def parse_pdf(file_path: str) -> Dict:
    """解析 PDF → 文本 + 表格，区分原生/扫描"""
    text = ""
    tables = []
    images = []
    metadata = {"source": file_path, "type": "pdf", "is_scanned": False}
    
    try:
        import pdfplumber
        
        with pdfplumber.open(file_path) as pdf:
            all_text = []
            for i, page in enumerate(pdf.pages):
                # 提取文本
                page_text = page.extract_text()
                if page_text:
                    all_text.append(f"[Page {i+1}]\n{page_text}")
                
                # 提取表格
                page_tables = page.extract_tables()
                for j, tbl in enumerate(page_tables):
                    if tbl:
                        # 转 Markdown 表格
                        md_table = _table_to_markdown(tbl)
                        tables.append({
                            "page": i + 1,
                            "table_index": j,
                            "markdown": md_table,
                            "rows": len(tbl),
                            "header": tbl[0] if tbl else []
                        })
            
            text = '\n\n'.join(all_text)
            
            # 判断是否扫描版（文本极少但有图片）
            if len(text.strip()) < 100:
                metadata["is_scanned"] = True
                # 扫描版用 OCR
                try:
                    ocr_text = await _ocr_pdf(file_path)
                    if ocr_text:
                        text = ocr_text
                except Exception as e:
                    logger.warning(f"OCR failed: {e}")
    
    except ImportError:
        logger.warning("pdfplumber not installed, falling back to PyMuPDF")
        try:
            import fitz
            doc = fitz.open(file_path)
            all_text = []
            for i, page in enumerate(doc):
                all_text.append(f"[Page {i+1}]\n{page.get_text()}")
            text = '\n\n'.join(all_text)
        except Exception:
            text = ""
    
    return {
        "type": "pdf",
        "text": text.strip(),
        "metadata": metadata,
        "tables": tables,
        "images": images
    }


async def _ocr_pdf(file_path: str) -> str:
    """OCR 扫描版 PDF"""
    try:
        import pytesseract
        from pdf2image import convert_from_path
        images = convert_from_path(file_path, first_page=1, last_page=20)
        texts = []
        for i, img in enumerate(images):
            t = pytesseract.image_to_string(img, lang='chi_sim+eng')
            if t.strip():
                texts.append(f"[Page {i+1}]\n{t}")
        return '\n\n'.join(texts)
    except ImportError:
        return ""


def _table_to_markdown(rows: List[List]) -> str:
    """将二维列表转为 Markdown 表格"""
    if not rows:
        return ""
    
    # 清理 None 和空值
    clean_rows = []
    for r in rows:
        clean_rows.append([str(c) if c else "" for c in r])
    
    lines = []
    header = clean_rows[0]
    lines.append('| ' + ' | '.join(header) + ' |')
    lines.append('|' + '|'.join(['---' for _ in header]) + '|')
    
    for row in clean_rows[1:]:
        lines.append('| ' + ' | '.join(row) + ' |')
    
    return '\n'.join(lines)


# ============================================================
# 表格解析器 (CSV/XLSX)
# ============================================================

async def parse_table(file_path: str) -> Dict:
    """解析 CSV/XLSX → 结构化 JSON"""
    ext = Path(file_path).suffix.lower()
    tables = []
    metadata = {"source": file_path, "type": "table", "format": ext}
    
    try:
        if ext == '.csv':
            tables = await _parse_csv(file_path)
        elif ext in ('.xlsx', '.xls'):
            tables = await _parse_xlsx(file_path)
    except Exception as e:
        logger.error(f"Table parse failed for {file_path}: {e}")
    
    # 生成描述文本（用于 BM25 检索）
    text_parts = []
    for t in tables:
        text_parts.append(f"Sheet: {t['sheet']} | Rows: {t['rows']} | Cols: {t['cols']}")
        if t.get('header'):
            text_parts.append('Columns: ' + ', '.join(t['header'][:20]))
        if t.get('sample_rows'):
            text_parts.append('Sample: ' + json.dumps(t['sample_rows'][:5], ensure_ascii=False))
    text = '\n'.join(text_parts)
    
    return {
        "type": "table",
        "text": text,
        "metadata": metadata,
        "tables": tables,
        "images": []
    }


async def _parse_csv(file_path: str) -> List[Dict]:
    """解析 CSV 文件"""
    tables = []
    for encoding in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
        try:
            with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                reader = csv.reader(f)
                rows = list(reader)
                if rows:
                    header = rows[0]
                    data = rows[1:] if len(rows) > 1 else []
                    tables.append({
                        "sheet": Path(file_path).stem,
                        "rows": len(rows),
                        "cols": len(header) if header else 0,
                        "header": header,
                        "sample_rows": data[:50],
                        "all_rows": data
                    })
            break
        except UnicodeDecodeError:
            continue
    
    return tables


async def _parse_xlsx(file_path: str) -> List[Dict]:
    """解析 XLSX 文件"""
    import openpyxl
    tables = []
    
    wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append([str(c) if c is not None else "" for c in row])
        
        if not rows:
            continue
        
        # 过滤全空行
        rows = [r for r in rows if any(c.strip() for c in r)]
        if not rows:
            continue
        
        header = rows[0]
        data = rows[1:] if len(rows) > 1 else []
        
        # 转 Markdown
        md_table = _table_to_markdown(rows[:50])  # 最多50行
        
        tables.append({
            "sheet": sheet_name,
            "rows": len(rows),
            "cols": len(header),
            "header": list(header),
            "sample_rows": data[:50],
            "all_rows": data,
            "markdown": md_table
        })
    
    wb.close()
    return tables


# ============================================================
# 图片解析器
# ============================================================

async def parse_image(file_path: str) -> Dict:
    """OCR 提取图片中文字"""
    text = ""
    metadata = {"source": file_path, "type": "image"}
    
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img, lang='chi_sim+eng')
        metadata["size"] = img.size
    except ImportError:
        # 无 OCR 时，只记录元数据
        try:
            from PIL import Image
            img = Image.open(file_path)
            metadata["size"] = img.size
        except Exception as e:

            logger.warning(f"[{module}] suppressed exception", exc_info=True)
        text = f"[图片: {Path(file_path).name}]"
    
    return {
        "type": "image",
        "text": text.strip(),
        "metadata": metadata,
        "tables": [],
        "images": [file_path]
    }


# ============================================================
# 压缩包处理
# ============================================================

async def parse_archive(file_path: str) -> List[Dict]:
    """解压压缩包，递归解析内部所有文件"""
    import tempfile
    import shutil
    import zipfile
    
    results = []
    extract_dir = tempfile.mkdtemp(prefix='fuxi_archive_')
    
    try:
        ext = Path(file_path).suffix.lower()
        
        if ext == '.zip':
            with zipfile.ZipFile(file_path, 'r') as zf:
                zf.extractall(extract_dir)
        elif ext in ('.rar', '.7z'):
            # 需要额外工具
            logger.warning(f"Archive format {ext} requires external tools")
            return results
        
        # 递归解析
        for root, dirs, files in os.walk(extract_dir):
            for fname in files:
                fpath = os.path.join(root, fname)
                ftype = identify_file(fpath)
                
                if ftype == 'archive':
                    continue  # 嵌套压缩包跳过
                elif ftype == 'text':
                    result = await parse_text(fpath)
                elif ftype == 'pdf':
                    result = await parse_pdf(fpath)
                elif ftype == 'table':
                    result = await parse_table(fpath)
                elif ftype == 'image':
                    result = await parse_image(fpath)
                else:
                    continue
                
                result["metadata"]["original_archive"] = file_path
                result["metadata"]["archive_path"] = Path(fpath).relative_to(extract_dir).as_posix()
                results.append(result)
    
    finally:
        shutil.rmtree(extract_dir, ignore_errors=True)
    
    return results


# ============================================================
# 统一入口
# ============================================================

async def parse_file(file_path: str) -> Dict:
    """统一解析入口：识别类型 → 分发解析器"""
    ftype = identify_file(file_path)
    
    if ftype == 'text':
        return await parse_text(file_path)
    elif ftype == 'pdf':
        return await parse_pdf(file_path)
    elif ftype == 'table':
        return await parse_table(file_path)
    elif ftype == 'image':
        return await parse_image(file_path)
    elif ftype == 'archive':
        results = await parse_archive(file_path)
        return {"type": "archive", "text": "", "metadata": {}, "tables": [], "images": [], "children": results}
    else:
        return {"type": "unknown", "text": "", "metadata": {"source": file_path}, "tables": [], "images": []}
